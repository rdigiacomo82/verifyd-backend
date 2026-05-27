# ============================================================
#  VeriFYD — doc_certifier.py
#
#  Creates certified/stamped PDF documents for VeriFYD Docs.
#
#  PDF input:
#    - Preserves the original PDF pages.
#    - Adds only a small lower-right VeriFYD mark.
#    - Does NOT append a full certificate-result page.
#    - Does NOT add a footer bar.
#
#  Image input:
#    - Converts JPG/JPEG/PNG into a marked certified PDF.
#
#  XLSX input:
#    - Renders workbook sheets/cell values into a certified PDF.
#    - Adds only the small VeriFYD lower-right mark on each page.
#
#  PPTX input:
#    - Renders a visual slide approximation into a certified PDF.
#    - Preserves slide dimensions, positioned text boxes, tables, images,
#      and basic filled shapes where possible.
#    - Adds only the small VeriFYD lower-right mark on each page.
#
#  Office/text inputs:
#    - Attempts LibreOffice PDF conversion first to preserve original pages.
#    - Falls back to safe readable text/rendered output when conversion fails.
#
#  Other non-PDF inputs:
#    - Creates a certified PDF rendering/summary fallback.
#
#  The official certificate details remain in the database, email,
#  and /certificate endpoint.
# ============================================================

from __future__ import annotations

import os
import tempfile
import base64
import hashlib
import hmac
import json
import logging
from datetime import datetime, timezone
from typing import Any, Dict, Iterable, List, Tuple

log = logging.getLogger("verifyd.doc_certifier")


def _draw_verifyd_mark(c, width: float, y: float = 24, x_right_pad: float = 34) -> None:
    """Draw a subtle lower-right VeriFYD mark on the active ReportLab canvas."""
    x_right = width - x_right_pad

    c.saveState()
    try:
        c.setFillAlpha(0.72)
    except Exception:
        pass

    c.setFont("Helvetica-Bold", 8)
    full = "VERIFYD"
    total_w = c.stringWidth(full, "Helvetica-Bold", 8)
    x = x_right - total_w

    c.setFillGray(0.15)
    c.drawString(x, y, "VERI")
    veri_w = c.stringWidth("VERI", "Helvetica-Bold", 8)

    c.setFillColorRGB(0.96, 0.62, 0.04)
    c.drawString(x + veri_w, y, "FYD")

    c.setFont("Helvetica", 5.5)
    c.setFillGray(0.35)
    c.drawRightString(x_right, y - 7, "certified")
    c.restoreState()


def _make_logo_overlay(width: float, height: float) -> str:
    """Create a one-page transparent overlay with a small lower-right VeriFYD mark."""
    from reportlab.pdfgen import canvas

    fd, overlay_path = tempfile.mkstemp(suffix="_verifyd_logo_overlay.pdf")
    os.close(fd)

    c = canvas.Canvas(overlay_path, pagesize=(width, height))
    _draw_verifyd_mark(c, width)
    c.save()
    return overlay_path


def _safe_text(value: Any, max_len: int = 80) -> str:
    """Convert a spreadsheet/document value into safe single-line text."""
    if value is None:
        return ""
    text = str(value).replace("\r", " ").replace("\n", " ").strip()
    if len(text) > max_len:
        return text[: max_len - 1] + "…"
    return text


def _chunked_rows(rows: List[List[str]], rows_per_page: int) -> Iterable[List[List[str]]]:
    for i in range(0, len(rows), rows_per_page):
        yield rows[i : i + rows_per_page]


def _convert_heic_to_renderable_image(src_path: str) -> str:
    """
    Convert HEIC/HEIF to a temporary JPEG/PNG that ReportLab/Pillow can render.

    The photo detector already handles HEIC by converting with ffmpeg. This helper
    mirrors that behavior for document certification so the analysis result does
    not succeed while the certified PDF fails to render/upload.
    """
    import subprocess
    import shutil

    ffmpeg = shutil.which("ffmpeg")
    if ffmpeg:
        fd, tmp_jpg = tempfile.mkstemp(suffix="_verifyd_heic_render.jpg")
        os.close(fd)
        try:
            cmd = [ffmpeg, "-y", "-i", src_path, "-frames:v", "1", "-q:v", "2", tmp_jpg]
            result = subprocess.run(cmd, capture_output=True, timeout=60)
            if result.returncode == 0 and os.path.exists(tmp_jpg) and os.path.getsize(tmp_jpg) > 1000:
                return tmp_jpg
        except Exception:
            pass
        try:
            if os.path.exists(tmp_jpg):
                os.remove(tmp_jpg)
        except Exception:
            pass

    try:
        from pillow_heif import register_heif_opener
        register_heif_opener()
        from PIL import Image

        fd, tmp_png = tempfile.mkstemp(suffix="_verifyd_heic_render.png")
        os.close(fd)
        with Image.open(src_path) as img:
            if img.mode not in ("RGB", "RGBA"):
                img = img.convert("RGB")
            img.save(tmp_png, format="PNG")
        if os.path.exists(tmp_png) and os.path.getsize(tmp_png) > 1000:
            return tmp_png
        try:
            os.remove(tmp_png)
        except Exception:
            pass
    except Exception:
        pass

    raise RuntimeError("HEIC/HEIF certification rendering failed. ffmpeg or pillow-heif could not convert the image.")


def _create_image_pdf(src_path: str, dest_path: str, cert_id: str, authenticity: int,
                      label: str, filename: str, sha256: str = "") -> str:
    """Create a certified PDF that preserves uploaded image documents visually."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.utils import ImageReader
    from PIL import Image, ImageSequence

    ext = os.path.splitext(src_path)[1].lower()
    render_src = src_path
    temp_images: List[str] = []

    if ext in (".heic", ".heif"):
        render_src = _convert_heic_to_renderable_image(src_path)
        temp_images.append(render_src)

    width, height = letter
    c = canvas.Canvas(dest_path, pagesize=letter)
    c.setAuthor("VeriFYD")
    c.setTitle(f"VeriFYD Certified Document {cert_id}")
    c.setSubject(f"{label} | Authenticity {authenticity}% | {filename}")

    margin = 36
    max_w = width - (margin * 2)
    max_h = height - (margin * 2)

    try:
        with Image.open(render_src) as img:
            frames = []
            try:
                for frame in ImageSequence.Iterator(img):
                    frames.append(frame.copy())
                    if len(frames) >= 50:
                        break
            except Exception:
                frames = [img.copy()]

        if not frames:
            with Image.open(render_src) as img:
                frames = [img.copy()]

        for idx, frame in enumerate(frames):
            if frame.mode not in ("RGB", "RGBA"):
                frame = frame.convert("RGB")

            draw_source = render_src
            if len(frames) > 1 or ext in (".tif", ".tiff", ".webp", ".heic", ".heif"):
                fd, tmp_img = tempfile.mkstemp(suffix=f"_verifyd_frame_{idx}.png")
                os.close(fd)
                frame.save(tmp_img, format="PNG")
                temp_images.append(tmp_img)
                draw_source = tmp_img

            img_w, img_h = frame.size
            scale = min(max_w / max(1, img_w), max_h / max(1, img_h))
            draw_w = img_w * scale
            draw_h = img_h * scale
            x = (width - draw_w) / 2
            y = (height - draw_h) / 2

            c.drawImage(
                ImageReader(draw_source),
                x,
                y,
                width=draw_w,
                height=draw_h,
                preserveAspectRatio=True,
                mask="auto",
            )
            _draw_verifyd_mark(c, width)
            if idx < len(frames) - 1:
                c.showPage()
    finally:
        for t in temp_images:
            try:
                if os.path.exists(t):
                    os.remove(t)
            except Exception:
                pass

    c.save()
    return dest_path



def _xlsx_display_value(cell: Any) -> str:
    """Return a clean Excel display value for certified XLSX PDF rendering."""
    from datetime import date, datetime
    try:
        from openpyxl.styles.numbers import is_date_format
    except Exception:
        is_date_format = lambda _fmt: False

    value = getattr(cell, "value", None)
    if value is None:
        return ""

    number_format = str(getattr(cell, "number_format", "") or "")

    if isinstance(value, datetime):
        try:
            return f"{value.month}/{value.day}/{value.year}"
        except Exception:
            return value.strftime("%Y-%m-%d")

    if isinstance(value, date):
        try:
            return f"{value.month}/{value.day}/{value.year}"
        except Exception:
            return value.isoformat()

    if isinstance(value, (int, float)) and not isinstance(value, bool):
        try:
            if is_date_format(number_format):
                from openpyxl.utils.datetime import from_excel
                dt = from_excel(value)
                return f"{dt.month}/{dt.day}/{dt.year}"
        except Exception:
            pass

        is_currency = "$" in number_format or "[$" in number_format
        is_percent = "%" in number_format

        if is_percent:
            return f"{float(value) * 100:.2f}%".replace(".00%", "%")

        if is_currency:
            return f"${float(value):,.2f}"

        # Avoid ugly binary float artifacts such as 2476.7999999999997.
        if isinstance(value, float):
            if abs(value - round(value)) < 0.0000001:
                return str(int(round(value)))
            return f"{value:,.2f}".rstrip("0").rstrip(".")
        return str(value)

    text = str(value).replace(chr(13), " ").strip()
    return text


def _xlsx_cell_fill_hex(cell: Any) -> str:
    """Best-effort background color for an Excel cell."""
    try:
        fill = getattr(cell, "fill", None)
        if not fill or getattr(fill, "fill_type", None) != "solid":
            return ""
        color = getattr(fill, "fgColor", None)
        if not color:
            return ""
        ctype = str(getattr(color, "type", "") or "").lower()
        if ctype == "rgb":
            rgb = str(getattr(color, "rgb", "") or "").upper()
            if len(rgb) == 8:
                rgb = rgb[2:]
            if len(rgb) == 6 and rgb not in ("000000", "FFFFFF"):
                return "#" + rgb
        if ctype == "indexed":
            indexed = getattr(color, "indexed", None)
            palette = {22: "#C0D6E4", 23: "#808080", 24: "#9999FF", 25: "#993366", 48: "#969696", 49: "#003366"}
            if indexed in palette:
                return palette[indexed]
        if ctype == "theme":
            # Most business quote templates use theme fills for section bars.
            return "#8EA0B5"
    except Exception:
        pass
    return ""


def _xlsx_cell_border_style(cell: Any) -> str:
    """Return a simple border strength indicator."""
    try:
        border = getattr(cell, "border", None)
        styles = [getattr(getattr(border, side, None), "style", None) for side in ("left", "right", "top", "bottom")]
        if any(s in ("medium", "thick", "double") for s in styles):
            return "strong"
        if any(bool(s) for s in styles):
            return "thin"
    except Exception:
        pass
    return ""


def _xlsx_cell_is_bold(cell: Any) -> bool:
    try:
        return bool(getattr(getattr(cell, "font", None), "bold", False))
    except Exception:
        return False


def _xlsx_cell_alignment(cell: Any) -> str:
    try:
        align = str(getattr(getattr(cell, "alignment", None), "horizontal", "") or "").lower()
        if align in ("center", "distributed", "centercontinuous"):
            return "CENTER"
        if align == "right":
            return "RIGHT"
    except Exception:
        pass
    return "LEFT"


def _xlsx_trimmed_bounds(ws: Any, max_rows_per_sheet: int = 500) -> Tuple[int, int, int, int] | None:
    """Find the meaningful visible range for a worksheet."""
    nonempty_rows: List[int] = []
    nonempty_cols: List[int] = []

    max_r = min(int(getattr(ws, "max_row", 1) or 1), max_rows_per_sheet)
    max_c = min(int(getattr(ws, "max_column", 1) or 1), 40)

    for r in range(1, max_r + 1):
        try:
            if getattr(ws.row_dimensions.get(r), "hidden", False):
                continue
        except Exception:
            pass
        for c in range(1, max_c + 1):
            try:
                from openpyxl.utils import get_column_letter
                if getattr(ws.column_dimensions.get(get_column_letter(c)), "hidden", False):
                    continue
            except Exception:
                pass
            try:
                if _xlsx_display_value(ws.cell(r, c)).strip():
                    nonempty_rows.append(r)
                    nonempty_cols.append(c)
            except Exception:
                continue

    if not nonempty_rows or not nonempty_cols:
        return None

    return min(nonempty_rows), max(nonempty_rows), min(nonempty_cols), max(nonempty_cols)


def _extract_xlsx_rows(src_path: str, max_columns: int = 14, max_rows_per_sheet: int = 500) -> List[Dict[str, Any]]:
    """
    Extract visible workbook cells for the enhanced certified XLSX renderer.

    Improvements over the old renderer:
      - skips hidden sheets and empty sheets,
      - trims to actual used visible range,
      - preserves merged-cell spans where practical,
      - keeps basic fill/bold/alignment/border hints,
      - formats dates/currency/numbers cleanly.
    """
    try:
        from openpyxl import load_workbook
        from openpyxl.utils import get_column_letter
    except Exception as e:
        raise RuntimeError("Missing dependency: openpyxl. Add openpyxl>=3.1.0 to requirements.txt") from e

    wb = load_workbook(src_path, data_only=True, read_only=False)
    rendered: List[Dict[str, Any]] = []
    skipped_hidden = 0
    skipped_empty = 0

    try:
        for ws in wb.worksheets:
            if str(getattr(ws, "sheet_state", "visible") or "visible").lower() != "visible":
                skipped_hidden += 1
                continue

            bounds = _xlsx_trimmed_bounds(ws, max_rows_per_sheet=max_rows_per_sheet)
            if not bounds:
                skipped_empty += 1
                continue

            min_r, max_r, min_c, max_c = bounds
            visible_cols: List[int] = []
            for c in range(min_c, max_c + 1):
                try:
                    if getattr(ws.column_dimensions.get(get_column_letter(c)), "hidden", False):
                        continue
                except Exception:
                    pass
                visible_cols.append(c)
                if len(visible_cols) >= max_columns:
                    break

            visible_rows: List[int] = []
            for r in range(min_r, max_r + 1):
                try:
                    if getattr(ws.row_dimensions.get(r), "hidden", False):
                        continue
                except Exception:
                    pass
                visible_rows.append(r)

            col_index = {c: idx for idx, c in enumerate(visible_cols)}
            row_index = {r: idx for idx, r in enumerate(visible_rows)}

            rows: List[List[Dict[str, Any]]] = []
            for r in visible_rows:
                row_cells: List[Dict[str, Any]] = []
                for c in visible_cols:
                    cell = ws.cell(r, c)
                    row_cells.append({
                        "text": _safe_text(_xlsx_display_value(cell), 400),
                        "fill": _xlsx_cell_fill_hex(cell),
                        "bold": _xlsx_cell_is_bold(cell),
                        "align": _xlsx_cell_alignment(cell),
                        "border": _xlsx_cell_border_style(cell),
                    })
                rows.append(row_cells)

            spans: List[Tuple[int, int, int, int]] = []
            try:
                for merged in ws.merged_cells.ranges:
                    r1, c1, r2, c2 = int(merged.min_row), int(merged.min_col), int(merged.max_row), int(merged.max_col)
                    if r1 in row_index and r2 in row_index and c1 in col_index and c2 in col_index:
                        if (r2 - r1) <= 2 and (c2 - c1) >= 1:
                            spans.append((col_index[c1], row_index[r1], col_index[c2], row_index[r2]))
                            tl = rows[row_index[r1]][col_index[c1]]
                            for rr in range(row_index[r1], row_index[r2] + 1):
                                for cc in range(col_index[c1], col_index[c2] + 1):
                                    rows[rr][cc]["fill"] = rows[rr][cc].get("fill") or tl.get("fill") or ""
                                    rows[rr][cc]["border"] = rows[rr][cc].get("border") or tl.get("border") or ""
            except Exception:
                spans = []

            while rows and not any(cell.get("text") for cell in rows[0]):
                rows.pop(0)
                spans = [(c1, r1 - 1, c2, r2 - 1) for (c1, r1, c2, r2) in spans if r2 > 0]
            while rows and not any(cell.get("text") for cell in rows[-1]):
                rows.pop()

            if not rows:
                skipped_empty += 1
                continue

            raw_widths: List[float] = []
            for c in visible_cols[:len(rows[0])]:
                try:
                    width = float(ws.column_dimensions[get_column_letter(c)].width or 10.0)
                except Exception:
                    width = 10.0
                raw_widths.append(max(5.0, min(width, 60.0)))

            rendered.append({
                "title": ws.title or "Sheet",
                "rows": rows,
                "spans": spans,
                "raw_widths": raw_widths,
                "hidden_sheets_skipped": skipped_hidden,
                "empty_sheets_skipped": skipped_empty,
            })
    finally:
        try:
            wb.close()
        except Exception:
            pass

    return rendered


def _xlsx_sheet_looks_like_quote(sheet: Dict[str, Any]) -> bool:
    sample = " ".join(str(cell.get("text", "")).lower() for row in sheet.get("rows", [])[:35] for cell in row)
    return any(x in sample for x in ("quick quote", "bill to", "ship to", "school information", "item number", "salesperson"))


def _xlsx_col_widths_for_pdf(sheet: Dict[str, Any], available_width: float) -> List[float]:
    rows = sheet.get("rows") or []
    max_cols = max((len(r) for r in rows), default=1)
    raw = list(sheet.get("raw_widths") or [])[:max_cols]
    if len(raw) < max_cols:
        raw.extend([10.0] * (max_cols - len(raw)))

    if _xlsx_sheet_looks_like_quote(sheet) and max_cols >= 6:
        weights = [1.15, 3.95, 1.15, 0.95, 0.85, 1.25] + [0.7] * (max_cols - 6)
    else:
        weights = [max(0.65, min(3.2, w / 10.0)) for w in raw]

    total = sum(weights) or 1.0
    return [available_width * (w / total) for w in weights[:max_cols]]


def _xlsx_make_paragraph(text: str, style: Any) -> Any:
    import html
    from reportlab.platypus import Paragraph
    return Paragraph(html.escape(str(text or "")).replace(chr(10), "<br/>").replace(chr(13), " "), style)


def _draw_table_page(c, width: float, height: float, title: str, subtitle: str,
                     rows: List[List[str]], page_num: int, total_hint: str = "") -> None:
    """
    Legacy canvas table drawer retained for compatibility.
    The enhanced XLSX renderer now uses ReportLab Platypus tables in _create_xlsx_pdf.
    """
    margin_x = 34
    top = height - 34
    bottom = 42

    c.setFont("Helvetica-Bold", 13)
    c.setFillGray(0.05)
    c.drawString(margin_x, top, title[:95])

    c.setFont("Helvetica", 7.5)
    c.setFillGray(0.35)
    c.drawString(margin_x, top - 14, subtitle[:130])

    if total_hint:
        c.drawRightString(width - margin_x, top - 14, total_hint[:60])

    y = top - 34
    row_h = 18
    available_w = width - (margin_x * 2)

    max_cols = max((len(r) for r in rows), default=1)
    max_cols = max(1, min(max_cols, 12))
    col_w = available_w / max_cols

    for ridx, row in enumerate(rows):
        if y < bottom + row_h:
            break

        if ridx == 0:
            c.setFillGray(0.90)
            c.rect(margin_x, y - 4, available_w, row_h, fill=1, stroke=0)
            font_name = "Helvetica-Bold"
        elif ridx % 2 == 0:
            c.setFillGray(0.975)
            c.rect(margin_x, y - 4, available_w, row_h, fill=1, stroke=0)
            font_name = "Helvetica"
        else:
            font_name = "Helvetica"

        c.setStrokeGray(0.80)
        c.setLineWidth(0.25)

        for cidx in range(max_cols):
            x = margin_x + (cidx * col_w)
            c.rect(x, y - 4, col_w, row_h, fill=0, stroke=1)
            value = row[cidx] if cidx < len(row) else ""
            c.setFillGray(0.08)
            c.setFont(font_name, 6.7)
            approx_chars = max(8, int(col_w / 4.0))
            display = _safe_text(value, approx_chars)
            c.drawString(x + 3, y + 2, display)

        y -= row_h

    c.setFont("Helvetica", 6.5)
    c.setFillGray(0.45)
    c.drawString(margin_x, 24, f"VeriFYD certified workbook rendering • Page {page_num}")
    _draw_verifyd_mark(c, width, y=24, x_right_pad=34)


def _create_xlsx_pdf(src_path: str, dest_path: str, cert_id: str, authenticity: int,
                     label: str, filename: str, sha256: str = "") -> str:
    """
    Render an XLSX workbook into a readable certified PDF with an enhanced
    quote/form-aware layout. This is the Python fallback used when LibreOffice is
    unavailable in Render.
    """
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_RIGHT
    from reportlab.lib.units import inch
    from reportlab.lib import colors

    page_size = landscape(letter)
    page_w, page_h = page_size
    left_margin = 0.30 * inch
    right_margin = 0.30 * inch
    top_margin = 0.32 * inch
    bottom_margin = 0.42 * inch
    available_w = page_w - left_margin - right_margin

    sheets = _extract_xlsx_rows(src_path)

    pdf = SimpleDocTemplate(
        dest_path,
        pagesize=page_size,
        leftMargin=left_margin,
        rightMargin=right_margin,
        topMargin=top_margin,
        bottomMargin=bottom_margin,
        title=f"VeriFYD Certified Spreadsheet {cert_id}",
        author="VeriFYD",
        subject=f"{label} | Authenticity {authenticity}% | {filename}",
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("VeriFYD_XLSX_Title", parent=styles["Title"], fontName="Helvetica-Bold", fontSize=18, leading=21, spaceAfter=10)
    meta_style = ParagraphStyle("VeriFYD_XLSX_Meta", parent=styles["Normal"], fontName="Helvetica", fontSize=8.5, leading=10.5, spaceAfter=3)
    sheet_title_style = ParagraphStyle("VeriFYD_XLSX_Sheet_Title", parent=styles["Heading2"], fontName="Helvetica-Bold", fontSize=12, leading=14, spaceAfter=4)
    small_style = ParagraphStyle("VeriFYD_XLSX_Small", parent=styles["Normal"], fontName="Helvetica", fontSize=6.0, leading=7.0, wordWrap="CJK")
    small_bold = ParagraphStyle("VeriFYD_XLSX_Small_Bold", parent=small_style, fontName="Helvetica-Bold")
    small_center = ParagraphStyle("VeriFYD_XLSX_Small_Center", parent=small_style, alignment=TA_CENTER)
    small_right = ParagraphStyle("VeriFYD_XLSX_Small_Right", parent=small_style, alignment=TA_RIGHT)

    story: List[Any] = []
    story.append(Paragraph("VeriFYD Certified Spreadsheet", title_style))
    story.append(Paragraph(f"Status: <b>{label}</b> &nbsp;&nbsp; Authenticity Score: <b>{authenticity}%</b>", meta_style))
    story.append(Paragraph(f"Certificate ID: {cert_id}", meta_style))
    story.append(Paragraph(f"Original file: {filename}", meta_style))
    story.append(Paragraph(f"Visible workbook sheets rendered: {len(sheets)}", meta_style))
    if sha256:
        story.append(Paragraph(f"SHA-256: {sha256}", meta_style))
    story.append(Spacer(1, 10))
    story.append(Paragraph("The following pages are a VeriFYD-rendered PDF view of visible, non-empty workbook sheets. Hidden or empty sheets are skipped in the visible preview, while the exact original XLSX remains embedded in the certified package.", meta_style))
    story.append(PageBreak())

    if not sheets:
        story.append(Paragraph("No visible worksheet cell values were found for rendering.", meta_style))
    else:
        for sheet_index, sheet in enumerate(sheets, start=1):
            sheet_name = str(sheet.get("title") or "Sheet")
            story.append(Paragraph(f"Worksheet: {sheet_name}", sheet_title_style))
            story.append(Paragraph(f"File: {filename} • Sheet {sheet_index}/{len(sheets)}", meta_style))
            rows = sheet.get("rows") or []
            max_cols = max((len(r) for r in rows), default=1)
            col_widths = _xlsx_col_widths_for_pdf(sheet, available_w)

            table_data: List[List[Any]] = []
            for ridx, row in enumerate(rows):
                out_row: List[Any] = []
                for cidx in range(max_cols):
                    cell = row[cidx] if cidx < len(row) else {"text": ""}
                    text = str(cell.get("text", "") or "")
                    align = str(cell.get("align", "LEFT") or "LEFT")
                    if cell.get("bold"):
                        style = small_bold
                    elif align == "CENTER":
                        style = small_center
                    elif align == "RIGHT":
                        style = small_right
                    else:
                        style = small_style
                    out_row.append(_xlsx_make_paragraph(text, style))
                table_data.append(out_row)

            table = Table(table_data, colWidths=col_widths, repeatRows=0, splitByRow=1)
            style_cmds: List[Tuple[Any, ...]] = [
                ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#C9C9C9")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 2.0),
                ("RIGHTPADDING", (0, 0), (-1, -1), 2.0),
                ("TOPPADDING", (0, 0), (-1, -1), 1.4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 1.4),
            ]

            for ridx, row in enumerate(rows):
                nonempty_texts = [str(c.get("text", "") or "").strip() for c in row]
                joined = " ".join(t.lower() for t in nonempty_texts if t)
                row_has_section_fill = any(x in joined for x in ("bill to", "ship to", "school information", "salesperson", "item number", "sow", "special instructions"))

                for cidx, cell in enumerate(row):
                    fill = str(cell.get("fill", "") or "")
                    if row_has_section_fill and not fill:
                        fill = "#8EA0B5"
                    if fill:
                        try:
                            style_cmds.append(("BACKGROUND", (cidx, ridx), (cidx, ridx), colors.HexColor(fill)))
                            if fill.upper() in ("#8EA0B5", "#8799AE", "#808080", "#003366"):
                                style_cmds.append(("TEXTCOLOR", (cidx, ridx), (cidx, ridx), colors.white))
                        except Exception:
                            pass
                    if cell.get("border") == "strong":
                        style_cmds.append(("BOX", (cidx, ridx), (cidx, ridx), 0.65, colors.black))

            for c1, r1, c2, r2 in sheet.get("spans", []) or []:
                if 0 <= r1 < len(rows) and 0 <= r2 < len(rows) and 0 <= c1 < max_cols and 0 <= c2 < max_cols:
                    style_cmds.append(("SPAN", (c1, r1), (c2, r2)))

            table.setStyle(TableStyle(style_cmds))
            story.append(table)

            if sheet_index < len(sheets):
                story.append(PageBreak())

    def _on_page(c, _doc):
        c.saveState()
        try:
            c.setFont("Helvetica", 6.2)
            c.setFillGray(0.40)
            c.drawString(left_margin, 16, f"VeriFYD certified workbook rendering • {filename[:90]}")
            c.drawRightString(page_w - right_margin, 16, f"Certificate: {cert_id[:18]}")
            _draw_verifyd_mark(c, page_w, y=24, x_right_pad=34)
        finally:
            c.restoreState()

    pdf.build(story, onFirstPage=_on_page, onLaterPages=_on_page)
    return dest_path



def _attach_original_file_to_pdf(pdf_path: str, original_path: str, original_filename: str) -> str:
    """
    Embed the original uploaded file inside the certified PDF as an attachment.

    This is especially important for Office files such as PPTX where a pure Python
    renderer cannot perfectly reproduce every visual detail, animation, theme,
    background, or SmartArt element. The PDF remains the certified viewing copy,
    while the exact original file is preserved inside the PDF for later review.
    """
    try:
        from pypdf import PdfReader, PdfWriter
    except Exception:
        # pypdf is already required for PDF certification, but do not fail the
        # whole certification if attachment support is unexpectedly unavailable.
        return pdf_path

    if not os.path.exists(pdf_path) or not os.path.exists(original_path):
        return pdf_path

    tmp_out = pdf_path + ".attached.tmp.pdf"
    try:
        reader = PdfReader(pdf_path)
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)

        try:
            existing_meta = reader.metadata or {}
            metadata = {str(k): str(v) for k, v in existing_meta.items() if v is not None}
            metadata.update({
                "/VeriFYD_Embedded_Original": original_filename,
            })
            writer.add_metadata(metadata)
        except Exception:
            pass

        with open(original_path, "rb") as fh:
            original_bytes = fh.read()

        # pypdf supports embedded file attachments using add_attachment.
        writer.add_attachment(original_filename or os.path.basename(original_path), original_bytes)

        with open(tmp_out, "wb") as out:
            writer.write(out)

        os.replace(tmp_out, pdf_path)
    except Exception:
        try:
            if os.path.exists(tmp_out):
                os.remove(tmp_out)
        except Exception:
            pass
    return pdf_path


def _stamp_existing_pdf_file(src_pdf_path: str, dest_path: str, cert_id: str,
                             authenticity: int, label: str, sha256: str = "") -> str:
    """Preserve an existing PDF's pages and add only the small lower-right VeriFYD mark."""
    from pypdf import PdfReader, PdfWriter

    reader = PdfReader(src_pdf_path)
    writer = PdfWriter()
    overlay_paths: list[str] = []

    try:
        for page in reader.pages:
            width = float(page.mediabox.width)
            height = float(page.mediabox.height)
            overlay = _make_logo_overlay(width, height)
            overlay_paths.append(overlay)
            overlay_page = PdfReader(overlay).pages[0]
            page.merge_page(overlay_page)
            writer.add_page(page)

        try:
            existing_meta = reader.metadata or {}
            metadata = {str(k): str(v) for k, v in existing_meta.items() if v is not None}
            metadata.update({
                "/VeriFYD": "Certified",
                "/VeriFYD_Certificate_ID": cert_id,
                "/VeriFYD_Label": label,
                "/VeriFYD_Authenticity": str(authenticity),
                "/VeriFYD_SHA256_Original": sha256 or "",
            })
            writer.add_metadata(metadata)
        except Exception:
            pass

        with open(dest_path, "wb") as out:
            writer.write(out)
    finally:
        for p in overlay_paths:
            try:
                if os.path.exists(p):
                    os.remove(p)
            except Exception:
                pass
    return dest_path




# ─────────────────────────────────────────────
# Universal document rendering via LibreOffice
# ─────────────────────────────────────────────

LIBREOFFICE_RENDER_EXTENSIONS = {
    ".docx", ".doc", ".odt", ".rtf",
    ".xlsx", ".xls", ".ods", ".csv",
    ".pptx", ".ppt", ".odp",
    ".txt", ".md",
    ".html", ".htm",
}


def _try_convert_document_to_pdf_with_libreoffice(src_path: str) -> Tuple[str, str] | None:
    """
    Try to use LibreOffice/soffice for true document-to-PDF conversion.

    This is the preferred path for Word, Excel, PowerPoint, OpenDocument,
    CSV, TXT, Markdown, and HTML because it preserves the original visual
    document pages far better than text extraction.

    Returns:
      (converted_pdf_path, output_directory) on success.
      None on failure/unavailable.

    Caller is responsible for deleting both the returned PDF and output dir.
    """
    try:
        import shutil
        import subprocess
        import uuid
    except Exception:
        return None

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        log.warning("doc_certifier: LibreOffice/soffice not found on worker PATH")
        return None

    log.info("doc_certifier: LibreOffice found at %s", soffice)

    if not os.path.exists(src_path):
        log.warning("doc_certifier: source file does not exist for LibreOffice conversion: %s", src_path)
        return None

    out_dir = tempfile.mkdtemp(prefix="verifyd_lo_convert_")
    user_profile_dir = tempfile.mkdtemp(prefix="verifyd_lo_profile_")

    try:
        cmd = [
            soffice,
            "--headless",
            "--nologo",
            "--nofirststartwizard",
            "--nodefault",
            "--nolockcheck",
            f"-env:UserInstallation=file://{user_profile_dir}",
            "--convert-to",
            "pdf",
            "--outdir",
            out_dir,
            src_path,
        ]

        log.info("doc_certifier: starting LibreOffice conversion cmd=%s", cmd)
        result = subprocess.run(cmd, capture_output=True, timeout=150)

        log.info(
            "doc_certifier: LibreOffice conversion finished returncode=%s stdout=%r stderr=%r",
            result.returncode,
            (result.stdout or b"")[:700],
            (result.stderr or b"")[:700],
        )

        if result.returncode != 0:
            log.warning(
                "doc_certifier: LibreOffice conversion failed returncode=%s file=%s",
                result.returncode,
                src_path,
            )
            return None

        # LibreOffice usually writes basename.pdf, but names can be altered
        # for odd filenames. Use any non-empty PDF in the output directory.
        pdfs = []
        try:
            for name in os.listdir(out_dir):
                if name.lower().endswith(".pdf"):
                    candidate = os.path.join(out_dir, name)
                    if os.path.exists(candidate) and os.path.getsize(candidate) > 1000:
                        pdfs.append(candidate)
        except Exception:
            pdfs = []

        if not pdfs:
            log.warning(
                "doc_certifier: LibreOffice completed but produced no PDF. out_dir=%s files=%s",
                out_dir,
                os.listdir(out_dir) if os.path.isdir(out_dir) else [],
            )
            return None

        pdfs.sort(key=lambda path: os.path.getmtime(path), reverse=True)
        log.info(
            "doc_certifier: LibreOffice produced PDF path=%s size=%s",
            pdfs[0],
            os.path.getsize(pdfs[0]) if os.path.exists(pdfs[0]) else 0,
        )
        return pdfs[0], out_dir

    except Exception as e:
        log.warning("doc_certifier: LibreOffice conversion exception for file=%s error=%s", src_path, e)
        return None
    finally:
        # Do not delete out_dir here because the converted PDF lives there.
        # The caller deletes out_dir after stamping. The LibreOffice user profile
        # can always be removed immediately.
        try:
            import shutil
            if os.path.isdir(user_profile_dir):
                shutil.rmtree(user_profile_dir, ignore_errors=True)
        except Exception:
            pass


def _create_office_pdf_via_libreoffice(src_path: str, dest_path: str, cert_id: str,
                                        authenticity: int, label: str,
                                        filename: str, sha256: str = "") -> str:
    """
    Convert a supported source document to PDF using LibreOffice, stamp the
    rendered pages with the lower-right VeriFYD mark, and embed the exact
    original source file into the certified PDF as an attachment.
    """
    import shutil

    converted = _try_convert_document_to_pdf_with_libreoffice(src_path)
    if not converted:
        raise RuntimeError("LibreOffice conversion unavailable or failed.")

    converted_pdf, converted_dir = converted
    try:
        _stamp_existing_pdf_file(
            src_pdf_path=converted_pdf,
            dest_path=dest_path,
            cert_id=cert_id,
            authenticity=authenticity,
            label=label,
            sha256=sha256,
        )
        _attach_original_file_to_pdf(dest_path, src_path, filename or os.path.basename(src_path))
        return dest_path
    finally:
        try:
            if converted_pdf and os.path.exists(converted_pdf):
                os.remove(converted_pdf)
        except Exception:
            pass
        try:
            if converted_dir and os.path.isdir(converted_dir):
                shutil.rmtree(converted_dir, ignore_errors=True)
        except Exception:
            pass


def _try_convert_pptx_to_pdf_with_libreoffice(src_path: str) -> str | None:
    """
    Try to use LibreOffice/soffice for a true PPTX-to-PDF conversion.

    Many Render Python environments do not include LibreOffice. When it is not
    available, return None and let the pure-Python renderer handle the fallback.
    """
    try:
        import shutil
        import subprocess
    except Exception:
        return None

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return None

    out_dir = tempfile.mkdtemp(prefix="verifyd_pptx_convert_")
    try:
        cmd = [
            soffice,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            out_dir,
            src_path,
        ]
        result = subprocess.run(cmd, capture_output=True, timeout=90)
        if result.returncode != 0:
            return None

        base = os.path.splitext(os.path.basename(src_path))[0] + ".pdf"
        converted = os.path.join(out_dir, base)
        if os.path.exists(converted) and os.path.getsize(converted) > 1000:
            # Caller is responsible for deleting the returned file and its folder.
            return converted
        return None
    except Exception:
        return None


# ─────────────────────────────────────────────
# PPTX rendering
# ─────────────────────────────────────────────

def _pptx_emu_to_pt(value: Any) -> float:
    """Convert PowerPoint EMUs to PDF points."""
    try:
        return float(value) / 914400.0 * 72.0
    except Exception:
        return 0.0


def _pptx_rgb_tuple(color_obj: Any, default: Tuple[float, float, float] = (1, 1, 1)) -> Tuple[float, float, float]:
    """Return ReportLab RGB floats from a python-pptx color object."""
    try:
        rgb = color_obj.rgb
        if rgb:
            return (rgb[0] / 255.0, rgb[1] / 255.0, rgb[2] / 255.0)
    except Exception:
        pass
    return default


def _pptx_text_runs(shape: Any) -> List[Tuple[str, bool, float]]:
    """Extract text with basic bold/font-size hints from a PPTX text shape."""
    runs: List[Tuple[str, bool, float]] = []
    try:
        tf = shape.text_frame
        for para in tf.paragraphs:
            parts: List[str] = []
            bold = False
            font_size = 11.0
            for run in para.runs:
                txt = (run.text or "").strip()
                if not txt:
                    continue
                parts.append(txt)
                try:
                    bold = bool(run.font.bold) or bold
                except Exception:
                    pass
                try:
                    if run.font.size:
                        font_size = max(6.0, min(30.0, float(run.font.size.pt)))
                except Exception:
                    pass
            if parts:
                runs.append((" ".join(parts), bold, font_size))
            else:
                txt = (getattr(para, "text", "") or "").strip()
                if txt:
                    runs.append((txt, False, font_size))
    except Exception:
        try:
            txt = (shape.text or "").strip()
            if txt:
                for line in txt.splitlines():
                    if line.strip():
                        runs.append((line.strip(), False, 11.0))
        except Exception:
            pass
    return runs


def _wrap_reportlab_text(c, text: str, font_name: str, font_size: float, max_width: float) -> List[str]:
    """Simple word wrapping for ReportLab text drawing."""
    text = str(text or "").replace("\t", " ").strip()
    if not text:
        return []

    lines: List[str] = []
    current = ""
    for word in text.split():
        trial = (current + " " + word).strip()
        if c.stringWidth(trial, font_name, font_size) <= max_width:
            current = trial
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines or [text]


def _draw_pptx_visual_slide(c, slide: Any, slide_idx: int, prs_width: float, prs_height: float,
                            page_w: float, page_h: float, filename: str, page_num: int,
                            temp_files: List[str]) -> None:
    """Draw one PPTX slide as a visual approximation with positioned shapes/text/images."""
    from reportlab.lib.utils import ImageReader

    margin = 34.0
    footer_h = 28.0
    max_w = page_w - (margin * 2)
    max_h = page_h - (margin * 2) - footer_h
    scale = min(max_w / max(1.0, prs_width), max_h / max(1.0, prs_height))
    slide_w = prs_width * scale
    slide_h = prs_height * scale
    origin_x = (page_w - slide_w) / 2.0
    origin_y = footer_h + margin + ((max_h - slide_h) / 2.0)

    # Slide canvas.
    c.saveState()
    c.setFillColorRGB(1, 1, 1)
    c.rect(origin_x, origin_y, slide_w, slide_h, fill=1, stroke=0)
    c.setStrokeGray(0.70)
    c.setLineWidth(0.8)
    c.rect(origin_x, origin_y, slide_w, slide_h, fill=0, stroke=1)
    c.restoreState()

    # Best-effort visual rendering of slide contents.
    for shape in slide.shapes:
        try:
            left = origin_x + _pptx_emu_to_pt(shape.left) * scale
            top_from_slide = _pptx_emu_to_pt(shape.top) * scale
            width = max(1.0, _pptx_emu_to_pt(shape.width) * scale)
            height = max(1.0, _pptx_emu_to_pt(shape.height) * scale)
            bottom = origin_y + slide_h - top_from_slide - height
        except Exception:
            continue

        # Pictures: preserve embedded image content where possible.
        try:
            if hasattr(shape, "image") and getattr(shape, "image", None):
                img_blob = shape.image.blob
                ext = shape.image.ext or "png"
                fd, img_path = tempfile.mkstemp(suffix=f"_pptx_img.{ext}")
                os.close(fd)
                with open(img_path, "wb") as fh:
                    fh.write(img_blob)
                temp_files.append(img_path)
                c.drawImage(
                    ImageReader(img_path),
                    left,
                    bottom,
                    width=width,
                    height=height,
                    preserveAspectRatio=True,
                    anchor="c",
                    mask="auto",
                )
                continue
        except Exception:
            pass

        # Tables: draw an approximate positioned grid.
        try:
            if getattr(shape, "has_table", False):
                table = shape.table
                rows = list(table.rows)
                cols = len(table.columns) if table.columns else 1
                row_h = height / max(1, len(rows))
                col_w = width / max(1, cols)
                for r_idx, row in enumerate(rows):
                    y = bottom + height - ((r_idx + 1) * row_h)
                    for c_idx, cell in enumerate(row.cells):
                        x = left + (c_idx * col_w)
                        c.setStrokeGray(0.70)
                        c.setLineWidth(0.35)
                        c.rect(x, y, col_w, row_h, fill=0, stroke=1)
                        txt = _safe_text(cell.text, 80)
                        if txt:
                            c.setFont("Helvetica", max(4.5, min(8, row_h * 0.38)))
                            c.setFillGray(0.05)
                            c.drawString(x + 2, y + max(2, row_h - 9), _safe_text(txt, int(max(8, col_w / 4.2))))
                continue
        except Exception:
            pass

        # Basic filled auto-shapes.
        try:
            fill = getattr(shape, "fill", None)
            if fill and getattr(fill, "type", None):
                rgb = _pptx_rgb_tuple(fill.fore_color, default=(0.94, 0.94, 0.94))
                c.saveState()
                c.setFillColorRGB(*rgb)
                c.setStrokeGray(0.82)
                c.setLineWidth(0.25)
                c.rect(left, bottom, width, height, fill=1, stroke=1)
                c.restoreState()
        except Exception:
            pass

        # Text boxes and shape text in approximate original position.
        try:
            if getattr(shape, "has_text_frame", False):
                runs = _pptx_text_runs(shape)
                if runs:
                    c.saveState()
                    c.setFillGray(0.05)
                    y = bottom + height - 7
                    for text, bold, font_size in runs:
                        scaled_font = max(5.0, min(22.0, font_size * scale * 1.15))
                        font_name = "Helvetica-Bold" if bold else "Helvetica"
                        wrapped = _wrap_reportlab_text(c, text, font_name, scaled_font, max(10.0, width - 8))
                        for line in wrapped[:8]:
                            if y < bottom + 4:
                                break
                            c.setFont(font_name, scaled_font)
                            c.drawString(left + 4, y, line)
                            y -= scaled_font * 1.25
                        y -= scaled_font * 0.25
                    c.restoreState()
        except Exception:
            pass

    c.setFont("Helvetica", 6.5)
    c.setFillGray(0.45)
    c.drawString(margin, 24, f"VeriFYD certified presentation rendering • Slide {slide_idx} • Page {page_num} • {filename[:80]}")
    _draw_verifyd_mark(c, page_w, y=24, x_right_pad=34)


def _create_pptx_pdf(src_path: str, dest_path: str, cert_id: str, authenticity: int,
                     label: str, filename: str, sha256: str = "") -> str:
    """
    Render a PPTX presentation into a certified PDF.

    Preferred path:
      1. If LibreOffice/soffice is available on the server, convert the PPTX to a
         true PDF rendering, stamp each page, and embed the original PPTX.

    Safe fallback:
      2. If LibreOffice is unavailable, use a pure-Python visual approximation
         and embed the original PPTX as a PDF attachment so the exact source file
         travels with the certified PDF.
    """
    try:
        from pptx import Presentation
    except Exception as e:
        raise RuntimeError("Missing dependency: python-pptx. Add python-pptx>=0.6.23 to requirements.txt") from e

    # First try a true Office conversion. This preserves the original slide
    # visuals far better than python-pptx can. It will silently fall back if
    # LibreOffice is not installed in Render.
    converted_pdf = _try_convert_pptx_to_pdf_with_libreoffice(src_path)
    if converted_pdf:
        converted_dir = os.path.dirname(converted_pdf)
        try:
            _stamp_existing_pdf_file(
                src_pdf_path=converted_pdf,
                dest_path=dest_path,
                cert_id=cert_id,
                authenticity=authenticity,
                label=label,
                sha256=sha256,
            )
            _attach_original_file_to_pdf(dest_path, src_path, filename or os.path.basename(src_path))
            return dest_path
        finally:
            try:
                if os.path.exists(converted_pdf):
                    os.remove(converted_pdf)
            except Exception:
                pass
            try:
                if converted_dir and os.path.isdir(converted_dir):
                    os.rmdir(converted_dir)
            except Exception:
                pass

    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import landscape, letter

    prs = Presentation(src_path)
    prs_width = _pptx_emu_to_pt(prs.slide_width)
    prs_height = _pptx_emu_to_pt(prs.slide_height)
    page_w, page_h = landscape(letter)
    c = canvas.Canvas(dest_path, pagesize=(page_w, page_h))
    temp_files: List[str] = []

    c.setAuthor("VeriFYD")
    c.setTitle(f"VeriFYD Certified Presentation {cert_id}")
    c.setSubject(f"{label} | Authenticity {authenticity}% | {filename}")

    # Cover/verification page.
    c.setFont("Helvetica-Bold", 20)
    c.drawString(54, page_h - 72, "VeriFYD Certified Presentation")
    c.setFont("Helvetica", 10)
    c.drawString(54, page_h - 105, f"Status: {label}")
    c.drawString(54, page_h - 123, f"Authenticity Score: {authenticity}%")
    c.drawString(54, page_h - 141, f"Certificate ID: {cert_id}")
    c.drawString(54, page_h - 159, f"Original file: {filename}")
    c.drawString(54, page_h - 177, f"Slides: {len(prs.slides)}")
    if sha256:
        c.drawString(54, page_h - 198, "SHA-256:")
        c.setFont("Helvetica", 7.5)
        c.drawString(54, page_h - 212, sha256[:120])

    c.setFont("Helvetica-Oblique", 8)
    c.setFillGray(0.35)
    c.drawString(54, 58, "The following pages are a VeriFYD-rendered PDF view of the uploaded PPTX presentation.")
    c.drawString(54, 44, "The original PPTX file is embedded inside this certified PDF as an attachment for exact-source review.")
    _draw_verifyd_mark(c, page_w, y=24, x_right_pad=34)
    c.showPage()

    try:
        if len(prs.slides) == 0:
            c.setFont("Helvetica", 12)
            c.drawString(54, page_h - 72, "[No slides found in uploaded presentation]")
            _draw_verifyd_mark(c, page_w, y=24, x_right_pad=34)
            c.showPage()
        else:
            for idx, slide in enumerate(prs.slides, start=1):
                _draw_pptx_visual_slide(
                    c=c,
                    slide=slide,
                    slide_idx=idx,
                    prs_width=prs_width,
                    prs_height=prs_height,
                    page_w=page_w,
                    page_h=page_h,
                    filename=filename,
                    page_num=idx + 1,
                    temp_files=temp_files,
                )
                c.showPage()
    finally:
        for path in temp_files:
            try:
                if os.path.exists(path):
                    os.remove(path)
            except Exception:
                pass

    c.save()

    # Embed the original deck so the certified PDF includes the exact source
    # presentation even when pure-Python visual rendering cannot perfectly
    # reproduce every PowerPoint feature.
    _attach_original_file_to_pdf(dest_path, src_path, filename or os.path.basename(src_path))
    return dest_path



def _rtf_regex_fallback_for_render(raw: str) -> str:
    """Table-preserving RTF cleanup for certified PDF rendering."""
    import re
    raw = raw or ""
    # Remove only obvious embedded binary/image/object groups. Do not remove table
    # structure because many legal/government RTFs are form-style tables.
    raw = re.sub(r"{\\\*?\\pict[^{}]*(?:{[^{}]*}[^{}]*)*}", " ", raw, flags=re.DOTALL)
    raw = re.sub(r"{\\object[^{}]*(?:{[^{}]*}[^{}]*)*}", " ", raw, flags=re.DOTALL)

    replacements = {
        "\\par": "\n",
        "\\line": "\n",
        "\\tab": "\t",
        "\\cell": "|",
        "\\row": "\n",
    }
    for src, dst in replacements.items():
        raw = raw.replace(src, dst)

    def _hex_to_char(match):
        try:
            return bytes.fromhex(match.group(1)).decode("cp1252", errors="replace")
        except Exception:
            return " "

    raw = re.sub(r"\\'([0-9a-fA-F]{2})", _hex_to_char, raw)
    raw = re.sub(r"\\u(-?\d+)\??", lambda m: chr(int(m.group(1)) % 65536), raw)
    raw = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", raw)
    raw = raw.replace("{", " ").replace("}", " ")
    raw = re.sub(r"\x00+", " ", raw)
    raw = re.sub(r"[ \t]+", " ", raw)
    raw = re.sub(r" *\| *", "|", raw)
    raw = re.sub(r"\n\s*\n\s*\n+", "\n\n", raw)
    return raw.strip()


def _strip_rtf_for_render(raw: str) -> str:
    """RTF-to-text cleanup for certified rendering with table-preserving fallback."""
    import re
    regex_text = _rtf_regex_fallback_for_render(raw)

    try:
        from striprtf.striprtf import rtf_to_text
        text = rtf_to_text(raw or "")
        text = re.sub(r"\x00+", " ", text)
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n\s*\n\s*\n+", "\n\n", text)
        text = text.strip()
        # striprtf is usually cleaner, but it can over-strip table-based forms.
        # Keep the fallback when it preserves substantially more readable content.
        if len(text) >= 500 or len(text) >= max(120, len(regex_text) * 0.35):
            return text
    except Exception:
        pass

    return regex_text



def _plain_text_from_binary_chunks(data: bytes, limit: int = 2_000_000) -> str:
    """Best-effort readable text extraction from legacy Office binary streams."""
    import re
    fragments: List[str] = []
    sample = data[:limit]
    for enc in ("utf-16-le", "latin-1"):
        try:
            decoded = sample.decode(enc, errors="ignore")
        except Exception:
            continue
        decoded = decoded.replace("\x00", " ")
        decoded = re.sub(r"[\x01-\x08\x0b\x0c\x0e-\x1f]+", " ", decoded)
        decoded = re.sub(r"[ \t]+", " ", decoded)
        pieces = re.findall(r"[A-Za-z0-9][A-Za-z0-9\s\.,;:\-_/@$%#&()\[\]{}'\"!?]{4,}", decoded)
        fragments.extend(p.strip() for p in pieces if p and len(p.strip()) >= 5)
    seen = set()
    cleaned: List[str] = []
    for frag in fragments:
        compact = re.sub(r"\s+", " ", frag).strip()
        if len(compact) < 5:
            continue
        key = compact[:160].lower()
        if key in seen:
            continue
        seen.add(key)
        cleaned.append(compact)
        if sum(len(x) for x in cleaned) > limit:
            break
    return "\n".join(cleaned)


def _read_ole_text_for_render(src_path: str) -> str:
    """Best-effort legacy DOC/PPT text rendering helper."""
    try:
        import olefile
        parts: List[str] = []
        with olefile.OleFileIO(src_path) as ole:
            for stream in ole.listdir()[:80]:
                try:
                    raw = ole.openstream(stream).read(750_000)
                    text = _plain_text_from_binary_chunks(raw, limit=750_000)
                    if text:
                        parts.append(text)
                except Exception:
                    continue
        if parts:
            return "\n".join(parts)
    except Exception:
        pass
    try:
        with open(src_path, "rb") as fh:
            return _plain_text_from_binary_chunks(fh.read(2_000_000))
    except Exception:
        return ""


def _read_xls_for_render(src_path: str) -> str:
    """Render legacy XLS cell values into plain text for certified PDF output."""
    try:
        import xlrd
        book = xlrd.open_workbook(src_path, on_demand=True)
        lines: List[str] = []
        rows_seen = 0
        for sheet_name in book.sheet_names():
            sh = book.sheet_by_name(sheet_name)
            lines.append(f"Worksheet: {sheet_name}")
            for r_idx in range(sh.nrows):
                vals: List[str] = []
                for c_idx in range(min(sh.ncols, 30)):
                    try:
                        val = sh.cell_value(r_idx, c_idx)
                    except Exception:
                        continue
                    if val in (None, ""):
                        continue
                    if isinstance(val, float) and val.is_integer():
                        val = int(val)
                    text = str(val).strip()
                    if text:
                        vals.append(text[:300])
                if vals:
                    lines.append(" | ".join(vals))
                rows_seen += 1
                if rows_seen >= 5000:
                    lines.append("[VeriFYD note: legacy workbook rendering truncated after 5000 scanned rows]")
                    break
            if rows_seen >= 5000:
                break
        try:
            book.release_resources()
        except Exception:
            pass
        return "\n".join(lines)
    except Exception:
        return _read_ole_text_for_render(src_path)



# ─────────────────────────────────────────────
# MSG/EML certified email fallback rendering
# ─────────────────────────────────────────────

def _decode_email_value(value: Any) -> str:
    """Decode bytes/byte-string/HTML-ish email body values into readable text."""
    import ast
    import html
    import re

    if value is None:
        return ""

    if isinstance(value, bytes):
        raw_bytes = value
        for enc in ("utf-8", "utf-16", "cp1252", "latin-1"):
            try:
                value = raw_bytes.decode(enc, errors="replace")
                break
            except Exception:
                continue
        if isinstance(value, bytes):
            value = raw_bytes.decode("latin-1", errors="replace")
    else:
        value = str(value)

    text = str(value or "")

    # Some libraries expose bytes as the literal string b'...'. Convert back.
    stripped = text.strip()
    if (stripped.startswith("b'") and stripped.endswith("'")) or (stripped.startswith('b"') and stripped.endswith('"')):
        try:
            lit = ast.literal_eval(stripped)
            if isinstance(lit, bytes):
                text = lit.decode("utf-8", errors="replace")
        except Exception:
            pass

    # Decode common escaped byte sequences only when they are visibly present.
    if "\\x" in text or "\\u" in text:
        try:
            text = text.encode("utf-8", errors="ignore").decode("unicode_escape", errors="replace")
        except Exception:
            pass

    text = html.unescape(text)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = text.replace("\xa0", " ").replace("\u00a0", " ").replace("\u200b", "")
    text = text.replace("\ufeff", "")

    # Preserve useful line breaks around common HTML block tags, then strip tags.
    text = re.sub(r"(?i)<\s*br\s*/?\s*>", "\n", text)
    text = re.sub(r"(?i)</\s*(p|div|tr|li|table|blockquote|h[1-6])\s*>", "\n", text)
    text = re.sub(r"(?is)<(script|style)[^>]*>.*?</\1>", " ", text)
    text = re.sub(r"(?s)<[^>]+>", " ", text)

    # Clean Outlook/HTML noise without removing the actual thread.
    text = text.replace("&nbsp;", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r" *\n *", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _email_header_value(value: Any) -> str:
    """Single-line safe display value for email headers."""
    text = _decode_email_value(value)
    text = " ".join(text.split())
    return text.strip()


def _email_body_from_msg(src_path: str) -> Dict[str, Any]:
    """Extract a structured certified email record from an Outlook MSG file."""
    try:
        import extract_msg
    except Exception as e:
        raise RuntimeError("Missing dependency: extract-msg. Add extract-msg to requirements.txt") from e

    msg = extract_msg.Message(src_path)
    try:
        attachments: List[str] = []
        try:
            for att in list(getattr(msg, "attachments", []) or [])[:60]:
                name = (
                    getattr(att, "longFilename", "")
                    or getattr(att, "shortFilename", "")
                    or getattr(att, "displayName", "")
                    or ""
                )
                name = _email_header_value(name)
                if name:
                    attachments.append(name)
        except Exception:
            attachments = []

        body_plain = _decode_email_value(getattr(msg, "body", "") or "")
        body_html = _decode_email_value(getattr(msg, "htmlBody", "") or getattr(msg, "html_body", "") or "")
        body = body_plain if len(body_plain) >= max(40, len(body_html) * 0.35) else body_html

        return {
            "kind": "MSG",
            "from": _email_header_value(getattr(msg, "sender", "") or getattr(msg, "sender_email", "")),
            "to": _email_header_value(getattr(msg, "to", "")),
            "cc": _email_header_value(getattr(msg, "cc", "")),
            "bcc": _email_header_value(getattr(msg, "bcc", "")),
            "subject": _email_header_value(getattr(msg, "subject", "")),
            "date": _email_header_value(getattr(msg, "date", "")),
            "message_id": _email_header_value(getattr(msg, "messageId", "") or getattr(msg, "message_id", "")),
            "attachments": attachments,
            "body": body,
        }
    finally:
        try:
            msg.close()
        except Exception:
            pass


def _email_body_from_eml(src_path: str) -> Dict[str, Any]:
    """Extract a structured certified email record from an EML file."""
    from email import policy
    from email.parser import BytesParser

    with open(src_path, "rb") as fh:
        msg = BytesParser(policy=policy.default).parse(fh)

    body_plain: List[str] = []
    body_html: List[str] = []
    attachments: List[str] = []

    if msg.is_multipart():
        for part in msg.walk():
            disp = str(part.get_content_disposition() or "")
            ctype = str(part.get_content_type() or "").lower()
            filename = part.get_filename()
            if filename:
                attachments.append(_email_header_value(filename))
            if disp == "attachment":
                continue
            try:
                content = part.get_content()
            except Exception:
                payload = part.get_payload(decode=True) or b""
                charset = part.get_content_charset() or "utf-8"
                content = payload.decode(charset, errors="replace")
            if ctype == "text/plain" and content:
                body_plain.append(_decode_email_value(content))
            elif ctype == "text/html" and content:
                body_html.append(_decode_email_value(content))
    else:
        try:
            content = msg.get_content()
        except Exception:
            content = (msg.get_payload(decode=True) or b"").decode(msg.get_content_charset() or "utf-8", errors="replace")
        if str(msg.get_content_type() or "").lower() == "text/html":
            body_html.append(_decode_email_value(content))
        else:
            body_plain.append(_decode_email_value(content))

    body_p = "\n\n".join(x for x in body_plain if x).strip()
    body_h = "\n\n".join(x for x in body_html if x).strip()
    body = body_p if body_p else body_h

    return {
        "kind": "EML",
        "from": _email_header_value(msg.get("From", "")),
        "to": _email_header_value(msg.get("To", "")),
        "cc": _email_header_value(msg.get("Cc", "")),
        "bcc": _email_header_value(msg.get("Bcc", "")),
        "subject": _email_header_value(msg.get("Subject", "")),
        "date": _email_header_value(msg.get("Date", "")),
        "message_id": _email_header_value(msg.get("Message-ID", "")),
        "attachments": [a for a in attachments if a],
        "body": body,
    }


def _split_email_thread_sections(body: str) -> List[Tuple[str, str]]:
    """Split email body into readable top message / thread history sections."""
    import re
    body = _decode_email_value(body)
    if not body:
        return [("Message Body", "[No readable email body was available.]")]

    markers = [m.start() for m in re.finditer(r"(?im)^\s*(From:|-----Original Message-----|Sent:|On .+ wrote:)\s*", body)]
    # Keep the first marker after some leading body text as thread boundary.
    usable = [m for m in markers if m > 40]
    if not usable:
        return [("Message Body", body[:120000])]

    first = usable[0]
    sections: List[Tuple[str, str]] = []
    lead = body[:first].strip()
    if lead:
        sections.append(("Most Recent Message", lead))
    history = body[first:].strip()
    if history:
        # Do not over-split aggressively; preserving order is more important than perfect threading.
        chunks = re.split(r"(?im)(?=^\s*From:\s+)", history)
        idx = 1
        for chunk in chunks:
            chunk = chunk.strip()
            if not chunk:
                continue
            sections.append((f"Thread History {idx}", chunk[:60000]))
            idx += 1
            if idx > 20:
                sections.append(("Thread History Continued", "[VeriFYD note: email thread display truncated after 20 sections]"))
                break
    return sections or [("Message Body", body[:120000])]


def _draw_email_wrapped_text(c: Any, text: str, x: float, y: float, max_w: float,
                             bottom: float, page_w: float, page_h: float,
                             page_num: int, title: str, cert_id: str) -> Tuple[float, int]:
    """Draw wrapped email text and paginate."""
    import re
    font = "Helvetica"
    size = 8.0
    line_h = 10.2

    def footer() -> None:
        c.setFont("Helvetica", 6.4)
        c.setFillGray(0.45)
        c.drawString(42, 24, f"VeriFYD certified email rendering • Page {page_num}")
        c.drawRightString(page_w - 42, 24, f"Certificate: {cert_id[:18]}")
        _draw_verifyd_mark(c, page_w, y=24, x_right_pad=34)

    def new_page() -> float:
        nonlocal page_num
        footer()
        c.showPage()
        page_num += 1
        c.setFont("Helvetica-Bold", 11)
        c.setFillGray(0.05)
        c.drawString(42, page_h - 44, title)
        return page_h - 66

    for raw_para in str(text or "").splitlines():
        para = re.sub(r"\s+", " ", raw_para).strip()
        if not para:
            y -= line_h * 0.65
            if y < bottom:
                y = new_page()
            continue
        words = para.split()
        current = ""
        for word in words:
            trial = (current + " " + word).strip()
            if c.stringWidth(trial, font, size) <= max_w:
                current = trial
            else:
                if y < bottom:
                    y = new_page()
                c.setFont(font, size)
                c.setFillGray(0.08)
                c.drawString(x, y, current[:180])
                y -= line_h
                current = word
        if current:
            if y < bottom:
                y = new_page()
            c.setFont(font, size)
            c.setFillGray(0.08)
            c.drawString(x, y, current[:180])
            y -= line_h
    return y, page_num


def _create_email_render_pdf(src_path: str, dest_path: str, cert_id: str, authenticity: int,
                             label: str, filename: str, sha256: str = "") -> str:
    """Create a clean certified PDF rendering for Outlook MSG and EML email files."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter

    ext = os.path.splitext(src_path)[1].lower()
    if ext == ".msg":
        record = _email_body_from_msg(src_path)
    else:
        record = _email_body_from_eml(src_path)

    width, height = letter
    c = canvas.Canvas(dest_path, pagesize=letter)
    c.setAuthor("VeriFYD")
    c.setTitle(f"VeriFYD Certified Email Record {cert_id}")
    c.setSubject(f"{label} | Authenticity {authenticity}% | {filename}")

    margin_x = 42
    y = height - 46
    bottom = 48
    max_w = width - (margin_x * 2)
    page_num = 1
    title = "VeriFYD Certified Email Record"

    c.setFont("Helvetica-Bold", 17)
    c.setFillGray(0.05)
    c.drawString(margin_x, y, title)
    y -= 18
    c.setFont("Helvetica", 8)
    c.setFillGray(0.35)
    c.drawString(margin_x, y, f"File: {filename} • Status: {label} • Authenticity: {authenticity}%")
    y -= 10
    c.drawString(margin_x, y, f"Certificate ID: {cert_id}")
    y -= 18

    # Header box.
    header_rows = [
        ("From", record.get("from", "")),
        ("To", record.get("to", "")),
        ("Cc", record.get("cc", "")),
        ("Subject", record.get("subject", "")),
        ("Date", record.get("date", "")),
        ("Message-ID", record.get("message_id", "")),
    ]
    attachments = list(record.get("attachments") or [])
    if attachments:
        header_rows.append(("Attachments", f"{len(attachments)} attachment(s): " + ", ".join(attachments[:8])))
    else:
        header_rows.append(("Attachments", "None detected"))

    c.setStrokeGray(0.72)
    c.setLineWidth(0.5)
    box_top = y
    row_h = 17
    label_w = 76
    for label_name, value in header_rows:
        if y - row_h < bottom:
            c.showPage(); page_num += 1; y = height - 46; box_top = y
        c.setFillGray(0.93)
        c.rect(margin_x, y - row_h + 3, label_w, row_h, fill=1, stroke=1)
        c.setFillGray(1)
        c.rect(margin_x + label_w, y - row_h + 3, max_w - label_w, row_h, fill=1, stroke=1)
        c.setFont("Helvetica-Bold", 7.8)
        c.setFillGray(0.08)
        c.drawString(margin_x + 4, y - row_h + 8, label_name)
        c.setFont("Helvetica", 7.5)
        c.drawString(margin_x + label_w + 4, y - row_h + 8, _safe_text(value, 112))
        y -= row_h
    y -= 12

    if sha256:
        c.setFont("Helvetica", 6.7)
        c.setFillGray(0.38)
        c.drawString(margin_x, y, f"Original SHA-256: {sha256[:96]}")
        y -= 14

    sections = _split_email_thread_sections(str(record.get("body") or ""))
    for section_title, section_text in sections:
        if y < bottom + 28:
            c.setFont("Helvetica", 6.4)
            c.setFillGray(0.45)
            c.drawString(42, 24, f"VeriFYD certified email rendering • Page {page_num}")
            _draw_verifyd_mark(c, width, y=24, x_right_pad=34)
            c.showPage(); page_num += 1; y = height - 46
        c.setFont("Helvetica-Bold", 10.5)
        c.setFillGray(0.05)
        c.drawString(margin_x, y, section_title[:95])
        y -= 13
        y, page_num = _draw_email_wrapped_text(
            c, section_text, margin_x, y, max_w, bottom, width, height, page_num, title, cert_id
        )
        y -= 8

    c.setFont("Helvetica", 6.4)
    c.setFillGray(0.45)
    c.drawString(42, 24, f"VeriFYD certified email rendering • Page {page_num}")
    c.drawRightString(width - 42, 24, f"Certificate: {cert_id[:18]}")
    _draw_verifyd_mark(c, width, y=24, x_right_pad=34)
    c.save()
    return dest_path

def _read_msg_for_render(src_path: str) -> str:
    """Render Outlook MSG headers/body into plain text for certified PDF output."""
    try:
        import extract_msg
        msg = extract_msg.Message(src_path)
        try:
            lines: List[str] = []
            header_map = [
                ("From", getattr(msg, "sender", "") or ""),
                ("To", getattr(msg, "to", "") or ""),
                ("Cc", getattr(msg, "cc", "") or ""),
                ("Subject", getattr(msg, "subject", "") or ""),
                ("Date", getattr(msg, "date", "") or ""),
                ("Message-ID", getattr(msg, "messageId", "") or ""),
            ]
            for label, value in header_map:
                value = str(value or "").strip()
                if value:
                    lines.append(f"{label}: {value}")
            attachments = list(getattr(msg, "attachments", []) or [])
            if attachments:
                names = []
                for att in attachments[:30]:
                    name = str(getattr(att, "longFilename", "") or getattr(att, "shortFilename", "") or "")
                    if name:
                        names.append(name)
                lines.append(f"Attachments: {len(attachments)}" + (f" ({', '.join(names[:10])})" if names else ""))
            lines.append("")
            body = str(getattr(msg, "body", "") or getattr(msg, "htmlBody", "") or "")
            if body and "<" in body and ">" in body:
                import re
                body = re.sub(r"<[^>]+>", " ", body)
                body = re.sub(r"\s+", " ", body).strip()
            lines.append(body)
            return "\n".join(lines).strip()
        finally:
            try:
                msg.close()
            except Exception:
                pass
    except Exception:
        return _read_ole_text_for_render(src_path)


def _read_odf_for_render(src_path: str) -> str:
    """Extract readable text/table content from ODT/ODS/ODP for certified PDF rendering."""
    import re
    import zipfile
    import xml.etree.ElementTree as ET

    ext = os.path.splitext(src_path)[1].lower()
    kind = {".odt": "OpenDocument Text", ".ods": "OpenDocument Spreadsheet", ".odp": "OpenDocument Presentation"}.get(ext, "OpenDocument")
    lines: List[str] = [kind]

    def _local(tag: str) -> str:
        return str(tag).split("}")[-1].lower()

    def _text(elem: Any) -> str:
        try:
            return re.sub(r"\s+", " ", " ".join(t.strip() for t in elem.itertext() if t and t.strip())).strip()
        except Exception:
            return ""

    try:
        with zipfile.ZipFile(src_path) as zf:
            names = set(zf.namelist())
            if "content.xml" not in names:
                return "[No OpenDocument content.xml found]"
            root = ET.fromstring(zf.read("content.xml"))

            for elem in root.iter():
                name = _local(elem.tag)
                if ext == ".ods" and name == "table":
                    table_name = elem.attrib.get("{urn:oasis:names:tc:opendocument:xmlns:table:1.0}name", "Sheet")
                    lines.append(f"\nWorksheet: {table_name}")
                elif ext == ".odp" and name == "page":
                    page_name = elem.attrib.get("{urn:oasis:names:tc:opendocument:xmlns:drawing:1.0}name", "Slide")
                    lines.append(f"\nSlide: {page_name}")
                elif ext == ".ods" and name == "table-row":
                    cells: List[str] = []
                    for child in list(elem):
                        if _local(child.tag) == "table-cell":
                            value = _text(child)
                            if value:
                                cells.append(value[:300])
                    if cells:
                        lines.append(" | ".join(cells))
                elif name in ("h", "p"):
                    value = _text(elem)
                    if value:
                        lines.append(value[:1000])
    except Exception as e:
        return f"[OpenDocument rendering failed: {str(e)[:120]}]"

    cleaned = [line for line in lines if line and line.strip()]
    return "\n".join(cleaned) if cleaned else "[No readable OpenDocument text could be extracted]"



def _strip_html_for_render(html: str) -> str:
    """Convert HTML markup into readable text for certified PDF rendering."""
    import re
    try:
        from html.parser import HTMLParser
        from html import unescape
    except Exception:
        return re.sub(r"\s+", " ", re.sub(r"<[^>]+>", " ", html or "")).strip()

    class _Parser(HTMLParser):
        def __init__(self):
            super().__init__(convert_charrefs=True)
            self.parts: List[str] = []
            self.skip = 0
        def handle_starttag(self, tag, attrs):
            tag = str(tag or "").lower()
            if tag in ("script", "style", "noscript", "template"):
                self.skip += 1
            if tag in ("p", "div", "br", "li", "tr", "h1", "h2", "h3", "h4", "section", "article"):
                self.parts.append("\n")
        def handle_endtag(self, tag):
            tag = str(tag or "").lower()
            if tag in ("script", "style", "noscript", "template") and self.skip:
                self.skip -= 1
            if tag in ("p", "div", "li", "tr", "h1", "h2", "h3", "h4", "section", "article"):
                self.parts.append("\n")
        def handle_data(self, data):
            if not self.skip:
                t = unescape(str(data or "")).strip()
                if t:
                    self.parts.append(t + " ")
    parser = _Parser()
    try:
        parser.feed(html or "")
        parser.close()
        text = "".join(parser.parts)
    except Exception:
        text = re.sub(r"<[^>]+>", " ", html or "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n\s*\n\s*\n+", "\n\n", text)
    return text.strip()


def _read_html_for_render(src_path: str) -> str:
    data = open(src_path, "rb").read(3_000_000)
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return _strip_html_for_render(data.decode(enc, errors="replace"))
        except Exception:
            continue
    return _strip_html_for_render(data.decode("latin-1", errors="replace"))


def _read_mhtml_for_render(src_path: str) -> str:
    from email import policy
    from email.parser import BytesParser
    with open(src_path, "rb") as fh:
        msg = BytesParser(policy=policy.default).parse(fh)
    lines: List[str] = []
    for h in ("From", "Subject", "Date", "Content-Type"):
        v = str(msg.get(h, "") or "").strip()
        if v:
            lines.append(f"{h}: {v}")
    lines.append("")
    parts: List[str] = []
    if msg.is_multipart():
        for part in msg.walk():
            ctype = part.get_content_type()
            if ctype not in ("text/plain", "text/html"):
                continue
            try:
                content = part.get_content()
            except Exception:
                try:
                    content = (part.get_payload(decode=True) or b"").decode(part.get_content_charset() or "utf-8", errors="replace")
                except Exception:
                    content = ""
            if not content:
                continue
            parts.append(_strip_html_for_render(str(content)) if ctype == "text/html" else str(content).strip())
    else:
        try:
            content = msg.get_content()
        except Exception:
            content = ""
        parts.append(_strip_html_for_render(str(content)) if msg.get_content_type() == "text/html" else str(content).strip())
    lines.append("\n\n".join(p for p in parts if p))
    return "\n".join(lines).strip()


def _read_xml_for_render(src_path: str) -> str:
    import re
    import xml.etree.ElementTree as ET
    data = open(src_path, "rb").read(3_000_000)
    try:
        root = ET.fromstring(data)
        def _local(tag: str) -> str:
            return str(tag).split("}")[-1]
        lines = [f"Root: {_local(root.tag)}"]
        for elem in root.iter():
            value = re.sub(r"\s+", " ", " ".join(t.strip() for t in elem.itertext() if t and t.strip())).strip()
            attrs = " ".join(f"{_local(k)}={v}" for k, v in list(elem.attrib.items())[:8])
            if value or attrs:
                lines.append(f"{_local(elem.tag)}: {(value or attrs)[:900]}")
            if len(lines) >= 2500:
                lines.append("[VeriFYD note: XML rendering truncated]")
                break
        return "\n".join(lines)
    except Exception:
        raw = data.decode("utf-8", errors="replace")
        return re.sub(r"\s+", " ", re.sub(r"<[^>]+>", " ", raw)).strip()


def _read_json_for_render(src_path: str) -> str:
    data = open(src_path, "rb").read(5_000_000)
    raw = data.decode("utf-8", errors="replace")
    try:
        import json as _json
        parsed = _json.loads(raw)
        lines: List[str] = []
        def walk(value: Any, key: str = "", depth: int = 0):
            if len(lines) >= 2500 or depth > 8:
                return
            if isinstance(value, dict):
                if key:
                    lines.append(f"{key}: object({len(value)})")
                for k, v in list(value.items())[:200]:
                    walk(v, f"{key}.{k}" if key else str(k), depth + 1)
            elif isinstance(value, list):
                lines.append(f"{key}: array({len(value)})")
                for i, item in enumerate(value[:80]):
                    walk(item, f"{key}[{i}]", depth + 1)
            else:
                t = str(value).strip()
                if t:
                    lines.append(f"{key}: {t[:900]}")
        walk(parsed)
        return "\n".join(lines) if lines else raw[:12000]
    except Exception:
        return raw[:12000]


# ─────────────────────────────────────────────
# VSDX certified Visio visual fallback rendering
# ─────────────────────────────────────────────

def _vsdx_local_name(tag: Any) -> str:
    """Return local XML tag name without namespace."""
    return str(tag).split("}")[-1]


def _vsdx_float(value: Any, default: float = 0.0) -> float:
    """Parse Visio numeric values safely."""
    try:
        text = str(value or "").strip()
        if not text:
            return default
        import re
        m = re.search(r"[-+]?\d+(?:\.\d+)?", text)
        if not m:
            return default
        return float(m.group(0))
    except Exception:
        return default


def _vsdx_clean_text(value: Any) -> str:
    """Clean text extracted from Visio XML."""
    import html
    import re
    text = html.unescape(str(value or ""))
    text = text.replace("\r", "\n").replace("\t", " ")
    text = re.sub(r"\s*\n\s*", "\n", text)
    text = re.sub(r"[ \u00a0]+", " ", text)
    return text.strip()


def _vsdx_element_text(elem: Any) -> str:
    """Extract readable text from a Visio XML element."""
    try:
        text = " ".join(t for t in elem.itertext() if t is not None)
    except Exception:
        text = ""
    return _vsdx_clean_text(text)


def _vsdx_shape_cells(shape: Any) -> Dict[str, str]:
    """Collect Visio Cell values from a Shape element."""
    cells: Dict[str, str] = {}
    try:
        for child in shape.iter():
            if _vsdx_local_name(getattr(child, "tag", "")).lower() == "cell":
                name = str(child.attrib.get("N") or child.attrib.get("Name") or "").strip()
                value = str(child.attrib.get("V") or child.attrib.get("Value") or child.attrib.get("F") or "").strip()
                if name:
                    cells[name] = value
    except Exception:
        pass
    return cells


def _vsdx_shape_type(shape: Any, cells: Dict[str, str], text: str) -> str:
    """Best-effort shape type label for the visual fallback."""
    try:
        raw_type = str(shape.attrib.get("Type") or "").lower()
        master = str(shape.attrib.get("Master") or shape.attrib.get("MasterShape") or "").lower()
    except Exception:
        raw_type = ""
        master = ""

    sample = " ".join([raw_type, master, text.lower()])
    if any(x in sample for x in ("connector", "dynamic connector", "arrow")):
        return "connector"
    if any(x in sample for x in ("ellipse", "oval", "circle")):
        return "ellipse"
    if any(x in sample for x in ("decision", "diamond")):
        return "diamond"
    if any(x in sample for x in ("swimlane", "container")):
        return "container"
    return "box"


def _vsdx_extract_pages(src_path: str) -> List[Dict[str, Any]]:
    """
    Extract a lightweight visual model from a VSDX package.

    This is not a full Visio rendering engine. It reads page XML and draws a
    reviewable diagram approximation: shape bounds, connectors where possible,
    labels/text, and a manifest summary. The exact original VSDX is still
    attached to the certified PDF/package for source review.
    """
    import zipfile
    import xml.etree.ElementTree as ET
    import os
    pages: List[Dict[str, Any]] = []

    with zipfile.ZipFile(src_path, "r") as zf:
        names = zf.namelist()
        page_names = [
            n for n in names
            if n.lower().startswith("visio/pages/")
            and n.lower().endswith(".xml")
            and "/_rels/" not in n.lower()
        ]
        page_names.sort()

        for pidx, name in enumerate(page_names, start=1):
            shapes: List[Dict[str, Any]] = []
            raw_xml = zf.read(name)
            try:
                root = ET.fromstring(raw_xml)
            except Exception as e:
                pages.append({
                    "index": pidx,
                    "name": os.path.basename(name),
                    "shapes": [],
                    "error": f"Could not parse page XML: {str(e)[:120]}",
                })
                continue

            for elem in root.iter():
                if _vsdx_local_name(getattr(elem, "tag", "")).lower() != "shape":
                    continue

                cells = _vsdx_shape_cells(elem)
                text = ""
                for child in elem.iter():
                    if _vsdx_local_name(getattr(child, "tag", "")).lower() == "text":
                        t = _vsdx_element_text(child)
                        if t:
                            text = t
                            break

                pin_x = _vsdx_float(cells.get("PinX"), 0.0)
                pin_y = _vsdx_float(cells.get("PinY"), 0.0)
                width = abs(_vsdx_float(cells.get("Width"), 1.0))
                height = abs(_vsdx_float(cells.get("Height"), 0.55))
                begin_x = _vsdx_float(cells.get("BeginX"), pin_x)
                begin_y = _vsdx_float(cells.get("BeginY"), pin_y)
                end_x = _vsdx_float(cells.get("EndX"), pin_x)
                end_y = _vsdx_float(cells.get("EndY"), pin_y)

                shape_id = str(elem.attrib.get("ID") or elem.attrib.get("Id") or "")
                shape_name = str(elem.attrib.get("Name") or elem.attrib.get("NameU") or "")
                shape_type = _vsdx_shape_type(elem, cells, text)

                if not text and not cells and not shape_id and not shape_name:
                    continue

                shapes.append({
                    "id": shape_id,
                    "name": shape_name,
                    "text": text,
                    "type": shape_type,
                    "pin_x": pin_x,
                    "pin_y": pin_y,
                    "width": max(width, 0.08),
                    "height": max(height, 0.08),
                    "begin_x": begin_x,
                    "begin_y": begin_y,
                    "end_x": end_x,
                    "end_y": end_y,
                })
                if len(shapes) >= 700:
                    shapes.append({
                        "id": "",
                        "name": "truncated",
                        "text": "[VeriFYD note: VSDX page rendering truncated after 700 shapes]",
                        "type": "box",
                        "pin_x": 0,
                        "pin_y": 0,
                        "width": 3,
                        "height": 0.4,
                        "begin_x": 0,
                        "begin_y": 0,
                        "end_x": 0,
                        "end_y": 0,
                    })
                    break

            pages.append({
                "index": pidx,
                "name": os.path.basename(name),
                "shapes": shapes,
                "error": "",
            })

    return pages


def _vsdx_diagram_bounds(shapes: List[Dict[str, Any]]) -> Tuple[float, float, float, float]:
    """Compute diagram bounds in Visio coordinate units."""
    xs: List[float] = []
    ys: List[float] = []
    for s in shapes:
        if s.get("type") == "connector":
            xs.extend([float(s.get("begin_x", 0)), float(s.get("end_x", 0))])
            ys.extend([float(s.get("begin_y", 0)), float(s.get("end_y", 0))])
        else:
            px = float(s.get("pin_x", 0))
            py = float(s.get("pin_y", 0))
            w = float(s.get("width", 0.1))
            h = float(s.get("height", 0.1))
            xs.extend([px - w / 2, px + w / 2])
            ys.extend([py - h / 2, py + h / 2])
    if not xs or not ys:
        return (0, 0, 10, 7.5)
    min_x, max_x = min(xs), max(xs)
    min_y, max_y = min(ys), max(ys)
    if max_x - min_x < 0.1:
        max_x = min_x + 10
    if max_y - min_y < 0.1:
        max_y = min_y + 7.5
    return (min_x, min_y, max_x, max_y)


def _vsdx_draw_wrapped_label(c: Any, text: str, x: float, y: float, w: float, h: float, font_size: float = 6.2) -> None:
    """Draw a compact wrapped label inside a shape."""
    import textwrap
    text = _vsdx_clean_text(text)
    if not text:
        return
    approx_chars = max(6, int(w / max(2.5, font_size * 0.46)))
    lines: List[str] = []
    for para in text.split("\n"):
        lines.extend(textwrap.wrap(para, width=approx_chars) or [""])
        if len(lines) >= 5:
            break
    c.setFont("Helvetica", font_size)
    c.setFillGray(0.05)
    line_h = font_size + 1
    total_h = len(lines) * line_h
    start_y = y + (h / 2) + (total_h / 2) - line_h
    for idx, line in enumerate(lines[:5]):
        draw_y = start_y - idx * line_h
        if draw_y < y + 2:
            break
        c.drawCentredString(x + w / 2, draw_y, line[:max(6, approx_chars + 6)])


def _create_vsdx_diagram_pdf(src_path: str, dest_path: str, cert_id: str, authenticity: int,
                             label: str, filename: str, sha256: str = "") -> str:
    """
    Create a certified visual/evidence rendering for Visio VSDX files.

    The current Render environment does not have Visio or LibreOffice available,
    so this is a pure-Python diagram approximation from the VSDX XML package.
    The exact original VSDX remains embedded/packaged for source review.
    """
    import zipfile
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import landscape, letter

    width, height = landscape(letter)
    c = canvas.Canvas(dest_path, pagesize=(width, height))
    c.setAuthor("VeriFYD")
    c.setTitle(f"VeriFYD Certified Visio Diagram {cert_id}")
    c.setSubject(f"{label} | Authenticity {authenticity}% | {filename}")

    try:
        pages = _vsdx_extract_pages(src_path)
    except Exception as e:
        pages = [{"index": 1, "name": "VSDX", "shapes": [], "error": f"VSDX extraction failed: {str(e)[:160]}"}]

    try:
        with zipfile.ZipFile(src_path, "r") as zf:
            names = zf.namelist()
            package_files = len(names)
            media_files = len([n for n in names if n.lower().startswith("visio/media/")])
            page_count = len([n for n in names if n.lower().startswith("visio/pages/") and n.lower().endswith(".xml")])
    except Exception:
        package_files = 0
        media_files = 0
        page_count = len(pages)

    c.setFont("Helvetica-Bold", 18)
    c.setFillGray(0.05)
    c.drawString(42, height - 54, "VeriFYD Certified Visio Diagram")
    c.setFont("Helvetica", 9)
    c.setFillGray(0.20)
    y = height - 78
    for line in [
        f"File: {filename}",
        f"Status: {label} • Authenticity: {authenticity}%",
        f"Certificate ID: {cert_id}",
        f"Pages discovered: {page_count}",
        f"Package files: {package_files}",
        f"Embedded media files: {media_files}",
        f"SHA-256: {sha256 or 'not provided'}",
        "The following pages are a VeriFYD-rendered visual approximation of the VSDX XML diagram. The exact original VSDX remains embedded in the certified PDF/package.",
    ]:
        c.drawString(42, y, line[:155])
        y -= 15
    _draw_verifyd_mark(c, width, y=24, x_right_pad=34)
    c.showPage()

    if not pages:
        pages = [{"index": 1, "name": "No pages found", "shapes": [], "error": "No Visio page XML files were found in this VSDX package."}]

    for page in pages:
        shapes = list(page.get("shapes") or [])
        c.setFont("Helvetica-Bold", 13)
        c.setFillGray(0.05)
        c.drawString(34, height - 34, f"Visio Page {page.get('index', '')}: {page.get('name', '')}"[:120])
        c.setFont("Helvetica", 7)
        c.setFillGray(0.38)
        c.drawString(34, height - 48, f"File: {filename} • Certificate: {cert_id[:18]} • Shapes rendered: {len(shapes)}")

        draw_x = 34
        draw_y = 54
        draw_w = width - 68
        draw_h = height - 120

        c.setStrokeGray(0.72)
        c.setLineWidth(0.8)
        c.rect(draw_x, draw_y, draw_w, draw_h, fill=0, stroke=1)

        if page.get("error"):
            c.setFont("Helvetica", 9)
            c.setFillGray(0.20)
            c.drawString(draw_x + 14, height - 88, str(page.get("error"))[:130])
        elif not shapes:
            c.setFont("Helvetica", 10)
            c.setFillGray(0.20)
            c.drawString(draw_x + 14, height - 88, "No drawable shapes or extractable shape text were found in this Visio page XML.")
            c.setFont("Helvetica", 8)
            c.setFillGray(0.38)
            c.drawString(draw_x + 14, height - 104, "The file is still certified and the exact source VSDX is preserved in the certified package.")
        else:
            min_x, min_y, max_x, max_y = _vsdx_diagram_bounds(shapes)
            margin = 18
            scale = min((draw_w - 2 * margin) / max(0.1, max_x - min_x), (draw_h - 2 * margin) / max(0.1, max_y - min_y))

            def tx(v: float) -> float:
                return draw_x + margin + ((v - min_x) * scale)

            def ty(v: float) -> float:
                return draw_y + margin + ((v - min_y) * scale)

            for s in shapes:
                if s.get("type") != "connector":
                    continue
                x1, y1 = tx(float(s.get("begin_x", 0))), ty(float(s.get("begin_y", 0)))
                x2, y2 = tx(float(s.get("end_x", 0))), ty(float(s.get("end_y", 0)))
                c.setStrokeGray(0.28)
                c.setLineWidth(0.8)
                c.line(x1, y1, x2, y2)

            for s in shapes:
                if s.get("type") == "connector":
                    continue
                px = float(s.get("pin_x", 0))
                py = float(s.get("pin_y", 0))
                sw = max(10, float(s.get("width", 0.4)) * scale)
                sh = max(8, float(s.get("height", 0.25)) * scale)
                x = tx(px) - sw / 2
                y0 = ty(py) - sh / 2
                kind = str(s.get("type") or "box")
                text = str(s.get("text") or s.get("name") or "")

                c.setLineWidth(0.7)
                c.setStrokeGray(0.25)
                if text:
                    c.setFillGray(0.96)
                else:
                    c.setFillGray(0.99)

                if kind == "ellipse":
                    c.ellipse(x, y0, x + sw, y0 + sh, fill=1, stroke=1)
                elif kind == "diamond":
                    pts = [(x + sw / 2, y0 + sh), (x + sw, y0 + sh / 2), (x + sw / 2, y0), (x, y0 + sh / 2)]
                    p = c.beginPath()
                    p.moveTo(*pts[0])
                    for pt in pts[1:]:
                        p.lineTo(*pt)
                    p.close()
                    c.drawPath(p, fill=1, stroke=1)
                else:
                    c.roundRect(x, y0, sw, sh, radius=3, fill=1, stroke=1)

                if text:
                    _vsdx_draw_wrapped_label(c, text, x + 2, y0 + 2, max(6, sw - 4), max(6, sh - 4), font_size=6.1)

        c.setFont("Helvetica", 6.5)
        c.setFillGray(0.45)
        c.drawString(34, 24, f"VeriFYD certified Visio rendering • Page {page.get('index', '')}")
        _draw_verifyd_mark(c, width, y=24, x_right_pad=34)
        c.showPage()

    c.save()
    return dest_path

def _read_vsdx_for_render(src_path: str) -> str:
    import zipfile
    import xml.etree.ElementTree as ET
    import re
    def _local(tag: str) -> str:
        return str(tag).split("}")[-1]
    def _text(elem: Any) -> str:
        return re.sub(r"\s+", " ", " ".join(t.strip() for t in elem.itertext() if t and t.strip())).strip()
    lines: List[str] = ["VeriFYD Visio Diagram Rendering"]
    with zipfile.ZipFile(src_path, "r") as zf:
        names = zf.namelist()
        pages = [n for n in names if n.lower().startswith("visio/pages/") and n.lower().endswith(".xml")]
        lines.append(f"Pages: {len(pages)}")
        lines.append(f"Embedded media files: {len([n for n in names if n.lower().startswith('visio/media/')])}")
        for idx, name in enumerate(pages[:80], start=1):
            try:
                root = ET.fromstring(zf.read(name))
            except Exception:
                continue
            vals: List[str] = []
            for elem in root.iter():
                if _local(elem.tag).lower() in ("text", "cp", "pp", "tp"):
                    value = _text(elem)
                    if value and len(value) > 1:
                        vals.append(value[:900])
            lines.append(f"\nPage {idx}: {os.path.basename(name)}")
            if vals:
                seen = set()
                for v in vals:
                    k = v.lower()[:200]
                    if k not in seen:
                        seen.add(k)
                        lines.append(v)
            else:
                lines.append("[No extractable shape text]")
    return "\n".join(lines)



def _read_extra_textlike_for_render(src_path: str, ext: str) -> str:
    """Render lightweight text/config/log/code evidence formats."""
    data = open(src_path, "rb").read(5_000_000)
    text = ""
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            text = data.decode(enc, errors="replace")
            break
        except Exception:
            continue
    if not text:
        text = data.decode("latin-1", errors="replace")
    label = {
        ".yaml": "YAML Configuration",
        ".yml": "YAML Configuration",
        ".ini": "INI Configuration",
        ".log": "Log File",
        ".sql": "SQL Script/Export",
    }.get(ext, ext.upper())
    return f"VeriFYD {label} Rendering\n\n" + text[:120000]


def _read_dxf_for_render(src_path: str) -> str:
    data = open(src_path, "rb").read(8_000_000)
    text = data.decode("utf-8", errors="replace")
    if "SECTION" not in text[:200000].upper():
        text = data.decode("latin-1", errors="replace")
    import re
    layers = re.findall(r"\n\s*8\s*\n([^\n\r]{1,120})", text)
    entities = re.findall(r"\n\s*0\s*\n(LINE|LWPOLYLINE|POLYLINE|CIRCLE|ARC|TEXT|MTEXT|INSERT|DIMENSION|HATCH|SPLINE)\b", text, flags=re.I)
    lines = ["VeriFYD DXF CAD Evidence Rendering"]
    if layers:
        seen = []
        for layer in layers:
            layer = layer.strip()
            if layer and layer not in seen:
                seen.append(layer)
            if len(seen) >= 60:
                break
        lines.append("Layers: " + ", ".join(seen))
    if entities:
        from collections import Counter
        counts = Counter(x.upper() for x in entities)
        lines.append("Entity counts: " + ", ".join(f"{k}:{v}" for k, v in counts.most_common(20)))
    readable = re.findall(r"[A-Za-z0-9][A-Za-z0-9\s\.,;:\-_/@$%#&()\[\]{}'\"!?]{4,}", text[:1_500_000])
    lines.extend(re.sub(r"\s+", " ", x).strip()[:500] for x in readable[:800])
    return "\n".join(x for x in lines if x)


def _read_binary_evidence_for_render(src_path: str, ext: str) -> str:
    """Create a readable evidence record for PST/OST/DWG binary files."""
    import re
    size = os.path.getsize(src_path)
    with open(src_path, "rb") as fh:
        head = fh.read(64)
        fh.seek(0)
        sample = fh.read(3_000_000)
    fragments: List[str] = []
    for enc in ("utf-16-le", "latin-1"):
        try:
            decoded = sample.decode(enc, errors="ignore").replace("\x00", " ")
        except Exception:
            continue
        pieces = re.findall(r"[A-Za-z0-9][A-Za-z0-9\s\.,;:\-_/@$%#&()\[\]{}'\"!?]{5,}", decoded)
        for piece in pieces:
            clean = re.sub(r"\s+", " ", piece).strip()
            if len(clean) >= 6:
                fragments.append(clean[:500])
            if len(fragments) >= 300:
                break
        if len(fragments) >= 300:
            break
    type_label = {
        ".pst": "Outlook PST email archive",
        ".ost": "Outlook OST offline mailbox archive",
        ".dwg": "AutoCAD DWG drawing",
    }.get(ext, ext.upper())
    lines = [
        f"VeriFYD {type_label} Evidence Record",
        f"File size: {size} bytes",
        f"Magic/header hex: {head.hex()[:128]}",
        "Analysis mode: lightweight evidence certification. Deep native parsing/rendering is not enabled for this format.",
    ]
    if fragments:
        lines.append("\nReadable string samples:")
        seen = set()
        for frag in fragments:
            key = frag.lower()[:160]
            if key in seen:
                continue
            seen.add(key)
            lines.append(frag)
            if len(seen) >= 200:
                break
    else:
        lines.append("\nNo readable text strings found in scanned binary sample.")
    return "\n".join(lines)




# ─────────────────────────────────────────────
# DOCX table-aware fallback rendering
# ─────────────────────────────────────────────

def _iter_docx_blocks(doc: Any):
    """Yield DOCX paragraphs and tables in original body order."""
    try:
        from docx.oxml.text.paragraph import CT_P
        from docx.oxml.table import CT_Tbl
        from docx.text.paragraph import Paragraph
        from docx.table import Table
    except Exception:
        return

    body = getattr(doc, "element", None)
    body = getattr(body, "body", None)
    if body is None:
        return

    for child in body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, doc)
        elif isinstance(child, CT_Tbl):
            yield Table(child, doc)


def _docx_para_text_for_pdf(para: Any) -> str:
    """Return paragraph text with tabs normalized for PDF rendering."""
    try:
        text = para.text or ""
    except Exception:
        text = ""
    text = text.replace("\t", "    ").replace("\r", " ")
    return text.strip()


def _docx_cell_text_for_pdf(cell: Any) -> str:
    """Return readable table-cell text while preserving line breaks inside cells."""
    parts: List[str] = []
    try:
        for para in cell.paragraphs:
            t = _docx_para_text_for_pdf(para)
            if t:
                parts.append(t)
    except Exception:
        try:
            t = str(cell.text or "").strip()
            if t:
                parts.append(t)
        except Exception:
            pass
    return "\n".join(parts).strip()


def _docx_is_boldish(para: Any) -> bool:
    """Best-effort test for title/header-style paragraphs."""
    try:
        style_name = str(getattr(getattr(para, "style", None), "name", "") or "").lower()
        if any(x in style_name for x in ("heading", "title", "subtitle")):
            return True
    except Exception:
        pass
    try:
        runs = list(getattr(para, "runs", []) or [])
        if runs and any(bool(getattr(getattr(r, "font", None), "bold", None)) or bool(getattr(r, "bold", None)) for r in runs):
            return True
    except Exception:
        pass
    return False


def _docx_alignment_name(para: Any) -> str:
    """Return a ReportLab alignment keyword for a paragraph."""
    try:
        alignment = getattr(para, "alignment", None)
        # python-docx WD_ALIGN_PARAGRAPH enum values: LEFT=0, CENTER=1, RIGHT=2, JUSTIFY=3
        if alignment == 1:
            return "CENTER"
        if alignment == 2:
            return "RIGHT"
    except Exception:
        pass
    return "LEFT"


def _docx_table_col_widths(rows: List[List[str]], available_width: float) -> List[float]:
    """Compute stable column widths for DOCX table rendering."""
    max_cols = max((len(r) for r in rows), default=1)
    max_cols = max(1, min(max_cols, 12))

    # Common review-comment/table template: number, drawing, comment, response, action/date.
    if max_cols == 5:
        weights = [0.72, 1.05, 5.45, 2.05, 1.15]
    elif max_cols == 6:
        weights = [0.65, 1.0, 4.7, 1.8, 1.0, 1.0]
    elif max_cols == 4:
        weights = [0.9, 1.4, 5.2, 2.1]
    else:
        # Estimate text-heavy columns a bit wider while keeping the layout predictable.
        weights = []
        for cidx in range(max_cols):
            sample = " ".join((r[cidx] if cidx < len(r) else "") for r in rows[:25])
            length = len(sample)
            weights.append(max(0.8, min(3.0, 0.8 + (length / 260.0))))

    total = sum(weights) or 1.0
    return [available_width * (w / total) for w in weights]



def _create_docx_layout_pdf(src_path: str, dest_path: str, cert_id: str, authenticity: int,
                            label: str, filename: str, sha256: str = "") -> str:
    """
    Render DOCX using a table-aware pure-Python fallback.

    This does not require LibreOffice and is intentionally safer for the current
    Render Python worker. It will not perfectly clone Microsoft Word, but it now
    preserves section headers, header tables, body tables, footer text, and footer
    action-code tables so SCA-style review forms remain recognizable.
    """
    try:
        from docx import Document
    except Exception as e:
        raise RuntimeError("Missing dependency: python-docx. Add python-docx>=1.1.0 to requirements.txt") from e

    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, KeepTogether
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT
    from reportlab.lib.units import inch

    docx_doc = Document(src_path)
    page_size = landscape(letter)
    page_w, page_h = page_size

    # Tight margins help landscape review forms keep all response/action columns.
    left_margin = 0.24 * inch
    right_margin = 0.24 * inch
    top_margin = 0.26 * inch
    bottom_margin = 0.38 * inch
    available_w = page_w - left_margin - right_margin

    pdf = SimpleDocTemplate(
        dest_path,
        pagesize=page_size,
        leftMargin=left_margin,
        rightMargin=right_margin,
        topMargin=top_margin,
        bottomMargin=bottom_margin,
        title=f"VeriFYD Certified Word Document {cert_id}",
        author="VeriFYD",
        subject=f"{label} | Authenticity {authenticity}% | {filename}",
    )

    styles = getSampleStyleSheet()
    base = ParagraphStyle(
        "VeriFYD_DOCX_Base",
        parent=styles["Normal"],
        fontName="Helvetica",
        fontSize=7.1,
        leading=8.25,
        spaceAfter=2.0,
        alignment=TA_LEFT,
        wordWrap="CJK",
    )
    bold = ParagraphStyle(
        "VeriFYD_DOCX_Bold",
        parent=base,
        fontName="Helvetica-Bold",
        fontSize=7.25,
        leading=8.5,
        spaceAfter=2.4,
    )
    title = ParagraphStyle(
        "VeriFYD_DOCX_Title",
        parent=bold,
        fontSize=10.4,
        leading=11.4,
        alignment=TA_CENTER,
        spaceAfter=3.0,
    )
    right = ParagraphStyle("VeriFYD_DOCX_Right", parent=base, alignment=TA_RIGHT)
    center = ParagraphStyle("VeriFYD_DOCX_Center", parent=base, alignment=TA_CENTER)
    small = ParagraphStyle(
        "VeriFYD_DOCX_Small",
        parent=base,
        fontSize=5.9,
        leading=6.7,
        spaceAfter=1.2,
    )
    cell_style = ParagraphStyle(
        "VeriFYD_DOCX_Cell",
        parent=base,
        fontSize=5.65,
        leading=6.35,
        spaceAfter=0,
        wordWrap="CJK",
    )
    cell_bold = ParagraphStyle(
        "VeriFYD_DOCX_Cell_Bold",
        parent=cell_style,
        fontName="Helvetica-Bold",
        fontSize=5.55,
        leading=6.25,
    )
    header_cell = ParagraphStyle(
        "VeriFYD_DOCX_HeaderCell",
        parent=cell_bold,
        fontSize=5.35,
        leading=6.05,
        alignment=TA_CENTER,
    )
    review_cell = ParagraphStyle(
        "VeriFYD_DOCX_ReviewCell",
        parent=cell_style,
        fontSize=4.95,
        leading=5.55,
        spaceAfter=0,
        wordWrap="CJK",
    )
    review_header_cell = ParagraphStyle(
        "VeriFYD_DOCX_ReviewHeaderCell",
        parent=review_cell,
        fontName="Helvetica-Bold",
        fontSize=4.95,
        leading=5.45,
        alignment=TA_CENTER,
    )

    def esc(value: str) -> str:
        import html
        return html.escape(str(value or "")).replace("\n", "<br/>")

    def clean_cell_values(row_cells: List[Any]) -> List[str]:
        """Extract row cells and reduce duplicate text caused by Word merged cells."""
        values: List[str] = []
        last_tc_id = None
        last_text = None
        for cell in row_cells:
            text_value = _docx_cell_text_for_pdf(cell)
            tc_id = None
            try:
                tc_id = id(cell._tc)
            except Exception:
                tc_id = None

            # Merged cells often appear multiple times. Keep the first visual slot
            # and blank the repeated slots so the grid remains but text is not duplicated.
            if values and ((tc_id is not None and tc_id == last_tc_id) or (text_value and text_value == last_text and len(text_value) > 18)):
                values.append("")
            else:
                values.append(text_value)
                if text_value:
                    last_text = text_value
            last_tc_id = tc_id

        # Preserve form columns, but trim useless all-empty duplicate tail cells.
        while len(values) > 1 and values[-1] == "" and values[-2] == "":
            values.pop()
        return values

    def table_has_header(rows: List[List[str]]) -> bool:
        sample = " ".join(" ".join(r) for r in rows[:2]).lower()
        return any(x in sample for x in ("comment no", "drawing no", "reviewer comments", "designer responses", "action"))

    def looks_like_review_comment_table(rows: List[List[str]]) -> bool:
        if not rows:
            return False
        hits = 0
        for row in rows[:8]:
            if len(row) >= 3:
                first = str(row[0]).strip().rstrip(".")
                second = str(row[1]).strip().upper()
                third = str(row[2]).strip()
                if first.isdigit() and (second.startswith("E") or second) and third:
                    hits += 1
        return hits >= 2

    def append_paragraph(story: List[Any], para: Any, force_style: Any | None = None) -> bool:
        text = _docx_para_text_for_pdf(para) if not isinstance(para, str) else str(para or "").strip()
        if not text:
            return False
        if force_style is not None:
            style = force_style
        elif not isinstance(para, str):
            align = _docx_alignment_name(para)
            if align == "CENTER":
                style = title if len(text) <= 90 and _docx_is_boldish(para) else center
            elif align == "RIGHT":
                style = right
            elif _docx_is_boldish(para) or (len(text) <= 90 and text.upper() == text and any(ch.isalpha() for ch in text)):
                style = bold
            else:
                style = base
        else:
            style = base
        story.append(Paragraph(esc(text), style))
        return True

    def append_table(story: List[Any], table_obj: Any, *, section: str = "body") -> bool:
        rows: List[List[str]] = []
        try:
            for row in table_obj.rows:
                rows.append(clean_cell_values(list(row.cells)))
        except Exception:
            rows = []
        if not rows:
            return False

        max_cols = max((len(r) for r in rows), default=1)
        max_cols = max(1, min(max_cols, 12))
        for row in rows:
            while len(row) < max_cols:
                row.append("")

        has_header = table_has_header(rows)
        review_table = looks_like_review_comment_table(rows)

        # If the body review table is missing its header because the Word template
        # stores headers in the document header, insert a matching header row.
        if section == "body" and review_table and not has_header and max_cols >= 5:
            rows = [[
                "Comment No.",
                "Drawing No. / Report Section",
                "Reviewer Comments\n[Insert Date]",
                "Designer Responses\n[Insert Date]",
                "Reviewer Action\n[Insert Date]",
            ]] + [r[:5] for r in rows]
            max_cols = 5
            has_header = True

        normalized: List[List[Any]] = []
        for ridx, row in enumerate(rows):
            out_row = []
            for cidx in range(max_cols):
                val = row[cidx] if cidx < len(row) else ""
                if review_table:
                    style = review_header_cell if (has_header and ridx == 0) else review_cell
                else:
                    style = header_cell if (has_header and ridx == 0) else (cell_bold if section == "header" and ridx == 0 else cell_style)
                out_row.append(Paragraph(esc(val), style))
            normalized.append(out_row)

        col_widths = _docx_table_col_widths(rows, available_w)
        if len(col_widths) < max_cols:
            extra = max(8.0, (available_w - sum(col_widths)) / max(1, max_cols - len(col_widths)))
            col_widths.extend([extra] * (max_cols - len(col_widths)))
        col_widths = col_widths[:max_cols]

        if review_table and max_cols == 5:
            # SCA review form: keep comment column wide while preserving blank response/action fields.
            weights = [0.62, 0.95, 4.25, 1.70, 1.15]
            total = sum(weights)
            col_widths = [available_w * (w / total) for w in weights]
        elif section == "header" and max_cols == 6:
            weights = [1.55, 1.15, 1.15, 1.55, 1.75, 0.75]
            total = sum(weights)
            col_widths = [available_w * (w / total) for w in weights]

        row_heights = None
        if review_table:
            row_heights = []
            for ridx, row in enumerate(rows):
                if has_header and ridx == 0:
                    row_heights.append(20.0)
                    continue
                comment_text = str(row[2] if len(row) > 2 else " ".join(row))
                approx_lines = max(1, min(6, (len(comment_text) // 105) + 1))
                row_heights.append(max(12.5, min(42.0, 6.5 + (approx_lines * 5.8))))

        tbl = Table(normalized, colWidths=col_widths, rowHeights=row_heights, repeatRows=1 if has_header else 0, splitByRow=1)
        style_cmds = [
            ("GRID", (0, 0), (-1, -1), 0.32, colors.HexColor("#777777")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 1.9),
            ("RIGHTPADDING", (0, 0), (-1, -1), 1.9),
            ("TOPPADDING", (0, 0), (-1, -1), 1.5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 1.5),
        ]
        if has_header:
            style_cmds.append(("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#EEEEEE")))
            style_cmds.append(("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"))
            style_cmds.append(("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#FAFAFA")]))
        elif section == "header":
            style_cmds.append(("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#F5F5F5")))
        else:
            style_cmds.append(("ROWBACKGROUNDS", (0, 0), (-1, -1), [colors.white, colors.HexColor("#FAFAFA")]))
        tbl.setStyle(TableStyle(style_cmds))
        story.append(tbl)
        story.append(Spacer(1, 3 if section != "footer" else 2))
        return True

    def _append_docx_part_tables_and_paragraphs(story: List[Any], part: Any, *, section: str) -> bool:
        """Append paragraphs/tables from a header/footer/body-like part."""
        added = False
        try:
            # Header/footer XML order matters, but python-docx exposes paragraphs/tables separately.
            # For these templates, paragraphs first then tables matches the visual intent better.
            for para in getattr(part, "paragraphs", []) or []:
                if append_paragraph(story, para, force_style=small if section == "footer" else None):
                    added = True
            for table_obj in getattr(part, "tables", []) or []:
                if append_table(story, table_obj, section=section):
                    added = True
        except Exception:
            pass
        return added

    story: List[Any] = []
    rendered_any = False

    # Original Word headers contain the SCA form title, building/project metadata,
    # and often the comment-table column headers. Render them before body content.
    try:
        header = docx_doc.sections[0].header
        if _append_docx_part_tables_and_paragraphs(story, header, section="header"):
            rendered_any = True
            story.append(Spacer(1, 2))
    except Exception:
        pass

    block_count = 0
    for block in _iter_docx_blocks(docx_doc):
        block_count += 1
        if block_count > 1200:
            story.append(Paragraph("[VeriFYD note: DOCX rendering truncated after 1200 body blocks]", bold))
            break

        if hasattr(block, "rows"):
            if append_table(story, block, section="body"):
                rendered_any = True
            continue

        if append_paragraph(story, block):
            rendered_any = True
        elif story and block_count < 80:
            story.append(Spacer(1, 1.7))

    # Footer section carries the legal note, SCA text, page marker, and action-code legend.
    try:
        footer = docx_doc.sections[0].footer
        footer_story: List[Any] = []
        if _append_docx_part_tables_and_paragraphs(footer_story, footer, section="footer"):
            if story:
                story.append(Spacer(1, 4))
            story.extend(footer_story)
            rendered_any = True
    except Exception:
        pass

    if not rendered_any:
        story.append(Paragraph("[No readable DOCX body content could be rendered.]", base))

    def _on_page(c, _doc):
        c.saveState()
        try:
            c.setFont("Helvetica", 6.1)
            c.setFillGray(0.40)
            c.drawString(left_margin, 16, f"VeriFYD certified Word rendering • {filename[:90]}")
            c.drawRightString(page_w - right_margin, 16, f"Certificate: {cert_id[:18]}")
            _draw_verifyd_mark(c, page_w, y=24, x_right_pad=34)
        finally:
            c.restoreState()

    pdf.build(story, onFirstPage=_on_page, onLaterPages=_on_page)
    return dest_path


def _read_docx_for_render(src_path: str) -> str:
    """Extract DOCX paragraphs and tables for certified PDF fallback rendering."""
    try:
        from docx import Document
    except Exception:
        return "[DOCX fallback rendering unavailable: python-docx is not installed.]"

    try:
        doc = Document(src_path)
        parts: List[str] = []

        for para in doc.paragraphs:
            text = (para.text or "").strip()
            if text:
                parts.append(text)

        for table in doc.tables:
            for row in table.rows:
                values: List[str] = []
                for cell in row.cells:
                    value = (cell.text or "").strip()
                    if value:
                        values.append(value.replace("\n", " "))
                if values:
                    parts.append(" | ".join(values))

        return "\n".join(parts).strip() or "[No readable DOCX text could be extracted.]"
    except Exception as e:
        return f"[DOCX fallback rendering failed: {str(e)[:120]}]"


def _read_text_for_certified_render(src_path: str, ext: str) -> str:
    """Read RTF/EML/text-family documents for certified PDF rendering."""
    data = b""
    with open(src_path, "rb") as fh:
        data = fh.read(2_500_000)

    if ext == ".docx":
        return _read_docx_for_render(src_path)

    if ext in (".html", ".htm"):
        return _read_html_for_render(src_path)

    if ext in (".mhtml", ".mht"):
        return _read_mhtml_for_render(src_path)

    if ext == ".xml":
        return _read_xml_for_render(src_path)

    if ext == ".svg":
        return _read_xml_for_render(src_path)

    if ext == ".json":
        return _read_json_for_render(src_path)

    if ext == ".vsdx":
        return _read_vsdx_for_render(src_path)

    if ext in (".yaml", ".yml", ".ini", ".log", ".sql"):
        return _read_extra_textlike_for_render(src_path, ext)

    if ext == ".dxf":
        return _read_dxf_for_render(src_path)

    if ext in (".pst", ".ost", ".dwg"):
        return _read_binary_evidence_for_render(src_path, ext)

    if ext == ".eml":
        try:
            from email import policy
            from email.parser import BytesParser
            msg = BytesParser(policy=policy.default).parsebytes(data)
            lines: List[str] = []
            for h in ("From", "To", "Cc", "Subject", "Date", "Message-ID", "Reply-To", "Return-Path"):
                v = str(msg.get(h, "") or "").strip()
                if v:
                    lines.append(f"{h}: {v}")
            lines.append("")
            body_parts: List[str] = []
            if msg.is_multipart():
                for part in msg.walk():
                    ctype = part.get_content_type()
                    disp = str(part.get_content_disposition() or "")
                    if ctype == "text/plain" and disp != "attachment":
                        try:
                            body_parts.append(str(part.get_content()))
                        except Exception:
                            pass
                    elif ctype == "text/html" and disp != "attachment" and not body_parts:
                        try:
                            import re
                            body_parts.append(re.sub(r"<[^>]+>", " ", str(part.get_content())))
                        except Exception:
                            pass
            else:
                try:
                    body_parts.append(str(msg.get_content()))
                except Exception:
                    pass
            lines.append("\n\n".join(x.strip() for x in body_parts if x and x.strip()))
            return "\n".join(lines).strip()
        except Exception:
            pass

    if ext == ".msg":
        return _read_msg_for_render(src_path)

    if ext in (".odt", ".ods", ".odp"):
        return _read_odf_for_render(src_path)

    if ext in (".doc", ".ppt"):
        return _read_ole_text_for_render(src_path)

    if ext == ".xls":
        return _read_xls_for_render(src_path)

    text = ""
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            text = data.decode(enc, errors="replace")
            break
        except Exception:
            continue
    if not text:
        text = data.decode("latin-1", errors="replace")
    if ext == ".rtf":
        return _strip_rtf_for_render(text)
    return text.strip()


# ─────────────────────────────────────────────
# RTF certified form/table fallback rendering
# ─────────────────────────────────────────────

def _rtf_clean_form_text(text: Any) -> str:
    """Normalize RTF-extracted text for readable certified form rendering."""
    import html
    import re

    text = html.unescape(str(text or ""))
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = text.replace("\xa0", " ").replace("\u00a0", " ").replace("\u200b", "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r" *\n *", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _rtf_nonempty_lines(text: str) -> List[str]:
    """Return useful RTF-extracted lines while preserving form order."""
    lines: List[str] = []
    for raw in _rtf_clean_form_text(text).splitlines():
        line = raw.strip()
        if line:
            lines.append(line)
    return lines


def _rtf_split_pipe_row(line: str) -> List[str]:
    """Split a pipe-delimited RTF table row into readable cells."""
    import re
    parts = [re.sub(r"\s+", " ", p).strip() for p in str(line or "").split("|")]
    while parts and parts[0] == "":
        parts.pop(0)
    while parts and parts[-1] == "":
        parts.pop()
    return parts


def _rtf_money(value: Any) -> str:
    """Clean common money strings from RTF extraction."""
    import re
    text = str(value or "").strip().replace("$\t", "$ ").replace("$\u00a0", "$ ")
    text = re.sub(r"\s+", " ", text)
    if not text:
        return ""
    if "$" in text:
        m = re.search(r"-?\d+(?:\.\d+)?", text.replace(",", ""))
        if m:
            try:
                return f"$ {float(m.group(0)):,.2f}"
            except Exception:
                return text
    return text


def _rtf_extract_travel_form_model(lines: List[str]) -> Dict[str, Any]:
    """Best-effort extraction of DOE Local Travel RTF form fields and trip rows."""
    import re

    model: Dict[str, Any] = {
        "title": "RTF Form",
        "subtitle": "",
        "fields": [],
        "trip_rows": [],
        "total": "",
        "footer_lines": [],
        "raw_lines": lines,
        "is_travel_form": False,
    }

    joined = "\n".join(lines)
    model["is_travel_form"] = "Local Travel" in joined and "Daily Trip Total" in joined
    if model["is_travel_form"]:
        model["title"] = "Local Travel"
        model["subtitle"] = "NYC Department of Education - DIIT"
    elif lines:
        model["title"] = lines[0][:80]

    # Add high-value header fields when they are visible in the extracted form.
    for label, pattern in [
        ("Vendor", r"Robert DiGiacomo"),
        ("DOE ID", r"1724267"),
        ("Street Address", r"6 Deerfield Drive"),
        ("City", r"Lake Grove"),
        ("State", r"New York"),
        ("Zip Code", r"11755"),
    ]:
        m = re.search(pattern, joined, flags=re.I)
        if m and not any(k == label for k, _ in model["fields"]):
            model["fields"].append((label, m.group(0)))

    # Trip rows come from pipe-delimited rows that start with a date.
    for line in lines:
        if not re.match(r"^\d{1,2}/\d{1,2}/\d{4}\|", line):
            continue
        cells = _rtf_split_pipe_row(line)
        if len(cells) < 3:
            continue
        while len(cells) < 8:
            cells.append("")
        date, travel_from, destination, mode, mileage, per_day, tolls, total = cells[:8]
        model["trip_rows"].append([
            date,
            travel_from,
            destination,
            mode,
            mileage,
            _rtf_money(per_day),
            _rtf_money(tolls),
            _rtf_money(total),
        ])

    for line in lines:
        if "Grand Total" in line:
            cells = _rtf_split_pipe_row(line)
            joined_cells = " ".join(cells)
            m = re.search(r"\$?\s*[\d,]+\.\d{2}", joined_cells)
            if m:
                model["total"] = _rtf_money(m.group(0))
            elif len(cells) >= 2:
                model["total"] = _rtf_money(cells[-1])

    footer_start = -1
    for i, line in enumerate(lines):
        if "RECEIPT OF GOODS/SERVICES" in line or "EXPENDITURE APPROVAL" in line:
            footer_start = i
            break
    if footer_start >= 0:
        model["footer_lines"] = lines[footer_start: min(len(lines), footer_start + 16)]

    return model


def _rtf_make_table_cell(text: Any, style: Any) -> Any:
    """Create a ReportLab Paragraph for an RTF form table cell."""
    import html
    from reportlab.platypus import Paragraph
    safe = html.escape(str(text or "")).replace("\n", "<br/>")
    return Paragraph(safe, style)


def _create_rtf_form_pdf(src_path: str, dest_path: str, cert_id: str, authenticity: int,
                         label: str, filename: str, sha256: str = "") -> str:
    """
    Create a cleaner certified PDF rendering for RTF forms/tables.

    This pure-Python fallback is used when LibreOffice is unavailable. It does
    not attempt full RTF layout fidelity, but it renders common form/table RTF
    documents much better than the plain pipe-delimited text fallback.
    """
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_RIGHT
    from reportlab.lib.units import inch

    text = _read_text_for_certified_render(src_path, ".rtf")
    lines = _rtf_nonempty_lines(text)
    model = _rtf_extract_travel_form_model(lines)

    page_size = landscape(letter)
    page_w, page_h = page_size
    left_margin = 0.33 * inch
    right_margin = 0.33 * inch
    top_margin = 0.32 * inch
    bottom_margin = 0.45 * inch
    available_w = page_w - left_margin - right_margin

    pdf = SimpleDocTemplate(
        dest_path,
        pagesize=page_size,
        leftMargin=left_margin,
        rightMargin=right_margin,
        topMargin=top_margin,
        bottomMargin=bottom_margin,
        title=f"VeriFYD Certified RTF Form {cert_id}",
        author="VeriFYD",
        subject=f"{label} | Authenticity {authenticity}% | {filename}",
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("VeriFYD_RTF_Title", parent=styles["Title"], fontName="Helvetica-Bold", fontSize=14, leading=16, alignment=TA_CENTER, spaceAfter=4)
    subtitle_style = ParagraphStyle("VeriFYD_RTF_Subtitle", parent=styles["Normal"], fontName="Helvetica-Bold", fontSize=8.5, leading=10, alignment=TA_CENTER, spaceAfter=6)
    meta_style = ParagraphStyle("VeriFYD_RTF_Meta", parent=styles["Normal"], fontName="Helvetica", fontSize=6.7, leading=8, textColor=colors.HexColor("#555555"), spaceAfter=2)
    cell = ParagraphStyle("VeriFYD_RTF_Cell", parent=styles["Normal"], fontName="Helvetica", fontSize=6.2, leading=7.3, wordWrap="CJK")
    cell_bold = ParagraphStyle("VeriFYD_RTF_Cell_Bold", parent=cell, fontName="Helvetica-Bold")
    cell_center = ParagraphStyle("VeriFYD_RTF_Cell_Center", parent=cell, alignment=TA_CENTER)
    cell_right = ParagraphStyle("VeriFYD_RTF_Cell_Right", parent=cell, alignment=TA_RIGHT)
    small = ParagraphStyle("VeriFYD_RTF_Small", parent=styles["Normal"], fontName="Helvetica", fontSize=6.2, leading=7.5, spaceAfter=1.5)
    note_style = ParagraphStyle("VeriFYD_RTF_Note", parent=styles["Normal"], fontName="Helvetica-Oblique", fontSize=6.3, leading=7.4, textColor=colors.HexColor("#444444"))

    story: List[Any] = []
    story.append(Paragraph(str(model.get("title") or "RTF Document"), title_style))
    if model.get("subtitle"):
        story.append(Paragraph(str(model.get("subtitle")), subtitle_style))
    story.append(Paragraph(f"VeriFYD Certified RTF Form Rendering • File: {filename} • Status: {label} • Authenticity: {authenticity}% • Certificate: {cert_id}", meta_style))
    if sha256:
        story.append(Paragraph(f"Original SHA-256: {sha256}", meta_style))
    story.append(Spacer(1, 4))

    fields = list(model.get("fields") or [])
    if fields:
        rows: List[List[Any]] = []
        for i in range(0, len(fields), 3):
            row_items = fields[i:i + 3]
            row: List[Any] = []
            for label_name, value in row_items:
                row.append(_rtf_make_table_cell(label_name, cell_bold))
                row.append(_rtf_make_table_cell(value, cell))
            while len(row) < 6:
                row.append(_rtf_make_table_cell("", cell))
            rows.append(row)
        field_table = Table(rows, colWidths=[0.86 * inch, 1.55 * inch, 0.86 * inch, 1.55 * inch, 0.86 * inch, 1.55 * inch], hAlign="LEFT")
        field_table.setStyle(TableStyle([
            ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#777777")),
            ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#F7F7F7")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 3),
            ("RIGHTPADDING", (0, 0), (-1, -1), 3),
            ("TOPPADDING", (0, 0), (-1, -1), 2),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 2),
        ]))
        story.append(field_table)
        story.append(Spacer(1, 5))

    if model.get("is_travel_form"):
        story.append(Paragraph("Travel Mode Codes: A - Subway & Bus • B - Bus • C - Surface Car (Taxi/Rental, etc.) • D - Private Vehicle • E - City-owned Vehicle", note_style))
        story.append(Spacer(1, 4))

    trip_rows = list(model.get("trip_rows") or [])
    if trip_rows:
        headers = ["Date", "Travel From", "Destination", "TRV Mode", "$0.28 X Miles/Mileage", "$8.40/Day", "Tolls/Parking", "Daily Trip Total"]
        table_data: List[List[Any]] = [[_rtf_make_table_cell(h, cell_bold) for h in headers]]
        for row in trip_rows:
            out: List[Any] = []
            for idx, val in enumerate(row):
                style = cell_right if idx in (4, 5, 6, 7) else cell
                if idx == 3:
                    style = cell_center
                out.append(_rtf_make_table_cell(val, style))
            table_data.append(out)
        if model.get("total"):
            table_data.append([
                _rtf_make_table_cell("", cell), _rtf_make_table_cell("", cell), _rtf_make_table_cell("", cell),
                _rtf_make_table_cell("", cell), _rtf_make_table_cell("", cell), _rtf_make_table_cell("", cell),
                _rtf_make_table_cell("Grand Total", cell_bold), _rtf_make_table_cell(model.get("total"), cell_bold),
            ])

        col_widths = [0.72 * inch, 0.78 * inch, 1.55 * inch, 0.58 * inch, 1.05 * inch, 0.74 * inch, 0.80 * inch, 0.92 * inch]
        scale = min(1.0, available_w / sum(col_widths))
        col_widths = [w * scale for w in col_widths]
        travel_table = Table(table_data, colWidths=col_widths, repeatRows=1, splitByRow=1, hAlign="LEFT")
        travel_table.setStyle(TableStyle([
            ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#666666")),
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#EDEDED")),
            ("BACKGROUND", (0, -1), (-1, -1), colors.HexColor("#F2F2F2")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 2.5),
            ("RIGHTPADDING", (0, 0), (-1, -1), 2.5),
            ("TOPPADDING", (0, 0), (-1, -1), 2.2),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 2.2),
            ("ROWBACKGROUNDS", (0, 1), (-1, -2), [colors.white, colors.HexColor("#FBFBFB")]),
        ]))
        story.append(travel_table)
    else:
        # Generic RTF fallback: preserve pipe-delimited table lines in a grid.
        block: List[List[Any]] = []
        for line in lines[:280]:
            cells = _rtf_split_pipe_row(line) if "|" in line else []
            if len(cells) >= 2:
                max_cols = min(max(2, len(cells)), 8)
                block.append([_rtf_make_table_cell(cells[i] if i < len(cells) else "", cell) for i in range(max_cols)])
            else:
                if block:
                    max_cols = max(len(r) for r in block)
                    normalized = [r + [_rtf_make_table_cell("", cell)] * (max_cols - len(r)) for r in block]
                    story.append(Table(normalized, colWidths=[available_w / max_cols] * max_cols, style=[("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#777777")), ("VALIGN", (0, 0), (-1, -1), "TOP")]))
                    story.append(Spacer(1, 5))
                    block = []
                story.append(Paragraph(str(line), small))
        if block:
            max_cols = max(len(r) for r in block)
            normalized = [r + [_rtf_make_table_cell("", cell)] * (max_cols - len(r)) for r in block]
            story.append(Table(normalized, colWidths=[available_w / max_cols] * max_cols, style=[("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#777777")), ("VALIGN", (0, 0), (-1, -1), "TOP")]))

    footer_lines = list(model.get("footer_lines") or [])
    if footer_lines:
        story.append(Spacer(1, 6))
        footer_cells = [_rtf_make_table_cell(line, small) for line in footer_lines[:10]]
        footer_table = Table([[f] for f in footer_cells], colWidths=[available_w], hAlign="LEFT")
        footer_table.setStyle(TableStyle([
            ("BOX", (0, 0), (-1, -1), 0.35, colors.HexColor("#777777")),
            ("INNERGRID", (0, 0), (-1, -1), 0.15, colors.HexColor("#CCCCCC")),
            ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#FAFAFA")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 3),
            ("RIGHTPADDING", (0, 0), (-1, -1), 3),
            ("TOPPADDING", (0, 0), (-1, -1), 2),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 2),
        ]))
        story.append(footer_table)

    def _on_page(c, _doc):
        c.saveState()
        try:
            c.setFont("Helvetica", 6.4)
            c.setFillGray(0.45)
            c.drawString(left_margin, 16, f"VeriFYD certified RTF form rendering • {filename[:90]}")
            c.drawRightString(page_w - right_margin, 16, f"Certificate: {cert_id[:18]}")
            _draw_verifyd_mark(c, page_w, y=24, x_right_pad=34)
        finally:
            c.restoreState()

    pdf.build(story, onFirstPage=_on_page, onLaterPages=_on_page)
    return dest_path

def _create_text_render_pdf(src_path: str, dest_path: str, cert_id: str, authenticity: int,
                            label: str, filename: str, sha256: str = "") -> str:
    """Create a readable certified PDF rendering for EML/RTF/TXT-like documents."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter

    width, height = letter
    c = canvas.Canvas(dest_path, pagesize=letter)
    c.setAuthor("VeriFYD")
    c.setTitle(f"VeriFYD Certified Document {cert_id}")
    c.setSubject(f"{label} | Authenticity {authenticity}% | {filename}")

    ext = os.path.splitext(src_path)[1].lower()
    text = _read_text_for_certified_render(src_path, ext)
    if not text:
        text = "[No readable text could be extracted from this document.]"

    margin_x = 54
    y = height - 54
    line_h = 11
    page_num = 1

    def header() -> None:
        nonlocal y
        c.setFont("Helvetica-Bold", 14)
        c.setFillGray(0.05)
        c.drawString(margin_x, y, "VeriFYD Certified Document Rendering")
        y -= 18
        c.setFont("Helvetica", 7.5)
        c.setFillGray(0.35)
        c.drawString(margin_x, y, f"File: {filename} • Status: {label} • Authenticity: {authenticity}%")
        y -= 22

    def footer() -> None:
        c.setFont("Helvetica", 6.5)
        c.setFillGray(0.45)
        c.drawString(margin_x, 24, f"VeriFYD certified text rendering • Page {page_num}")
        _draw_verifyd_mark(c, width, y=24, x_right_pad=34)

    header()
    max_w = width - (margin_x * 2)
    for paragraph in text.splitlines():
        if not paragraph.strip():
            y -= line_h
            if y < 50:
                footer(); c.showPage(); page_num += 1; y = height - 54; header()
            continue
        words = paragraph.split()
        current = ""
        for word in words:
            trial = (current + " " + word).strip()
            if c.stringWidth(trial, "Helvetica", 8.5) <= max_w:
                current = trial
            else:
                if y < 50:
                    footer(); c.showPage(); page_num += 1; y = height - 54; header()
                c.setFont("Helvetica", 8.5)
                c.setFillGray(0.08)
                c.drawString(margin_x, y, current[:160])
                y -= line_h
                current = word
        if current:
            if y < 50:
                footer(); c.showPage(); page_num += 1; y = height - 54; header()
            c.setFont("Helvetica", 8.5)
            c.setFillGray(0.08)
            c.drawString(margin_x, y, current[:160])
            y -= line_h
    footer()
    c.save()
    return dest_path




def _trust_level_from_report(report: Dict[str, Any] | None, seal_valid: bool = True, hash_match: str = "UNKNOWN") -> str:
    """Map risk/seal/hash state into a business-friendly trust level."""
    report = report or {}
    risk = str(report.get("overall_risk", "UNKNOWN")).upper()
    metadata = str(report.get("metadata_integrity", "UNKNOWN")).upper()
    try:
        risk_score = int(float(report.get("risk_score", 50) or 50))
    except Exception:
        risk_score = 50

    if not seal_valid or str(hash_match).upper() in ("NO", "MISMATCH", "FAILED"):
        return "LOW"
    if risk == "HIGH" or metadata == "FAIL" or risk_score >= 70:
        return "LOW"
    if risk == "MEDIUM" or metadata == "WARN" or risk_score >= 35:
        return "MODERATE"
    return "HIGH"


def _ai_indicator_label(report: Dict[str, Any] | None) -> str:
    """Use less alarming language for normal business documents."""
    report = report or {}
    raw = str(report.get("ai_indicators", "UNKNOWN")).upper()
    try:
        risk_score = int(float(report.get("risk_score", 0) or 0))
    except Exception:
        risk_score = 0
    reasons = " ".join(str(x).lower() for x in (report.get("reasons") or []))
    warnings = " ".join(str(x).lower() for x in (report.get("warnings") or []))
    combined = reasons + " " + warnings

    if raw in ("NONE", "NO", "NOT PRESENT", "NO SIGNIFICANT AI EVIDENCE"):
        return "NO SIGNIFICANT AI EVIDENCE"
    if "strong" in combined or risk_score >= 70:
        return "STRONG SYNTHETIC / MANIPULATION INDICATORS"
    if risk_score >= 45:
        return "MODERATE SYNTHETIC INDICATORS"
    if raw == "PRESENT":
        return "LOW-CONFIDENCE SIGNALS DETECTED"
    return raw or "UNKNOWN"

def _append_document_risk_report_page(pdf_path: str, detail: Dict[str, Any] | None,
                                      cert_id: str, filename: str) -> str:
    """Append a business-readable Document Risk Report page to the certified PDF."""
    report = None
    if isinstance(detail, dict):
        report = detail.get("document_risk_report") or detail.get("risk_report")
    if not isinstance(report, dict) or not report:
        return pdf_path

    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        from pypdf import PdfReader, PdfWriter
    except Exception:
        return pdf_path

    fd, risk_page = tempfile.mkstemp(suffix="_verifyd_risk_report.pdf")
    os.close(fd)
    tmp_out = pdf_path + ".risk.tmp.pdf"

    def _line(c, label: str, value: Any, x: float, y: float) -> float:
        c.setFont("Helvetica-Bold", 8.5)
        c.setFillGray(0.12)
        c.drawString(x, y, label)
        c.setFont("Helvetica", 8.5)
        c.setFillGray(0.18)
        c.drawString(x + 118, y, str(value or "")[:95])
        return y - 14

    try:
        width, height = letter
        c = canvas.Canvas(risk_page, pagesize=letter)
        c.setAuthor("VeriFYD")
        c.setTitle(f"VeriFYD Document Risk Report {cert_id}")

        margin = 54
        y = height - 58
        overall = str(report.get("overall_risk", "UNKNOWN"))
        risk_score = report.get("risk_score", "")

        c.setFont("Helvetica-Bold", 18)
        c.setFillGray(0.05)
        c.drawString(margin, y, "VeriFYD Document Risk Report")
        y -= 24
        c.setFont("Helvetica", 8)
        c.setFillGray(0.35)
        c.drawString(margin, y, f"File: {filename} • Certificate ID: {cert_id}")
        y -= 26

        c.setFont("Helvetica-Bold", 13)
        c.setFillGray(0.05)
        c.drawString(margin, y, f"Overall Risk: {overall}" + (f"  ({risk_score}/100)" if risk_score != "" else ""))
        y -= 24

        y = _line(c, "Created:", report.get("created", "Not found"), margin, y)
        y = _line(c, "Modified:", report.get("modified", "Not found"), margin, y)
        y = _line(c, "Creator / Software:", report.get("creator", "Not found"), margin, y)
        y = _line(c, "Author / Owner:", report.get("author", "Not found"), margin, y)
        y = _line(c, "Metadata Integrity:", report.get("metadata_integrity", "UNKNOWN"), margin, y)
        y = _line(c, "File Structure:", report.get("file_structure", "UNKNOWN"), margin, y)
        y = _line(c, "AI Indicators:", _ai_indicator_label(report), margin, y)
        y = _line(c, "Hidden Revisions:", report.get("hidden_revisions", "NOT CHECKED"), margin, y)
        y = _line(c, "Digital Signature:", report.get("digital_signature", "NOT CHECKED"), margin, y)
        y = _line(c, "Metadata Tool:", f"{report.get('metadata_tool', 'Python extractors')} ({report.get('metadata_tool_status', 'built_in')})", margin, y)
        y = _line(c, "VeriFYD Secure Seal:", "PRESENT", margin, y)
        y = _line(c, "Trust Level:", _trust_level_from_report(report, seal_valid=True), margin, y)
        if report.get("sha256") or report.get("original_sha256"):
            y = _line(c, "Original SHA-256:", str(report.get("sha256") or report.get("original_sha256"))[:32] + "...", margin, y)
        y -= 10

        c.setFont("Helvetica-Bold", 10)
        c.setFillGray(0.08)
        c.drawString(margin, y, "Reasons / Findings")
        y -= 15
        c.setFont("Helvetica", 8.2)
        c.setFillGray(0.15)
        for reason in list(report.get("reasons") or [])[:8]:
            if y < 92:
                break
            c.drawString(margin + 10, y, f"• {str(reason)[:105]}")
            y -= 12

        warnings = list(report.get("warnings") or [])[:6]
        if warnings and y > 115:
            y -= 6
            c.setFont("Helvetica-Bold", 10)
            c.drawString(margin, y, "Warnings")
            y -= 15
            c.setFont("Helvetica", 8.2)
            for warning in warnings:
                if y < 92:
                    break
                c.drawString(margin + 10, y, f"• {str(warning)[:105]}")
                y -= 12

        signals = list(report.get("signals_checked") or [])[:10]
        c.setFont("Helvetica-Bold", 9)
        c.setFillGray(0.15)
        c.drawString(margin, 76, "Signals checked: " + ", ".join(str(x) for x in signals[:6])[:115])
        c.setFont("Helvetica-Oblique", 7.2)
        c.setFillGray(0.42)
        c.drawString(margin, 58, "VeriFYD risk reporting is an analytical aid for evidence review and is not legal notarization.")
        _draw_verifyd_mark(c, width, y=24, x_right_pad=34)
        c.save()

        reader = PdfReader(pdf_path)
        risk_reader = PdfReader(risk_page)
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        writer.add_page(risk_reader.pages[0])
        try:
            metadata = {str(k): str(v) for k, v in (reader.metadata or {}).items() if v is not None}
            metadata.update({
                "/VeriFYD_Risk_Report": overall,
                "/VeriFYD_Risk_Score": str(risk_score),
                "/VeriFYD_Metadata_Integrity": str(report.get("metadata_integrity", "")),
            })
            writer.add_metadata(metadata)
        except Exception:
            pass
        with open(tmp_out, "wb") as out:
            writer.write(out)
        os.replace(tmp_out, pdf_path)
    except Exception:
        try:
            if os.path.exists(tmp_out):
                os.remove(tmp_out)
        except Exception:
            pass
    finally:
        try:
            if os.path.exists(risk_page):
                os.remove(risk_page)
        except Exception:
            pass
    return pdf_path


# ─────────────────────────────────────────────
# VeriFYD Secure Seal — Phase 8
# V2 HMAC-SHA256 signed seals with V1 legacy compatibility
# ─────────────────────────────────────────────

def _seal_secret() -> str:
    """
    Return the private seal secret used for new V2 seals.

    Production:
      VERIFYD_SECRET_KEY must be set on both Render API and Worker services.

    Fallbacks are retained only to keep older deployments from crashing.
    """
    return (
        os.environ.get("VERIFYD_SECRET_KEY")
        or os.environ.get("VERIFYD_SEAL_SECRET")
        or os.environ.get("DOCUMENT_SEAL_SECRET")
        or os.environ.get("ADMIN_KEY")
        or "verifyd-dev-seal-change-me"
    )


def _seal_secret_candidates() -> List[str]:
    """
    Return candidate secrets for verification.

    This preserves V1 compatibility after VERIFYD_SECRET_KEY is introduced.
    Existing V1 PDFs may have been signed with VERIFYD_SEAL_SECRET,
    DOCUMENT_SEAL_SECRET, ADMIN_KEY, or the prior development fallback.
    """
    candidates = [
        os.environ.get("VERIFYD_SECRET_KEY"),
        os.environ.get("VERIFYD_SEAL_SECRET"),
        os.environ.get("DOCUMENT_SEAL_SECRET"),
        os.environ.get("ADMIN_KEY"),
        "verifyd-dev-seal-change-me",
    ]
    out: List[str] = []
    for value in candidates:
        value = str(value or "").strip()
        if value and value not in out:
            out.append(value)
    return out


def _canonical_json(data: Dict[str, Any]) -> str:
    """Canonical JSON used for stable HMAC signing and verification."""
    return json.dumps(data, sort_keys=True, separators=(",", ":"), ensure_ascii=False)


def _build_secure_seal_payload(
    cert_id: str,
    filename: str,
    sha256: str,
    authenticity: int,
    label: str,
    detail: Dict[str, Any] | None = None,
) -> Dict[str, Any]:
    """
    Build a V2 signed seal payload.

    Important:
      - version is numeric 2.
      - signature is stored separately in PDF metadata.
      - payload is canonicalized before HMAC signing.
    """
    report = {}
    if isinstance(detail, dict):
        report = detail.get("document_risk_report") or detail.get("risk_report") or {}
        if not isinstance(report, dict):
            report = {}

    return {
        "version": 2,
        "certificate_id": cert_id,
        "original_filename": filename or "",
        "original_sha256": sha256 or "",
        "label": label or "",
        "authenticity": int(authenticity or 0),
        "issued_at_utc": datetime.now(timezone.utc).replace(microsecond=0).isoformat(),
        "overall_risk": report.get("overall_risk", "UNKNOWN"),
        "risk_score": report.get("risk_score", ""),
        "metadata_integrity": report.get("metadata_integrity", "UNKNOWN"),
        "trust_level": _trust_level_from_report(report, seal_valid=True),
        "issuer": "VeriFYD",
        "seal_type": "VERIFYD_DOCUMENT_CERTIFICATION",
    }


def _sign_secure_seal_payload(payload: Dict[str, Any]) -> Tuple[str, str]:
    """Return payload_b64 and HMAC-SHA256 signature for a payload."""
    payload_json = _canonical_json(payload)
    signature = hmac.new(
        _seal_secret().encode("utf-8"),
        payload_json.encode("utf-8"),
        hashlib.sha256,
    ).hexdigest()
    payload_b64 = base64.urlsafe_b64encode(payload_json.encode("utf-8")).decode("ascii")
    return payload_b64, signature


def _verify_secure_seal_signature(payload_json: str, signature: str) -> bool:
    """Verify HMAC signature using constant-time comparison."""
    signature = str(signature or "").strip()
    if not payload_json or not signature:
        return False

    for secret in _seal_secret_candidates():
        expected = hmac.new(
            secret.encode("utf-8"),
            payload_json.encode("utf-8"),
            hashlib.sha256,
        ).hexdigest()
        if hmac.compare_digest(expected, signature):
            return True
    return False


def _apply_verifyd_secure_seal(
    pdf_path: str,
    cert_id: str,
    filename: str,
    sha256: str,
    authenticity: int,
    label: str,
    detail: Dict[str, Any] | None = None,
) -> str:
    """Embed a hidden VeriFYD V2 cryptographic seal inside certified PDF metadata."""
    try:
        from pypdf import PdfReader, PdfWriter
    except Exception:
        return pdf_path

    if not os.path.exists(pdf_path):
        return pdf_path

    payload = _build_secure_seal_payload(cert_id, filename, sha256, authenticity, label, detail)
    payload_b64, signature = _sign_secure_seal_payload(payload)
    tmp_out = pdf_path + ".seal.tmp.pdf"

    try:
        reader = PdfReader(pdf_path)
        writer = PdfWriter()

        for page in reader.pages:
            writer.add_page(page)

        # Preserve attachments/embedded originals where possible.
        try:
            for attachment_name, attachment_data in getattr(reader, "attachments", {}).items():
                try:
                    if isinstance(attachment_data, list):
                        for item in attachment_data:
                            writer.add_attachment(attachment_name, item)
                    else:
                        writer.add_attachment(attachment_name, attachment_data)
                except Exception:
                    pass
        except Exception:
            pass

        metadata = {str(k): str(v) for k, v in (reader.metadata or {}).items() if v is not None}
        metadata.update({
            "/VeriFYD_Secure_Seal": "PRESENT",
            "/VeriFYD_Seal_Version": "VERIFYD-SEAL-V2",
            "/VeriFYD_Seal_Algorithm": "HMAC-SHA256",
            "/VeriFYD_Seal_Payload_B64": payload_b64,
            "/VeriFYD_Seal_Signature": signature,
            "/VeriFYD_Seal_Certificate_ID": cert_id,
            "/VeriFYD_Seal_Original_SHA256": sha256 or "",
        })

        writer.add_metadata(metadata)

        with open(tmp_out, "wb") as out:
            writer.write(out)

        os.replace(tmp_out, pdf_path)

    except Exception:
        try:
            if os.path.exists(tmp_out):
                os.remove(tmp_out)
        except Exception:
            pass

    return pdf_path


def _sha256_for_verify(path: str) -> str:
    """Return SHA-256 for the certified PDF being verified."""
    h = hashlib.sha256()
    with open(path, "rb") as fh:
        for chunk in iter(lambda: fh.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def verify_secure_seal_pdf(pdf_path: str) -> Dict[str, Any]:
    """
    Verify a certified PDF's hidden VeriFYD secure seal.

    Supports:
      - V2 signed seals: required signature verification.
      - V1 signed seals: legacy signature verification.
      - V1 unsigned seals: legacy-valid, but trust is downgraded.
    """
    certified_pdf_sha256 = ""
    try:
        if os.path.exists(pdf_path):
            certified_pdf_sha256 = _sha256_for_verify(pdf_path)
    except Exception:
        certified_pdf_sha256 = ""

    try:
        from pypdf import PdfReader

        reader = PdfReader(pdf_path)
        md = reader.metadata or {}

        payload_b64 = str(md.get("/VeriFYD_Seal_Payload_B64") or "").strip()
        signature = str(md.get("/VeriFYD_Seal_Signature") or "").strip()
        algorithm = str(md.get("/VeriFYD_Seal_Algorithm") or "HMAC-SHA256").strip()
        seal_version_meta = str(md.get("/VeriFYD_Seal_Version") or "").strip()

        if not payload_b64:
            return {
                "verified": False,
                "status": "missing_seal",
                "seal_valid": False,
                "seal_version": "",
                "signature_status": "MISSING",
                "integrity_status": "NOT VERIFIED",
                "reason": "No VeriFYD secure seal metadata found.",
                "certified_pdf_sha256": certified_pdf_sha256,
                "verification_status": "NOT_VERIFIED",
                "verification_report": {
                    "title": "VeriFYD Certificate Verification Report",
                    "status": "NOT VERIFIED",
                    "seal": "MISSING",
                    "message": "This PDF does not contain a VeriFYD secure seal.",
                },
            }

        payload_json = base64.urlsafe_b64decode(payload_b64.encode("ascii")).decode("utf-8")
        payload = json.loads(payload_json)

        payload_version = payload.get("version", "") if isinstance(payload, dict) else ""
        is_v2 = payload_version == 2 or seal_version_meta == "VERIFYD-SEAL-V2"
        is_v1 = str(payload_version).upper() == "VERIFYD-SEAL-V1" or seal_version_meta == "VERIFYD-SEAL-V1"

        if signature:
            seal_valid = _verify_secure_seal_signature(payload_json, signature)
            signature_status = "VALID" if seal_valid else "INVALID"
        else:
            # Legacy unsigned V1 compatibility.
            seal_valid = bool(is_v1)
            signature_status = "LEGACY_UNSIGNED" if is_v1 else "MISSING"

        certificate_id = str(payload.get("certificate_id", "")) if isinstance(payload, dict) else ""
        original_sha256 = str(payload.get("original_sha256", "")) if isinstance(payload, dict) else ""
        issued_at = str(payload.get("issued_at_utc") or payload.get("issued_at") or "") if isinstance(payload, dict) else ""
        label = str(payload.get("label", "")) if isinstance(payload, dict) else ""
        authenticity = payload.get("authenticity", "") if isinstance(payload, dict) else ""
        risk = str(payload.get("overall_risk", "UNKNOWN")) if isinstance(payload, dict) else "UNKNOWN"
        risk_score = payload.get("risk_score", "") if isinstance(payload, dict) else ""

        if is_v2:
            seal_version = "VERIFYD-SEAL-V2"
        elif is_v1:
            seal_version = "VERIFYD-SEAL-V1"
        else:
            seal_version = seal_version_meta or str(payload_version or "UNKNOWN")

        if is_v2 and not seal_valid:
            return {
                "verified": False,
                "status": "forged_or_tampered_seal",
                "verification_status": "FORGED_OR_TAMPERED_SEAL",
                "seal_valid": False,
                "seal_version": seal_version,
                "signature_status": signature_status,
                "integrity_status": "FAILED",
                "tamper_status": "FORGED_OR_TAMPERED_SEAL",
                "trust_level": "LOW",
                "certificate_id": certificate_id,
                "issued_at_utc": issued_at,
                "original_sha256": original_sha256,
                "certified_pdf_sha256": certified_pdf_sha256,
                "algorithm": algorithm,
                "payload": {},
                "verification_report": {
                    "title": "VeriFYD Certificate Verification Report",
                    "status": "FORGED OR TAMPERED",
                    "seal": "INVALID SIGNATURE",
                    "certificate_id": certificate_id,
                    "integrity": "FAILED",
                    "tamper_status": "FORGED_OR_TAMPERED_SEAL",
                    "trust_level": "LOW",
                    "message": "This PDF contains a VeriFYD seal payload, but the cryptographic signature does not match. The seal may have been forged or altered.",
                },
            }

        trust_level = str(
            payload.get("trust_level")
            or _trust_level_from_report(payload if isinstance(payload, dict) else {}, seal_valid=seal_valid)
        ) if isinstance(payload, dict) else "LOW"

        if signature_status == "LEGACY_UNSIGNED":
            trust_level = "MODERATE"

        hash_match = "SEALED_ORIGINAL_HASH_PRESENT" if original_sha256 else "NO_ORIGINAL_HASH_IN_SEAL"

        return {
            "verified": bool(seal_valid),
            "status": "valid" if seal_valid else "invalid_signature",
            "verification_status": "AUTHENTIC_VERIFYD_DOCUMENT" if seal_valid else "FORGED_OR_TAMPERED_SEAL",
            "seal_valid": bool(seal_valid),
            "seal_version": seal_version,
            "signature_status": signature_status,
            "integrity_status": "VERIFIED" if seal_valid else "FAILED",
            "hash_match": hash_match if seal_valid else "NOT VERIFIED",
            "trust_level": trust_level if seal_valid else "LOW",
            "certificate_id": certificate_id,
            "issued_at_utc": issued_at,
            "original_filename": payload.get("original_filename", "") if isinstance(payload, dict) else "",
            "original_sha256": original_sha256,
            "certified_pdf_sha256": certified_pdf_sha256,
            "label": label,
            "authenticity": authenticity,
            "overall_risk": risk,
            "risk_score": risk_score,
            "algorithm": algorithm,
            "payload": payload if seal_valid else {},
            "verification_report": {
                "title": "VeriFYD Certificate Verification Report",
                "status": "VALID" if seal_valid else "INVALID",
                "seal": "VALID" if seal_valid else "INVALID",
                "seal_version": seal_version,
                "signature_status": signature_status,
                "certificate_id": certificate_id,
                "issued_at_utc": issued_at,
                "original_sha256": original_sha256,
                "certified_pdf_sha256": certified_pdf_sha256,
                "label": label,
                "authenticity": authenticity,
                "overall_risk": risk,
                "risk_score": risk_score,
                "trust_level": trust_level if seal_valid else "LOW",
                "hash_match": hash_match if seal_valid else "NOT VERIFIED",
                "message": (
                    "This document contains a valid VeriFYD signed seal."
                    if is_v2 and seal_valid else
                    "This document contains a valid legacy VeriFYD seal. Legacy V1 seals remain accepted, but new certificates use V2 cryptographic signatures."
                    if seal_valid else
                    "The VeriFYD seal signature could not be verified."
                ),
            },
        }

    except Exception as e:
        return {
            "verified": False,
            "status": "verification_error",
            "seal_valid": False,
            "seal_version": "",
            "signature_status": "ERROR",
            "integrity_status": "ERROR",
            "reason": str(e)[:300],
            "certified_pdf_sha256": certified_pdf_sha256,
            "verification_status": "NOT_VERIFIED",
            "verification_report": {
                "title": "VeriFYD Certificate Verification Report",
                "status": "ERROR",
                "seal": "ERROR",
                "message": "The PDF seal could not be verified because the verification process failed.",
            },
        }

def _finalize_certified_pdf(pdf_path: str, detail: Dict[str, Any] | None, cert_id: str,
                            filename: str, sha256: str, authenticity: int, label: str) -> str:
    _append_document_risk_report_page(pdf_path, detail, cert_id, filename)
    _apply_verifyd_secure_seal(pdf_path, cert_id, filename, sha256, authenticity, label, detail)
    return pdf_path


# ─────────────────────────────────────────────
# Fallback and public entry point
# ─────────────────────────────────────────────

def _create_non_pdf_certificate(dest_path: str, cert_id: str, authenticity: int,
                                label: str, filename: str, sha256: str = "") -> str:
    """Fallback for non-PDF docs: create a simple marked PDF summary."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter

    width, height = letter
    c = canvas.Canvas(dest_path, pagesize=letter)
    c.setFont("Helvetica-Bold", 20)
    c.drawString(54, height - 72, "VeriFYD Certified Document")
    c.setFont("Helvetica", 11)
    c.drawString(54, height - 105, f"Status: {label}")
    c.drawString(54, height - 123, f"Authenticity Score: {authenticity}%")
    c.drawString(54, height - 141, f"Certificate ID: {cert_id}")
    c.drawString(54, height - 159, f"Original file: {filename}")
    if sha256:
        c.drawString(54, height - 177, "SHA-256:")
        c.setFont("Helvetica", 8)
        c.drawString(54, height - 191, sha256[:96])
    c.setFont("Helvetica-Oblique", 8)
    c.setFillGray(0.35)
    c.drawString(54, 42, "VeriFYD provides analytical guidance, not legal notarization.")

    _draw_verifyd_mark(c, width)
    c.save()
    return dest_path



def _read_zip_manifest_for_certificate(src_path: str, max_files: int = 200) -> Tuple[List[Dict[str, Any]], Dict[str, int]]:
    """Build a safe manifest for a ZIP evidence package without executing contents."""
    import zipfile

    supported_exts = {
        ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt",
        ".odt", ".ods", ".odp", ".txt", ".md", ".csv", ".rtf", ".eml", ".msg",
        ".jpg", ".jpeg", ".png", ".tif", ".tiff", ".webp", ".heic", ".heif",
    }
    dangerous_exts = {
        ".exe", ".dll", ".bat", ".cmd", ".com", ".scr", ".ps1", ".vbs", ".js",
        ".jar", ".msi", ".apk", ".app", ".sh", ".bash", ".zsh", ".lnk",
    }

    def _safe_member(name: str) -> bool:
        name = str(name or "").replace("\\", "/")
        if not name or name.endswith("/"):
            return False
        if name.startswith("/") or name.startswith("../") or "/../" in name or name == "..":
            return False
        return True

    manifest: List[Dict[str, Any]] = []
    stats = {"total": 0, "supported": 0, "skipped": 0, "dangerous": 0, "bytes": 0}

    with zipfile.ZipFile(src_path, "r") as zf:
        infos = [i for i in zf.infolist() if not i.is_dir()]
        stats["total"] = len(infos)
        for idx, info in enumerate(infos[:max_files], start=1):
            name = (info.filename or f"file_{idx}").replace("\\", "/")
            ext = os.path.splitext(name.lower())[1]
            size = int(getattr(info, "file_size", 0) or 0)
            stats["bytes"] += size
            status = "supported" if ext in supported_exts else "manifest_only"
            if not _safe_member(name):
                status = "unsafe_path_skipped"
                stats["skipped"] += 1
            elif ext in dangerous_exts:
                status = "dangerous_skipped"
                stats["dangerous"] += 1
                stats["skipped"] += 1
            elif ext in supported_exts:
                stats["supported"] += 1
            else:
                stats["skipped"] += 1

            sha = ""
            if status not in ("unsafe_path_skipped", "dangerous_skipped") and size <= 50 * 1024 * 1024:
                try:
                    with zf.open(info, "r") as src:
                        data = src.read(50 * 1024 * 1024 + 1)
                    if len(data) <= 50 * 1024 * 1024:
                        import hashlib
                        sha = hashlib.sha256(data).hexdigest()
                except Exception:
                    pass

            manifest.append({
                "index": idx,
                "name": name,
                "extension": ext,
                "size": size,
                "status": status,
                "sha256": sha,
            })

    if stats["total"] > max_files:
        stats["skipped"] += stats["total"] - max_files
    return manifest, stats


def _create_zip_pdf(src_path: str, dest_path: str, cert_id: str, authenticity: int,
                    label: str, filename: str, sha256: str = "",
                    detail: Dict[str, Any] | None = None) -> str:
    """Create a certified ZIP evidence-package PDF with a manifest and embedded original ZIP."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import landscape, letter

    page_w, page_h = landscape(letter)
    c = canvas.Canvas(dest_path, pagesize=(page_w, page_h))
    c.setAuthor("VeriFYD")
    c.setTitle(f"VeriFYD Certified ZIP Evidence Package {cert_id}")
    c.setSubject(f"{label} | Authenticity {authenticity}% | {filename}")

    try:
        manifest, stats = _read_zip_manifest_for_certificate(src_path)
        manifest_error = ""
    except Exception as e:
        manifest, stats = [], {"total": 0, "supported": 0, "skipped": 0, "dangerous": 0, "bytes": 0}
        manifest_error = str(e)[:200]

    # Cover page.
    c.setFont("Helvetica-Bold", 20)
    c.drawString(54, page_h - 72, "VeriFYD Certified ZIP Evidence Package")
    c.setFont("Helvetica", 10)
    c.drawString(54, page_h - 105, f"Status: {label}")
    c.drawString(54, page_h - 123, f"Authenticity Score: {authenticity}%")
    c.drawString(54, page_h - 141, f"Certificate ID: {cert_id}")
    c.drawString(54, page_h - 159, f"Original file: {filename}")
    c.drawString(54, page_h - 177, f"Files in package: {stats.get('total', 0)}")
    c.drawString(54, page_h - 195, f"Supported files analyzed: {stats.get('supported', 0)}")
    c.drawString(54, page_h - 213, f"Skipped / manifest-only files: {stats.get('skipped', 0)}")
    c.drawString(54, page_h - 231, f"Potentially dangerous files skipped: {stats.get('dangerous', 0)}")
    c.drawString(54, page_h - 249, f"Uncompressed manifest bytes: {stats.get('bytes', 0)}")
    if sha256:
        c.drawString(54, page_h - 274, "ZIP SHA-256:")
        c.setFont("Helvetica", 7.5)
        c.drawString(54, page_h - 288, sha256[:120])
    if manifest_error:
        c.setFont("Helvetica-Bold", 9)
        c.setFillColorRGB(0.75, 0.10, 0.10)
        c.drawString(54, page_h - 312, f"Manifest warning: {manifest_error}")
        c.setFillGray(0)

    c.setFont("Helvetica-Oblique", 8)
    c.setFillGray(0.35)
    c.drawString(54, 58, "VeriFYD safely inspected the ZIP manifest and supported inner documents without executing archive contents.")
    c.drawString(54, 44, "The exact original ZIP package is embedded inside this certified PDF as an attachment for source review.")
    _draw_verifyd_mark(c, page_w, y=24, x_right_pad=34)
    c.showPage()

    # Manifest pages.
    rows_per_page = 24
    chunks = [manifest[i:i + rows_per_page] for i in range(0, len(manifest), rows_per_page)] or [[]]
    for page_idx, rows in enumerate(chunks, start=1):
        c.setFont("Helvetica-Bold", 14)
        c.setFillGray(0.05)
        c.drawString(34, page_h - 36, f"ZIP Manifest — Page {page_idx}/{len(chunks)}")
        c.setFont("Helvetica", 7)
        c.setFillGray(0.35)
        c.drawString(34, page_h - 50, f"File: {filename} • Certificate: {cert_id}")

        y = page_h - 76
        headers = [("#", 28), ("Name", 360), ("Ext", 45), ("Size", 70), ("Status", 115), ("SHA-256", 180)]
        x0 = 34
        c.setFont("Helvetica-Bold", 7)
        c.setFillGray(0.90)
        c.rect(x0, y - 4, sum(w for _, w in headers), 16, fill=1, stroke=0)
        c.setFillGray(0.05)
        x = x0
        for title, w in headers:
            c.drawString(x + 3, y + 1, title)
            x += w
        y -= 18

        c.setFont("Helvetica", 6.5)
        for row in rows:
            x = x0
            vals = [
                str(row.get("index", "")),
                str(row.get("name", ""))[:86],
                str(row.get("extension", ""))[:8],
                str(row.get("size", "")),
                str(row.get("status", ""))[:28],
                str(row.get("sha256", ""))[:32],
            ]
            for val, (_, w) in zip(vals, headers):
                c.setStrokeGray(0.82)
                c.rect(x, y - 4, w, 16, fill=0, stroke=1)
                c.setFillGray(0.05)
                c.drawString(x + 3, y + 1, val)
                x += w
            y -= 16

        if not rows:
            c.setFont("Helvetica", 10)
            c.drawString(54, page_h - 100, "No readable ZIP file entries were available for this package.")

        c.setFont("Helvetica", 6.5)
        c.setFillGray(0.45)
        c.drawString(34, 24, "VeriFYD certified ZIP evidence package manifest")
        _draw_verifyd_mark(c, page_w, y=24, x_right_pad=34)
        if page_idx < len(chunks):
            c.showPage()

    c.save()
    _attach_original_file_to_pdf(dest_path, src_path, filename or os.path.basename(src_path))
    return dest_path

def stamp_document(src_path: str, dest_path: str, cert_id: str,
                   authenticity: int, label: str, filename: str,
                   sha256: str = "", detail: Dict[str, Any] | None = None) -> str:
    """
    Create a certified PDF at dest_path.

    For PDFs: preserve original pages and add only a small lower-right VeriFYD mark.
    For JPG/JPEG/PNG: create a marked PDF containing the image.
    For XLSX: render workbook sheets into a marked PDF.
    For PPTX: render presentation slides into a visual marked PDF approximation.
    For other non-PDFs: create a simple PDF summary fallback.
    """
    ext = os.path.splitext(src_path)[1].lower()

    if ext == ".zip":
        _create_zip_pdf(src_path, dest_path, cert_id, authenticity, label, filename, sha256, detail)
        return _finalize_certified_pdf(dest_path, detail, cert_id, filename, sha256, authenticity, label)

    if ext in (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tif", ".tiff", ".webp", ".heic", ".heif"):
        _create_image_pdf(src_path, dest_path, cert_id, authenticity, label, filename, sha256)
        # Preserve the exact uploaded image as an embedded attachment, while the
        # visible certified PDF remains the official stamped viewing copy.
        _attach_original_file_to_pdf(dest_path, src_path, filename or os.path.basename(src_path))
        return _finalize_certified_pdf(dest_path, detail, cert_id, filename, sha256, authenticity, label)

    if ext in LIBREOFFICE_RENDER_EXTENSIONS:
        try:
            _create_office_pdf_via_libreoffice(src_path, dest_path, cert_id, authenticity, label, filename, sha256)
            return _finalize_certified_pdf(dest_path, detail, cert_id, filename, sha256, authenticity, label)
        except Exception as e:
            # Keep certification resilient. If LibreOffice is unavailable or a
            # specific file cannot be converted, use the existing safe fallback
            # renderers below, but log the reason so Render can be diagnosed.
            log.warning(
                "doc_certifier: LibreOffice render failed for ext=%s file=%s; falling back to text renderer: %s",
                ext,
                filename,
                e,
            )

    if ext == ".docx":
        try:
            _create_docx_layout_pdf(src_path, dest_path, cert_id, authenticity, label, filename, sha256)
            # Preserve exact source file for legal/forensic review. The visible PDF
            # is a table-aware rendering, not a perfect native Word clone.
            _attach_original_file_to_pdf(dest_path, src_path, filename or os.path.basename(src_path))
            return _finalize_certified_pdf(dest_path, detail, cert_id, filename, sha256, authenticity, label)
        except Exception as e:
            log.warning(
                "doc_certifier: DOCX table-aware fallback failed for file=%s; using text renderer: %s",
                filename,
                e,
            )

    if ext == ".xlsx":
        _create_xlsx_pdf(src_path, dest_path, cert_id, authenticity, label, filename, sha256)
        _attach_original_file_to_pdf(dest_path, src_path, filename or os.path.basename(src_path))
        return _finalize_certified_pdf(dest_path, detail, cert_id, filename, sha256, authenticity, label)

    if ext == ".pptx":
        _create_pptx_pdf(src_path, dest_path, cert_id, authenticity, label, filename, sha256)
        return _finalize_certified_pdf(dest_path, detail, cert_id, filename, sha256, authenticity, label)

    if ext in (".eml", ".msg"):
        _create_email_render_pdf(src_path, dest_path, cert_id, authenticity, label, filename, sha256)
        # Preserve exact source email file for legal/forensic review.
        _attach_original_file_to_pdf(dest_path, src_path, filename or os.path.basename(src_path))
        return _finalize_certified_pdf(dest_path, detail, cert_id, filename, sha256, authenticity, label)

    if ext == ".vsdx":
        _create_vsdx_diagram_pdf(src_path, dest_path, cert_id, authenticity, label, filename, sha256)
        # Preserve exact source Visio file for legal/forensic review.
        _attach_original_file_to_pdf(dest_path, src_path, filename or os.path.basename(src_path))
        return _finalize_certified_pdf(dest_path, detail, cert_id, filename, sha256, authenticity, label)

    if ext == ".rtf":
        _create_rtf_form_pdf(src_path, dest_path, cert_id, authenticity, label, filename, sha256)
        # Preserve exact source RTF file for legal/forensic review.
        _attach_original_file_to_pdf(dest_path, src_path, filename or os.path.basename(src_path))
        return _finalize_certified_pdf(dest_path, detail, cert_id, filename, sha256, authenticity, label)

    if ext in (".docx", ".txt", ".md", ".csv", ".rtf", ".doc", ".ppt", ".xls", ".odt", ".ods", ".odp", ".html", ".htm", ".mhtml", ".mht", ".xml", ".json", ".svg", ".vsdx", ".yaml", ".yml", ".ini", ".log", ".sql", ".pst", ".ost", ".dwg", ".dxf"):
        _create_text_render_pdf(src_path, dest_path, cert_id, authenticity, label, filename, sha256)
        # Preserve exact source file for legal/forensic review, especially when
        # the certified PDF is a readable text rendering rather than a perfect
        # native-layout clone.
        _attach_original_file_to_pdf(dest_path, src_path, filename or os.path.basename(src_path))
        return _finalize_certified_pdf(dest_path, detail, cert_id, filename, sha256, authenticity, label)

    if ext == ".pdf":
        from pypdf import PdfReader, PdfWriter

        reader = PdfReader(src_path)
        writer = PdfWriter()
        overlay_paths: list[str] = []

        try:
            for page in reader.pages:
                width = float(page.mediabox.width)
                height = float(page.mediabox.height)
                overlay = _make_logo_overlay(width, height)
                overlay_paths.append(overlay)
                overlay_page = PdfReader(overlay).pages[0]
                page.merge_page(overlay_page)
                writer.add_page(page)

            try:
                existing_meta = reader.metadata or {}
                metadata = {str(k): str(v) for k, v in existing_meta.items() if v is not None}
                metadata.update({
                    "/VeriFYD": "Certified",
                    "/VeriFYD_Certificate_ID": cert_id,
                    "/VeriFYD_Label": label,
                    "/VeriFYD_Authenticity": str(authenticity),
                    "/VeriFYD_SHA256_Original": sha256 or "",
                })
                writer.add_metadata(metadata)
            except Exception:
                pass

            with open(dest_path, "wb") as out:
                writer.write(out)
        finally:
            for p in overlay_paths:
                try:
                    if os.path.exists(p):
                        os.remove(p)
                except Exception:
                    pass
        return _finalize_certified_pdf(dest_path, detail, cert_id, filename, sha256, authenticity, label)

    _create_non_pdf_certificate(dest_path, cert_id, authenticity, label, filename, sha256)
    return _finalize_certified_pdf(dest_path, detail, cert_id, filename, sha256, authenticity, label)
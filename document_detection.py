# ============================================================
#  VeriFYD — document_detection.py  v1
#
#  MVP document authenticity engine for PDF / DOCX / TXT.
#  Returns the same tuple signature used by video/photo detection:
#      (authenticity:int, label:str, detail:dict)
#
#  Engines:
#    1. File/hash + metadata inspection
#    2. Text extraction + statistical writing signals
#    3. GPT semantic document analysis, if OPENAI_API_KEY is set
#    4. PDF embedded-image inventory / suspicious producer flags
# ============================================================

from __future__ import annotations

import hashlib
import json
import logging
import math
import os
import re
import shutil
import subprocess
from collections import Counter
from datetime import datetime
from typing import Any, Dict, List, Tuple

log = logging.getLogger("verifyd.document_detection")

THRESHOLD_REAL = 55
THRESHOLD_UNDETERMINED = 40
MAX_GPT_CHARS = 12000

AI_PRODUCER_KEYWORDS = (
    "chatgpt", "openai", "sora", "dall-e", "dalle", "midjourney",
    "stable diffusion", "stability ai", "runway", "kling", "pika",
    "firefly", "adobe firefly", "gemini", "bard", "claude", "anthropic",
    "canva", "copilot", "microsoft designer", "gamma", "beautiful.ai",
)

COMMON_GENERATOR_KEYWORDS = (
    "microsoft word", "word", "pages", "google docs", "adobe acrobat",
    "preview", "libreoffice", "powerpoint", "excel", "quartz pdfcontext",
)

AI_TEXT_PATTERNS = (
    "in today's", "ever-evolving", "delve into", "it's important to note",
    "in conclusion", "as an ai", "unlock the potential", "seamlessly",
    "robust solution", "cutting-edge", "game-changer", "leverage",
    "streamline", "comprehensive approach", "dynamic landscape",
)


def _sha256_file(path: str) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def _run_exiftool_metadata(path: str) -> Tuple[Dict[str, Any], Dict[str, Any]]:
    """
    Optional ExifTool metadata pass.

    ExifTool is intentionally optional: if the binary is not installed on Render,
    VeriFYD keeps using the existing Python extraction stack. When available,
    this adds a stronger forensic metadata view across PDF, Office, images,
    CAD/evidence files, and archive-like formats.
    """
    tool_info: Dict[str, Any] = {
        "name": "ExifTool",
        "available": False,
        "used": False,
        "status": "not_available",
    }
    exe = shutil.which("exiftool")
    if not exe:
        return {}, tool_info

    tool_info["available"] = True
    try:
        result = subprocess.run(
            [exe, "-j", "-a", "-G1", "-s", path],
            capture_output=True,
            text=True,
            timeout=25,
        )
        if result.returncode != 0:
            tool_info["status"] = f"error:{(result.stderr or '')[:120]}"
            return {}, tool_info
        parsed = json.loads(result.stdout or "[]")
        if not parsed or not isinstance(parsed, list) or not isinstance(parsed[0], dict):
            tool_info["status"] = "no_metadata_returned"
            return {}, tool_info
        raw_md = parsed[0]
        cleaned: Dict[str, Any] = {}
        for k, v in raw_md.items():
            if k in ("SourceFile",):
                continue
            if isinstance(v, (str, int, float, bool)):
                cleaned[str(k)] = str(v)[:1000]
            elif isinstance(v, list):
                cleaned[str(k)] = ", ".join(str(x) for x in v[:20])[:1000]
            elif isinstance(v, dict):
                cleaned[str(k)] = json.dumps(v, ensure_ascii=False)[:1000]
        tool_info.update({
            "used": True,
            "status": "ok",
            "field_count": len(cleaned),
        })
        return cleaned, tool_info
    except Exception as e:
        tool_info["status"] = f"exception:{str(e)[:120]}"
        return {}, tool_info


def _merge_exiftool_metadata(path: str, meta: Dict[str, Any]) -> Dict[str, Any]:
    """Merge optional ExifTool metadata into the existing meta object safely."""
    if not isinstance(meta, dict):
        meta = {"type": "document", "pages": 0, "embedded_images": 0, "metadata": {}}
    existing = dict(meta.get("metadata") or {})
    exif_md, tool_info = _run_exiftool_metadata(path)
    meta["external_metadata_tool"] = tool_info

    if exif_md:
        # Keep original parser keys untouched; add ExifTool keys under a prefix
        # to avoid changing current scoring behavior unexpectedly.
        for k, v in exif_md.items():
            prefixed = f"exiftool_{k}"
            if prefixed not in existing:
                existing[prefixed] = v

        # Promote common date/software fields only when the primary parser did
        # not already provide them. This improves the risk report without
        # overwriting pypdf/openpyxl/python-docx values.
        promotion_map = {
            "CreateDate": "created",
            "PDF:CreateDate": "created",
            "EXIF:CreateDate": "created",
            "File:FileModifyDate": "modified",
            "ModifyDate": "modified",
            "PDF:ModifyDate": "modified",
            "XMP:ModifyDate": "modified",
            "Producer": "producer",
            "PDF:Producer": "producer",
            "Creator": "creator",
            "XMP:CreatorTool": "creator",
            "Author": "author",
            "XMP:Author": "author",
            "Software": "software",
            "EXIF:Software": "software",
        }
        lowered_existing = {str(k).lower().replace("_", ""): k for k in existing.keys()}
        for src_key, dst_key in promotion_map.items():
            if dst_key.lower().replace("_", "") in lowered_existing:
                continue
            if src_key in exif_md and str(exif_md[src_key]).strip():
                existing[dst_key] = str(exif_md[src_key]).strip()[:500]

    meta["metadata"] = existing
    return meta


def _safe_int(value: Any, default: int = 0) -> int:
    try:
        return int(value)
    except Exception:
        return default


def _clip_text(text: str, limit: int = MAX_GPT_CHARS) -> str:
    text = re.sub(r"\s+", " ", text or "").strip()
    if len(text) <= limit:
        return text
    head = text[: limit // 2]
    tail = text[-limit // 2 :]
    return head + "\n\n[...middle omitted for analysis length...]\n\n" + tail


def _parse_pdf_date(value: Any) -> str:
    """Convert a PDF D:YYYYMMDDHHmmSS date to a readable ISO-like string."""
    if not value:
        return ""
    s = str(value).strip()
    if s.startswith("D:"):
        raw = s[2:16]
        try:
            return datetime.strptime(raw, "%Y%m%d%H%M%S").isoformat()
        except Exception:
            return s
    return s


def _read_pdf(path: str) -> Tuple[str, Dict[str, Any]]:
    meta: Dict[str, Any] = {"type": "pdf", "pages": 0, "embedded_images": 0, "metadata": {}}
    text_parts: List[str] = []
    try:
        from pypdf import PdfReader
    except Exception as e:
        raise RuntimeError("Missing dependency: pypdf. Add pypdf>=4.0.0 to requirements.txt") from e

    reader = PdfReader(path)
    meta["pages"] = len(reader.pages)

    raw_meta = reader.metadata or {}
    parsed_meta = {}
    for k, v in raw_meta.items():
        key = str(k).lstrip("/")
        parsed_meta[key] = _parse_pdf_date(v) if "date" in key.lower() else str(v)
    meta["metadata"] = parsed_meta

    # Text extraction
    for i, page in enumerate(reader.pages[:50]):
        try:
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(page_text)
        except Exception as e:
            log.debug("PDF text extraction failed page=%d: %s", i, e)

    # Embedded-image count, best-effort only
    image_count = 0
    for page in reader.pages[:50]:
        try:
            resources = page.get("/Resources", {}) or {}
            xobj = resources.get("/XObject", {}) or {}
            if hasattr(xobj, "get_object"):
                xobj = xobj.get_object()
            for _, obj in xobj.items():
                try:
                    resolved = obj.get_object() if hasattr(obj, "get_object") else obj
                    if resolved.get("/Subtype") == "/Image":
                        image_count += 1
                except Exception:
                    continue
        except Exception:
            continue
    meta["embedded_images"] = image_count
    return "\n".join(text_parts), meta


def _read_docx(path: str) -> Tuple[str, Dict[str, Any]]:
    meta: Dict[str, Any] = {"type": "docx", "pages": 0, "embedded_images": 0, "metadata": {}}
    try:
        from docx import Document
    except Exception as e:
        raise RuntimeError("Missing dependency: python-docx. Add python-docx>=1.1.0 to requirements.txt") from e

    doc = Document(path)
    text_parts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    # Tables often contain invoice/resume values
    for table in doc.tables:
        for row in table.rows:
            vals = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if vals:
                text_parts.append(" | ".join(vals))

    props = doc.core_properties
    meta["metadata"] = {
        "author": props.author or "",
        "last_modified_by": props.last_modified_by or "",
        "created": props.created.isoformat() if props.created else "",
        "modified": props.modified.isoformat() if props.modified else "",
        "title": props.title or "",
        "subject": props.subject or "",
        "keywords": props.keywords or "",
        "comments": props.comments or "",
        "revision": str(props.revision or ""),
    }
    # Inline shapes gives a decent embedded image count
    try:
        meta["embedded_images"] = len(doc.inline_shapes)
    except Exception:
        meta["embedded_images"] = 0
    return "\n".join(text_parts), meta


def _read_txt(path: str) -> Tuple[str, Dict[str, Any]]:
    meta = {"type": "txt", "pages": 0, "embedded_images": 0, "metadata": {}}
    with open(path, "rb") as f:
        data = f.read(2_000_000)
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(enc, errors="replace"), meta
        except Exception:
            continue
    return data.decode("latin-1", errors="replace"), meta




def _rtf_regex_fallback(raw: str) -> str:
    """Table-preserving RTF cleanup used when striprtf over-strips forms."""
    raw = raw or ""
    # Drop only the noisiest embedded binary/image groups. Keep table/text control flow.
    raw = re.sub(r"{\\\*?\\pict[^{}]*(?:{[^{}]*}[^{}]*)*}", " ", raw, flags=re.DOTALL)
    raw = re.sub(r"{\\object[^{}]*(?:{[^{}]*}[^{}]*)*}", " ", raw, flags=re.DOTALL)

    replacements = {
        "\\par": "\n",
        "\\line": "\n",
        "\\tab": "\t",
        "\\cell": "|",
        "\\row": "\n",
    }
    for src, dst in replacements.items():
        raw = raw.replace(src, dst)

    def _hex_to_char(match):
        try:
            return bytes.fromhex(match.group(1)).decode("cp1252", errors="replace")
        except Exception:
            return " "

    raw = re.sub(r"\\'([0-9a-fA-F]{2})", _hex_to_char, raw)
    raw = re.sub(r"\\u(-?\d+)\??", lambda m: chr(int(m.group(1)) % 65536), raw)
    raw = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", raw)
    raw = raw.replace("{", " ").replace("}", " ")
    raw = re.sub(r"\x00+", " ", raw)
    raw = re.sub(r"[ \t]+", " ", raw)
    raw = re.sub(r" *\| *", "|", raw)
    raw = re.sub(r"\n\s*\n\s*\n+", "\n\n", raw)
    return raw.strip()


def _strip_rtf_markup(raw: str) -> str:
    """RTF-to-text cleanup with a fallback that preserves forms/tables."""
    regex_text = _rtf_regex_fallback(raw)

    try:
        from striprtf.striprtf import rtf_to_text
        text = rtf_to_text(raw or "")
        text = re.sub(r"\x00+", " ", text)
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n\s*\n\s*\n+", "\n\n", text)
        text = text.strip()
        # striprtf can over-strip table-based government forms. If its output
        # is much shorter, keep the older table-preserving fallback.
        if len(text) >= 500 or len(text) >= max(120, len(regex_text) * 0.35):
            return text
    except Exception:
        pass

    return regex_text


def _read_rtf(path: str) -> Tuple[str, Dict[str, Any]]:
    """Best-effort RTF text extraction using safe built-in parsing."""
    meta = {"type": "rtf", "pages": 0, "embedded_images": 0, "metadata": {}}
    with open(path, "rb") as f:
        data = f.read(2_000_000)
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            raw = data.decode(enc, errors="replace")
            return _strip_rtf_markup(raw), meta
        except Exception:
            continue
    return _strip_rtf_markup(data.decode("latin-1", errors="replace")), meta


def _read_eml(path: str) -> Tuple[str, Dict[str, Any]]:
    """Best-effort EML email extraction: headers, text/plain body, and attachment inventory."""
    meta: Dict[str, Any] = {"type": "eml", "pages": 1, "embedded_images": 0, "metadata": {}}
    try:
        from email import policy
        from email.parser import BytesParser
    except Exception as e:
        raise RuntimeError("Missing Python email parser support.") from e

    with open(path, "rb") as f:
        msg = BytesParser(policy=policy.default).parse(f)

    headers = {
        "from": str(msg.get("From", ""))[:500],
        "to": str(msg.get("To", ""))[:500],
        "cc": str(msg.get("Cc", ""))[:500],
        "subject": str(msg.get("Subject", ""))[:500],
        "date": str(msg.get("Date", ""))[:200],
        "message_id": str(msg.get("Message-ID", ""))[:300],
        "return_path": str(msg.get("Return-Path", ""))[:300],
        "reply_to": str(msg.get("Reply-To", ""))[:300],
        "content_type": str(msg.get_content_type()),
    }

    text_parts: List[str] = []
    attachment_names: List[str] = []
    embedded_images = 0

    def _payload_to_text(part: Any) -> str:
        try:
            payload = part.get_content()
            if isinstance(payload, str):
                return payload
        except Exception:
            pass
        try:
            data = part.get_payload(decode=True) or b""
            charset = part.get_content_charset() or "utf-8"
            return data.decode(charset, errors="replace")
        except Exception:
            return ""

    if msg.is_multipart():
        for part in msg.walk():
            ctype = part.get_content_type()
            disposition = str(part.get_content_disposition() or "")
            filename = part.get_filename()
            if filename:
                attachment_names.append(str(filename)[:200])
            if ctype.startswith("image/"):
                embedded_images += 1
            if ctype == "text/plain" and disposition != "attachment":
                body = _payload_to_text(part).strip()
                if body:
                    text_parts.append(body)
            elif ctype == "text/html" and disposition != "attachment" and not text_parts:
                html = _payload_to_text(part)
                text = re.sub(r"<[^>]+>", " ", html)
                text = re.sub(r"\s+", " ", text).strip()
                if text:
                    text_parts.append(text)
    else:
        body = _payload_to_text(msg).strip()
        if body:
            text_parts.append(body)

    headers["attachment_count"] = str(len(attachment_names))
    headers["attachment_names"] = ", ".join(attachment_names[:20])
    meta["metadata"] = headers
    meta["embedded_images"] = embedded_images

    header_text = "\n".join(f"{k}: {v}" for k, v in headers.items() if v)
    body_text = "\n\n".join(text_parts)
    return (header_text + "\n\n" + body_text).strip(), meta



def _clean_ole_extracted_text(data: bytes, limit: int = 2_000_000) -> str:
    """Extract readable ASCII/UTF-16 text fragments from legacy OLE Office binaries."""
    fragments: List[str] = []
    sample = data[:limit]

    for enc in ("utf-16-le", "latin-1"):
        try:
            decoded = sample.decode(enc, errors="ignore")
        except Exception:
            continue
        decoded = decoded.replace("\x00", " ")
        decoded = re.sub(r"[\x01-\x08\x0b\x0c\x0e-\x1f]+", " ", decoded)
        decoded = re.sub(r"[ \t]+", " ", decoded)
        pieces = re.findall(r"[A-Za-z0-9][A-Za-z0-9\s\.,;:\-_/@$%#&()\[\]{}'\"!?]{4,}", decoded)
        fragments.extend(p.strip() for p in pieces if p and len(p.strip()) >= 5)

    seen = set()
    cleaned: List[str] = []
    running = 0
    for frag in fragments:
        compact = re.sub(r"\s+", " ", frag).strip()
        if len(compact) < 5:
            continue
        key = compact[:160].lower()
        if key in seen:
            continue
        seen.add(key)
        cleaned.append(compact)
        running += len(compact)
        if running > limit:
            break

    return "\n".join(cleaned)


def _read_ole_office(path: str, doc_type: str) -> Tuple[str, Dict[str, Any]]:
    """Best-effort metadata/text extraction for legacy .doc and .ppt OLE files."""
    meta: Dict[str, Any] = {"type": doc_type, "pages": 0, "embedded_images": 0, "metadata": {}}
    try:
        import olefile
    except Exception as e:
        raise RuntimeError("Missing dependency: olefile. Add olefile>=0.47 to requirements.txt") from e

    stream_names: List[str] = []
    text_parts: List[str] = []

    try:
        with olefile.OleFileIO(path) as ole:
            stream_names = ["/".join(s) for s in ole.listdir()]
            meta["metadata"] = {
                "ole_stream_count": str(len(stream_names)),
                "ole_streams": ", ".join(stream_names[:30]),
            }
            meta["embedded_images"] = sum(
                1 for name in stream_names
                if any(token in name.lower() for token in ("picture", "image", "jpeg", "png", "wmf", "emf"))
            )
            for stream in ole.listdir()[:80]:
                try:
                    raw = ole.openstream(stream).read(750_000)
                    txt = _clean_ole_extracted_text(raw, limit=750_000)
                    if txt:
                        text_parts.append(txt)
                except Exception:
                    continue
    except Exception as e:
        log.warning("OLE legacy Office extraction failed for %s: %s", path, e)

    if not text_parts:
        try:
            with open(path, "rb") as fh:
                text_parts.append(_clean_ole_extracted_text(fh.read(2_000_000)))
        except Exception:
            pass

    return "\n".join(x for x in text_parts if x).strip(), meta


def _read_doc(path: str) -> Tuple[str, Dict[str, Any]]:
    """Best-effort legacy Microsoft Word .doc extraction."""
    return _read_ole_office(path, "doc")


def _read_ppt(path: str) -> Tuple[str, Dict[str, Any]]:
    """Best-effort legacy Microsoft PowerPoint .ppt extraction."""
    return _read_ole_office(path, "ppt")


def _read_xls(path: str) -> Tuple[str, Dict[str, Any]]:
    """Best-effort legacy Excel .xls extraction using xlrd."""
    meta: Dict[str, Any] = {"type": "xls", "pages": 0, "embedded_images": 0, "metadata": {}}
    try:
        import xlrd
    except Exception as e:
        raise RuntimeError("Missing dependency: xlrd. Add xlrd==1.2.0 to requirements.txt") from e

    text_parts: List[str] = []
    book = xlrd.open_workbook(path, on_demand=True)
    sheet_names = book.sheet_names()
    meta["pages"] = len(sheet_names)
    meta["metadata"] = {
        "sheet_count": str(len(sheet_names)),
        "sheet_names": ", ".join(sheet_names[:30]),
    }

    rows_seen = 0
    max_rows_total = 10000
    for sheet_name in sheet_names:
        try:
            sh = book.sheet_by_name(sheet_name)
        except Exception:
            continue
        text_parts.append(f"Worksheet: {sheet_name}")
        for r_idx in range(sh.nrows):
            vals: List[str] = []
            for c_idx in range(min(sh.ncols, 50)):
                try:
                    val = sh.cell_value(r_idx, c_idx)
                except Exception:
                    continue
                if val in (None, ""):
                    continue
                if isinstance(val, float) and val.is_integer():
                    val = int(val)
                text = str(val).strip()
                if text:
                    vals.append(text[:500])
            if vals:
                text_parts.append(" | ".join(vals))
            rows_seen += 1
            if rows_seen >= max_rows_total:
                text_parts.append("[...legacy workbook truncated for analysis length...]")
                break
        if rows_seen >= max_rows_total:
            break

    try:
        book.release_resources()
    except Exception:
        pass
    return "\n".join(text_parts), meta


def _read_msg(path: str) -> Tuple[str, Dict[str, Any]]:
    """Best-effort Outlook .msg extraction: headers, body, and attachment inventory."""
    meta: Dict[str, Any] = {"type": "msg", "pages": 1, "embedded_images": 0, "metadata": {}}
    try:
        import extract_msg
    except Exception as e:
        raise RuntimeError("Missing dependency: extract-msg. Add extract-msg>=0.48.0 to requirements.txt") from e

    msg = extract_msg.Message(path)
    try:
        attachments = list(getattr(msg, "attachments", []) or [])
        attachment_names = []
        embedded_images = 0
        for att in attachments[:50]:
            name = str(getattr(att, "longFilename", "") or getattr(att, "shortFilename", "") or "")
            if name:
                attachment_names.append(name[:200])
                if os.path.splitext(name.lower())[1] in (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tif", ".tiff"):
                    embedded_images += 1
        headers = {
            "from": str(getattr(msg, "sender", "") or "")[:500],
            "to": str(getattr(msg, "to", "") or "")[:500],
            "cc": str(getattr(msg, "cc", "") or "")[:500],
            "subject": str(getattr(msg, "subject", "") or "")[:500],
            "date": str(getattr(msg, "date", "") or "")[:200],
            "message_id": str(getattr(msg, "messageId", "") or "")[:300],
            "attachment_count": str(len(attachments)),
            "attachment_names": ", ".join(attachment_names[:20]),
        }
        body = str(getattr(msg, "body", "") or getattr(msg, "htmlBody", "") or "")
        if body and "<" in body and ">" in body:
            body = re.sub(r"<[^>]+>", " ", body)
            body = re.sub(r"\s+", " ", body).strip()
        meta["metadata"] = headers
        meta["embedded_images"] = embedded_images
        header_text = "\n".join(f"{k}: {v}" for k, v in headers.items() if v)
        return (header_text + "\n\n" + body).strip(), meta
    finally:
        try:
            msg.close()
        except Exception:
            pass

def _read_xlsx(path: str) -> Tuple[str, Dict[str, Any]]:
    """Best-effort XLSX text and workbook metadata extraction."""
    meta: Dict[str, Any] = {"type": "xlsx", "pages": 0, "embedded_images": 0, "metadata": {}}
    try:
        from openpyxl import load_workbook
    except Exception as e:
        raise RuntimeError("Missing dependency: openpyxl. Add openpyxl>=3.1.0 to requirements.txt") from e

    wb = load_workbook(path, data_only=True, read_only=True)
    text_parts: List[str] = []
    sheet_names = list(wb.sheetnames)

    try:
        props = wb.properties
        meta["metadata"] = {
            "creator": props.creator or "",
            "last_modified_by": props.lastModifiedBy or "",
            "created": props.created.isoformat() if props.created else "",
            "modified": props.modified.isoformat() if props.modified else "",
            "title": props.title or "",
            "subject": props.subject or "",
            "keywords": props.keywords or "",
            "category": props.category or "",
            "description": props.description or "",
            "sheet_count": str(len(sheet_names)),
            "sheet_names": ", ".join(sheet_names[:20]),
        }
    except Exception:
        meta["metadata"] = {
            "sheet_count": str(len(sheet_names)),
            "sheet_names": ", ".join(sheet_names[:20]),
        }

    max_rows_total = 10000
    rows_seen = 0
    for ws in wb.worksheets:
        text_parts.append(f"Worksheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            vals = []
            for cell in row:
                if cell is None:
                    continue
                value = str(cell).strip()
                if value:
                    vals.append(value[:500])
            if vals:
                text_parts.append(" | ".join(vals))
            rows_seen += 1
            if rows_seen >= max_rows_total:
                text_parts.append("[...workbook truncated for analysis length...]")
                break
        if rows_seen >= max_rows_total:
            break

    try:
        wb.close()
    except Exception:
        pass

    meta["pages"] = len(sheet_names)
    meta["embedded_images"] = 0
    return "\n".join(text_parts), meta



def _read_pptx(path: str) -> Tuple[str, Dict[str, Any]]:
    """Best-effort PPTX slide text and presentation metadata extraction."""
    meta: Dict[str, Any] = {"type": "pptx", "pages": 0, "embedded_images": 0, "metadata": {}}
    try:
        from pptx import Presentation
    except Exception as e:
        raise RuntimeError("Missing dependency: python-pptx. Add python-pptx>=0.6.23 to requirements.txt") from e

    prs = Presentation(path)
    text_parts: List[str] = []
    embedded_images = 0

    try:
        props = prs.core_properties
        meta["metadata"] = {
            "author": props.author or "",
            "last_modified_by": props.last_modified_by or "",
            "created": props.created.isoformat() if props.created else "",
            "modified": props.modified.isoformat() if props.modified else "",
            "title": props.title or "",
            "subject": props.subject or "",
            "keywords": props.keywords or "",
            "comments": props.comments or "",
            "revision": str(props.revision or ""),
            "category": props.category or "",
        }
    except Exception:
        meta["metadata"] = {}

    max_slides = 80
    for slide_idx, slide in enumerate(prs.slides, start=1):
        if slide_idx > max_slides:
            text_parts.append("[...presentation truncated for analysis length...]")
            break

        slide_text: List[str] = []
        for shape in slide.shapes:
            try:
                if getattr(shape, "has_text_frame", False) and shape.text and shape.text.strip():
                    slide_text.append(shape.text.strip())
            except Exception:
                pass

            try:
                if getattr(shape, "shape_type", None) == 13:  # MSO_SHAPE_TYPE.PICTURE
                    embedded_images += 1
            except Exception:
                pass

            # Tables can hold quote/proposal values.
            try:
                if getattr(shape, "has_table", False):
                    table = shape.table
                    for row in table.rows:
                        vals = []
                        for cell in row.cells:
                            value = (cell.text or "").strip()
                            if value:
                                vals.append(value[:300])
                        if vals:
                            slide_text.append(" | ".join(vals))
            except Exception:
                pass

        text_parts.append(f"Slide {slide_idx}:")
        if slide_text:
            text_parts.extend(slide_text)
        else:
            text_parts.append("[No extractable slide text]")

    meta["pages"] = len(prs.slides)
    meta["embedded_images"] = embedded_images
    return "\n".join(text_parts), meta


def _xml_text_content(elem: Any) -> str:
    """Return collapsed text content from an XML element."""
    parts: List[str] = []
    try:
        for t in elem.itertext():
            if t and str(t).strip():
                parts.append(str(t).strip())
    except Exception:
        pass
    return re.sub(r"\s+", " ", " ".join(parts)).strip()


def _read_odf(path: str) -> Tuple[str, Dict[str, Any]]:
    """Best-effort OpenDocument extraction for ODT/ODS/ODP files."""
    ext = os.path.splitext(path)[1].lower()
    odf_type = {".odt": "odt", ".ods": "ods", ".odp": "odp"}.get(ext, "odf")
    meta: Dict[str, Any] = {"type": odf_type, "pages": 0, "embedded_images": 0, "metadata": {}}
    text_parts: List[str] = []

    try:
        import zipfile
        import xml.etree.ElementTree as ET
    except Exception as e:
        raise RuntimeError("Missing Python ZIP/XML support for OpenDocument parsing.") from e

    def _local(tag: str) -> str:
        return str(tag).split("}")[-1].lower()

    try:
        with zipfile.ZipFile(path) as zf:
            names = set(zf.namelist())
            try:
                if "mimetype" in names:
                    meta["metadata"]["mimetype"] = zf.read("mimetype").decode("utf-8", errors="replace")[:300]
            except Exception:
                pass

            if "meta.xml" in names:
                try:
                    root = ET.fromstring(zf.read("meta.xml"))
                    md: Dict[str, str] = dict(meta.get("metadata") or {})
                    for elem in root.iter():
                        name = _local(elem.tag)
                        if name in ("generator", "initial-creator", "creator", "creation-date", "date", "editing-duration", "document-statistic"):
                            value = _xml_text_content(elem) or " ".join(f"{_local(k)}={v}" for k, v in elem.attrib.items())
                            if value:
                                md[name] = value[:500]
                    meta["metadata"] = md
                except Exception as e:
                    log.debug("ODF meta.xml parse failed: %s", e)

            if "content.xml" in names:
                root = ET.fromstring(zf.read("content.xml"))
                page_count = 0
                for elem in root.iter():
                    name = _local(elem.tag)
                    if odf_type == "ods" and name == "table":
                        table_name = elem.attrib.get("{urn:oasis:names:tc:opendocument:xmlns:table:1.0}name", "Sheet")
                        text_parts.append(f"Worksheet: {table_name}")
                        page_count += 1
                    elif odf_type == "odp" and name == "page":
                        page_name = elem.attrib.get("{urn:oasis:names:tc:opendocument:xmlns:drawing:1.0}name", f"Slide {page_count + 1}")
                        text_parts.append(f"Slide: {page_name}")
                        page_count += 1
                    elif name in ("h", "p"):
                        value = _xml_text_content(elem)
                        if value:
                            text_parts.append(value[:1000])
                    elif odf_type == "ods" and name == "table-row":
                        cells: List[str] = []
                        for child in list(elem):
                            if _local(child.tag) == "table-cell":
                                value = _xml_text_content(child)
                                if value:
                                    cells.append(value[:300])
                        if cells:
                            text_parts.append(" | ".join(cells))
                meta["pages"] = page_count or 1

            try:
                meta["embedded_images"] = len([n for n in names if n.lower().startswith("pictures/")])
            except Exception:
                meta["embedded_images"] = 0
    except zipfile.BadZipFile as e:
        raise RuntimeError(f"Invalid OpenDocument file: {e}") from e

    if not text_parts:
        text_parts.append("[No readable OpenDocument text could be extracted]")
    return "\n".join(text_parts), meta


def _read_image(path: str) -> Tuple[str, Dict[str, Any]]:
    """Best-effort metadata extraction for JPG/JPEG/PNG/TIF/TIFF/WEBP/HEIC document images."""
    meta: Dict[str, Any] = {"type": "image", "pages": 1, "embedded_images": 1, "metadata": {}}
    try:
        from PIL import Image, ExifTags
    except Exception as e:
        raise RuntimeError("Missing dependency: Pillow. Add Pillow>=10.0.0 to requirements.txt") from e

    if os.path.splitext(path)[1].lower() in (".heic", ".heif"):
        try:
            from pillow_heif import register_heif_opener
            register_heif_opener()
        except Exception as e:
            raise RuntimeError("Missing dependency: pillow-heif. Add pillow-heif>=0.16.0 to requirements.txt") from e

    with Image.open(path) as img:
        md: Dict[str, Any] = {
            "format": img.format or "",
            "mode": img.mode or "",
            "width": str(img.width),
            "height": str(img.height),
        }

        # PNG text/info chunks and software fields.
        try:
            for k, v in (img.info or {}).items():
                if isinstance(v, (str, int, float)):
                    md[str(k)] = str(v)[:500]
        except Exception:
            pass

        # JPEG EXIF fields.
        try:
            exif = img.getexif()
            tag_map = getattr(ExifTags, "TAGS", {})
            for tag_id, value in exif.items():
                tag_name = tag_map.get(tag_id, str(tag_id))
                if isinstance(value, bytes):
                    value = value[:80].hex()
                md[str(tag_name)] = str(value)[:500]
        except Exception:
            pass

    meta["metadata"] = md
    return "", meta



def _strip_html_to_text(html: str) -> str:
    """Safely convert HTML markup into readable text without executing anything."""
    try:
        from html.parser import HTMLParser
        from html import unescape
    except Exception:
        return re.sub(r"<[^>]+>", " ", html or "")

    class _TextExtractor(HTMLParser):
        def __init__(self):
            super().__init__(convert_charrefs=True)
            self.parts: List[str] = []
            self.skip_depth = 0
            self.title = ""
            self._in_title = False

        def handle_starttag(self, tag, attrs):
            tag = str(tag or "").lower()
            if tag in ("script", "style", "noscript", "template"):
                self.skip_depth += 1
            if tag == "title":
                self._in_title = True
            if tag in ("p", "div", "br", "li", "tr", "h1", "h2", "h3", "h4", "h5", "h6", "section", "article"):
                self.parts.append("\n")

        def handle_endtag(self, tag):
            tag = str(tag or "").lower()
            if tag in ("script", "style", "noscript", "template") and self.skip_depth:
                self.skip_depth -= 1
            if tag == "title":
                self._in_title = False
            if tag in ("p", "div", "li", "tr", "h1", "h2", "h3", "h4", "h5", "h6", "section", "article"):
                self.parts.append("\n")

        def handle_data(self, data):
            if self.skip_depth:
                return
            text = unescape(str(data or "")).strip()
            if not text:
                return
            if self._in_title and not self.title:
                self.title = text[:300]
            self.parts.append(text)
            self.parts.append(" ")

    parser = _TextExtractor()
    try:
        parser.feed(html or "")
        parser.close()
    except Exception:
        return re.sub(r"\s+", " ", re.sub(r"<[^>]+>", " ", html or "")).strip()
    text = "".join(parser.parts)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n\s*\n\s*\n+", "\n\n", text)
    return text.strip()


def _read_html(path: str) -> Tuple[str, Dict[str, Any]]:
    """Best-effort HTML/HTM extraction for archived pages and web evidence."""
    meta: Dict[str, Any] = {"type": "html", "pages": 1, "embedded_images": 0, "metadata": {}}
    with open(path, "rb") as f:
        data = f.read(3_000_000)
    html_text = ""
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            html_text = data.decode(enc, errors="replace")
            break
        except Exception:
            continue
    if not html_text:
        html_text = data.decode("latin-1", errors="replace")

    title = ""
    try:
        m = re.search(r"<title[^>]*>(.*?)</title>", html_text, re.I | re.S)
        if m:
            title = re.sub(r"\s+", " ", re.sub(r"<[^>]+>", " ", m.group(1))).strip()[:300]
    except Exception:
        pass

    meta_tags: Dict[str, str] = {}
    for m in re.finditer(r"<meta\s+([^>]+)>", html_text[:500_000], re.I | re.S):
        attrs = m.group(1)
        name_m = re.search(r'(?:name|property)=["\']?([^"\'\s>]+)', attrs, re.I)
        content_m = re.search(r'content=["\']([^"\']{0,500})', attrs, re.I | re.S)
        if name_m and content_m:
            meta_tags[name_m.group(1).lower()[:80]] = re.sub(r"\s+", " ", content_m.group(1)).strip()[:500]
            if len(meta_tags) >= 20:
                break

    img_count = len(re.findall(r"<img\b", html_text, re.I))
    meta["embedded_images"] = img_count
    md = {"title": title, "html_bytes_scanned": str(len(data)), "image_tag_count": str(img_count)}
    md.update({f"meta_{k}": v for k, v in meta_tags.items()})
    meta["metadata"] = md
    body_text = _strip_html_to_text(html_text)
    if title and title not in body_text[:500]:
        body_text = f"Title: {title}\n\n{body_text}"
    return body_text, meta


def _read_mhtml(path: str) -> Tuple[str, Dict[str, Any]]:
    """Extract text and metadata from MHTML/MHT web archives."""
    meta: Dict[str, Any] = {"type": "mhtml", "pages": 1, "embedded_images": 0, "metadata": {}}
    from email import policy
    from email.parser import BytesParser

    with open(path, "rb") as f:
        msg = BytesParser(policy=policy.default).parse(f)

    headers = {
        "from": str(msg.get("From", ""))[:300],
        "subject": str(msg.get("Subject", ""))[:500],
        "date": str(msg.get("Date", ""))[:200],
        "content_type": str(msg.get_content_type()),
    }
    text_parts: List[str] = []
    html_parts: List[str] = []
    image_count = 0
    part_count = 0

    def _part_text(part: Any) -> str:
        try:
            payload = part.get_content()
            if isinstance(payload, str):
                return payload
        except Exception:
            pass
        try:
            data = part.get_payload(decode=True) or b""
            charset = part.get_content_charset() or "utf-8"
            return data.decode(charset, errors="replace")
        except Exception:
            return ""

    if msg.is_multipart():
        for part in msg.walk():
            part_count += 1
            ctype = part.get_content_type()
            if ctype.startswith("image/"):
                image_count += 1
            elif ctype == "text/html":
                html = _part_text(part)
                if html:
                    html_parts.append(_strip_html_to_text(html))
            elif ctype == "text/plain":
                txt = _part_text(part).strip()
                if txt:
                    text_parts.append(txt)
    else:
        content = _part_text(msg)
        if msg.get_content_type() == "text/html":
            html_parts.append(_strip_html_to_text(content))
        else:
            text_parts.append(content)

    meta["embedded_images"] = image_count
    headers.update({"part_count": str(part_count), "image_part_count": str(image_count)})
    meta["metadata"] = headers
    text = "\n\n".join(x for x in (text_parts + html_parts) if x and x.strip()).strip()
    return text or "[No readable MHTML text could be extracted]", meta


def _read_xml(path: str) -> Tuple[str, Dict[str, Any]]:
    """Best-effort XML extraction for government/system exports and evidence logs."""
    meta: Dict[str, Any] = {"type": "xml", "pages": 1, "embedded_images": 0, "metadata": {}}
    import xml.etree.ElementTree as ET
    data = open(path, "rb").read(3_000_000)
    text_raw = data.decode("utf-8", errors="replace")

    def _local(tag: str) -> str:
        return str(tag).split("}")[-1]

    try:
        root = ET.fromstring(data)
        tag_counts: Counter = Counter()
        lines: List[str] = [f"Root: {_local(root.tag)}"]
        for elem in root.iter():
            name = _local(elem.tag)
            tag_counts[name] += 1
            value = re.sub(r"\s+", " ", " ".join(t.strip() for t in elem.itertext() if t and t.strip())).strip()
            attrs = " ".join(f"{_local(k)}={v}" for k, v in list(elem.attrib.items())[:8])
            if value or attrs:
                line = f"{name}: " + (value[:800] if value else attrs[:800])
                lines.append(line)
            if len(lines) >= 2500:
                lines.append("[VeriFYD note: XML rendering truncated for analysis length]")
                break
        meta["metadata"] = {
            "root_tag": _local(root.tag),
            "element_count": str(sum(tag_counts.values())),
            "top_tags": ", ".join(f"{k}:{v}" for k, v in tag_counts.most_common(15)),
        }
        return "\n".join(lines), meta
    except Exception as e:
        meta["metadata"] = {"parse_error": str(e)[:200]}
        fallback = re.sub(r"\s+", " ", re.sub(r"<[^>]+>", " ", text_raw)).strip()
        return fallback or "[No readable XML text could be extracted]", meta


def _json_to_lines(obj: Any, prefix: str = "", limit: int = 2500) -> List[str]:
    """Flatten JSON into readable key/value lines for detection and certified rendering."""
    lines: List[str] = []
    def walk(value: Any, key: str, depth: int = 0):
        if len(lines) >= limit or depth > 8:
            return
        if isinstance(value, dict):
            if key:
                lines.append(f"{key}: object({len(value)})")
            for k, v in list(value.items())[:200]:
                walk(v, f"{key}.{k}" if key else str(k), depth + 1)
        elif isinstance(value, list):
            lines.append(f"{key}: array({len(value)})")
            for i, item in enumerate(value[:80]):
                walk(item, f"{key}[{i}]", depth + 1)
        else:
            txt = str(value)
            if txt.strip():
                lines.append(f"{key}: {txt[:800]}")
    walk(obj, prefix)
    return lines


def _read_json(path: str) -> Tuple[str, Dict[str, Any]]:
    """Best-effort JSON extraction for audit logs, API exports, and AI/system evidence."""
    meta: Dict[str, Any] = {"type": "json", "pages": 1, "embedded_images": 0, "metadata": {}}
    data = open(path, "rb").read(5_000_000)
    raw = data.decode("utf-8", errors="replace")
    try:
        parsed = json.loads(raw)
        root_type = type(parsed).__name__
        if isinstance(parsed, dict):
            root_keys = list(parsed.keys())[:50]
            meta["metadata"] = {"root_type": root_type, "root_keys": ", ".join(map(str, root_keys)), "byte_count": str(len(data))}
        elif isinstance(parsed, list):
            meta["metadata"] = {"root_type": root_type, "item_count": str(len(parsed)), "byte_count": str(len(data))}
        else:
            meta["metadata"] = {"root_type": root_type, "byte_count": str(len(data))}
        lines = _json_to_lines(parsed)
        return "\n".join(lines) if lines else raw[:12000], meta
    except Exception as e:
        meta["metadata"] = {"parse_error": str(e)[:200], "byte_count": str(len(data))}
        return raw[:12000], meta


def _read_vsdx(path: str) -> Tuple[str, Dict[str, Any]]:
    """Best-effort Visio VSDX extraction for network/engineering diagrams."""
    import zipfile
    import xml.etree.ElementTree as ET

    meta: Dict[str, Any] = {"type": "vsdx", "pages": 0, "embedded_images": 0, "metadata": {}}
    text_parts: List[str] = ["VeriFYD Visio Diagram Text Extraction"]

    def _local(tag: str) -> str:
        return str(tag).split("}")[-1]

    def _elem_text(elem: Any) -> str:
        try:
            return re.sub(r"\s+", " ", " ".join(t.strip() for t in elem.itertext() if t and t.strip())).strip()
        except Exception:
            return ""

    try:
        with zipfile.ZipFile(path, "r") as zf:
            names = zf.namelist()
            page_xmls = [n for n in names if n.lower().startswith("visio/pages/") and n.lower().endswith(".xml")]
            image_files = [n for n in names if n.lower().startswith("visio/media/")]
            meta["pages"] = len(page_xmls)
            meta["embedded_images"] = len(image_files)
            md: Dict[str, str] = {
                "page_count": str(len(page_xmls)),
                "embedded_media_count": str(len(image_files)),
                "package_file_count": str(len(names)),
            }

            for prop_name in ("docProps/core.xml", "docProps/app.xml"):
                if prop_name in names:
                    try:
                        root = ET.fromstring(zf.read(prop_name))
                        for elem in root.iter():
                            value = _elem_text(elem)
                            if value:
                                md[_local(elem.tag).lower()] = value[:500]
                    except Exception:
                        pass

            for idx, name in enumerate(page_xmls[:80], start=1):
                try:
                    root = ET.fromstring(zf.read(name))
                except Exception:
                    continue
                page_lines: List[str] = []
                for elem in root.iter():
                    lname = _local(elem.tag).lower()
                    # Visio shape text is commonly in Text elements, but some
                    # connectors/properties carry useful labels in attributes.
                    if lname in ("text", "cp", "pp", "tp"):
                        value = _elem_text(elem)
                        if value and len(value) > 1:
                            page_lines.append(value[:800])
                    elif lname in ("shape", "page"):
                        attrs = []
                        for k in ("Name", "NameU", "Text", "Master", "ID"):
                            if k in elem.attrib and str(elem.attrib[k]).strip():
                                attrs.append(f"{k}={elem.attrib[k]}")
                        if attrs:
                            page_lines.append("; ".join(attrs)[:500])
                if page_lines:
                    text_parts.append(f"\nPage {idx}: {os.path.basename(name)}")
                    # Deduplicate while preserving order.
                    seen = set()
                    for line in page_lines:
                        key = re.sub(r"\s+", " ", line).strip().lower()[:200]
                        if key and key not in seen:
                            seen.add(key)
                            text_parts.append(line)
                else:
                    text_parts.append(f"\nPage {idx}: {os.path.basename(name)} [No extractable shape text]")
            meta["metadata"] = md
    except zipfile.BadZipFile as e:
        raise RuntimeError(f"Invalid VSDX package: {e}") from e

    return "\n".join(text_parts), meta


ZIP_SUPPORTED_INNER_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt",
    ".odt", ".ods", ".odp", ".txt", ".md", ".csv", ".rtf", ".eml", ".msg",
    ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tif", ".tiff", ".webp", ".heic", ".heif",
    ".html", ".htm", ".mhtml", ".mht", ".xml", ".json", ".svg", ".vsdx",
    ".yaml", ".yml", ".ini", ".log", ".sql",
    ".pst", ".ost", ".dwg", ".dxf",
}
ZIP_DANGEROUS_EXTENSIONS = {
    ".exe", ".dll", ".bat", ".cmd", ".com", ".scr", ".ps1", ".vbs", ".js",
    ".jar", ".msi", ".apk", ".app", ".sh", ".bash", ".zsh", ".lnk",
}


def _is_safe_zip_member(name: str) -> bool:
    """Reject absolute paths, traversal, and Windows drive paths inside ZIP files."""
    name = str(name or "").replace("\\", "/")
    if not name or name.endswith("/"):
        return False
    if name.startswith("/") or name.startswith("../") or "/../" in name or name == "..":
        return False
    if re.match(r"^[A-Za-z]:", name):
        return False
    return True


def _read_zip(path: str) -> Tuple[str, Dict[str, Any]]:
    """
    Safely inspect a ZIP evidence package.

    This does not execute anything from the archive. It builds a manifest, hashes
    contained files, analyzes supported document/image/email formats by extracting
    them to temporary files, and skips scripts/executables or suspicious paths.
    """
    import zipfile
    import tempfile
    import shutil

    meta: Dict[str, Any] = {"type": "zip", "pages": 0, "embedded_images": 0, "metadata": {}}
    text_parts: List[str] = ["VeriFYD ZIP Evidence Package Manifest"]
    manifest: List[Dict[str, Any]] = []

    max_files = 100
    max_total_uncompressed = 250 * 1024 * 1024
    max_member_uncompressed = 50 * 1024 * 1024
    max_text_per_file = 2500

    supported_count = 0
    skipped_count = 0
    dangerous_count = 0
    image_count = 0
    total_uncompressed = 0

    tmp_dir = tempfile.mkdtemp(prefix="verifyd_zip_extract_")
    try:
        with zipfile.ZipFile(path, "r") as zf:
            infos = zf.infolist()
            file_infos = [i for i in infos if not i.is_dir()]
            meta["pages"] = len(file_infos)

            for idx, info in enumerate(file_infos, start=1):
                if idx > max_files:
                    skipped_count += max(0, len(file_infos) - max_files + 1)
                    text_parts.append(f"[VeriFYD note: ZIP manifest truncated after {max_files} files]")
                    break

                raw_name = info.filename or f"file_{idx}"
                safe_name = raw_name.replace("\\", "/")
                ext = os.path.splitext(safe_name.lower())[1]
                file_size = int(getattr(info, "file_size", 0) or 0)
                comp_size = int(getattr(info, "compress_size", 0) or 0)
                total_uncompressed += file_size

                entry: Dict[str, Any] = {
                    "name": safe_name[:500],
                    "extension": ext,
                    "size_bytes": file_size,
                    "compressed_bytes": comp_size,
                    "supported": False,
                    "status": "skipped",
                    "sha256": "",
                }

                if not _is_safe_zip_member(safe_name):
                    entry["status"] = "unsafe_path_skipped"
                    skipped_count += 1
                    manifest.append(entry)
                    continue

                if ext in ZIP_DANGEROUS_EXTENSIONS:
                    entry["status"] = "dangerous_file_skipped"
                    dangerous_count += 1
                    skipped_count += 1
                    manifest.append(entry)
                    continue

                if total_uncompressed > max_total_uncompressed:
                    entry["status"] = "zip_total_size_limit_reached"
                    skipped_count += 1
                    manifest.append(entry)
                    text_parts.append("[VeriFYD note: ZIP total uncompressed-size safety limit reached]")
                    break

                if file_size > max_member_uncompressed:
                    entry["status"] = "member_too_large_skipped"
                    skipped_count += 1
                    manifest.append(entry)
                    continue

                if ext in (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tif", ".tiff", ".webp", ".heic", ".heif"):
                    image_count += 1

                try:
                    with zf.open(info, "r") as src:
                        data = src.read(max_member_uncompressed + 1)
                    if len(data) > max_member_uncompressed:
                        entry["status"] = "member_read_limit_skipped"
                        skipped_count += 1
                        manifest.append(entry)
                        continue
                    entry["sha256"] = hashlib.sha256(data).hexdigest()
                except Exception as e:
                    entry["status"] = f"read_error: {str(e)[:80]}"
                    skipped_count += 1
                    manifest.append(entry)
                    continue

                if ext not in ZIP_SUPPORTED_INNER_EXTENSIONS:
                    entry["status"] = "unsupported_type_manifest_only"
                    skipped_count += 1
                    manifest.append(entry)
                    text_parts.append(f"{safe_name} | {file_size} bytes | unsupported type | sha256={entry['sha256'][:16]}")
                    continue

                member_path = os.path.join(tmp_dir, f"member_{idx}{ext}")
                try:
                    with open(member_path, "wb") as fh:
                        fh.write(data)
                    inner_text, inner_meta = _extract_document(member_path)
                    entry["supported"] = True
                    entry["status"] = "analyzed"
                    entry["inner_type"] = inner_meta.get("type", ext.lstrip("."))
                    supported_count += 1

                    text_parts.append(
                        f"\n--- ZIP Member {idx}: {safe_name} | {file_size} bytes | {ext} | sha256={entry['sha256'][:16]} ---"
                    )
                    cleaned = re.sub(r"\s+", " ", inner_text or "").strip()
                    if cleaned:
                        text_parts.append(cleaned[:max_text_per_file])
                    else:
                        text_parts.append("[No extractable text from this member]")
                except Exception as e:
                    entry["status"] = f"analysis_error: {str(e)[:100]}"
                    skipped_count += 1
                finally:
                    try:
                        if os.path.exists(member_path):
                            os.remove(member_path)
                    except Exception:
                        pass

                manifest.append(entry)

    except zipfile.BadZipFile as e:
        raise RuntimeError(f"Invalid ZIP archive: {e}") from e
    finally:
        try:
            shutil.rmtree(tmp_dir, ignore_errors=True)
        except Exception:
            pass

    meta["embedded_images"] = image_count
    meta["metadata"] = {
        "zip_file_count": str(len(manifest)),
        "zip_supported_count": str(supported_count),
        "zip_skipped_count": str(skipped_count),
        "zip_dangerous_count": str(dangerous_count),
        "zip_total_uncompressed_bytes": str(total_uncompressed),
        "zip_manifest_json": json.dumps(manifest[:80], ensure_ascii=False)[:12000],
    }

    if dangerous_count:
        text_parts.append(f"\n[VeriFYD warning: {dangerous_count} potentially dangerous file(s) were skipped and not executed]")
    if supported_count == 0:
        text_parts.append("\n[No supported documents were found inside this ZIP package; manifest-only analysis performed]")

    return "\n".join(text_parts), meta



def _read_textlike_config(path: str, doc_type: str) -> Tuple[str, Dict[str, Any]]:
    """Read plain-text evidence/config formats such as YAML, INI, LOG, and SQL."""
    meta: Dict[str, Any] = {"type": doc_type, "pages": 1, "embedded_images": 0, "metadata": {}}
    with open(path, "rb") as f:
        data = f.read(5_000_000)
    text = ""
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            text = data.decode(enc, errors="replace")
            break
        except Exception:
            continue
    if not text:
        text = data.decode("latin-1", errors="replace")

    line_count = text.count("\n") + 1 if text else 0
    nonblank = len([ln for ln in text.splitlines() if ln.strip()])
    meta["metadata"] = {
        "byte_count": str(len(data)),
        "line_count": str(line_count),
        "nonblank_line_count": str(nonblank),
    }
    if doc_type in ("yaml", "yml"):
        keys = re.findall(r"^\s*([A-Za-z0-9_\-.]+)\s*:", text, flags=re.MULTILINE)
        meta["metadata"]["top_keys"] = ", ".join(keys[:40])
    elif doc_type == "ini":
        sections = re.findall(r"^\s*\[([^\]]+)\]", text, flags=re.MULTILINE)
        meta["metadata"]["sections"] = ", ".join(sections[:40])
    elif doc_type == "sql":
        stmts = re.findall(r"\b(SELECT|INSERT|UPDATE|DELETE|CREATE|ALTER|DROP|WITH)\b", text, flags=re.IGNORECASE)
        meta["metadata"]["sql_statement_keywords"] = ", ".join(stmts[:50])
    return text[:120000], meta


def _read_svg(path: str) -> Tuple[str, Dict[str, Any]]:
    """Best-effort SVG/XML extraction for vector evidence and diagrams."""
    text, meta = _read_xml(path)
    meta["type"] = "svg"
    md = dict(meta.get("metadata") or {})
    md["format"] = "SVG vector image/XML"
    meta["metadata"] = md
    return text, meta


def _read_dxf(path: str) -> Tuple[str, Dict[str, Any]]:
    """Best-effort DXF extraction. DXF is commonly text-based CAD evidence."""
    meta: Dict[str, Any] = {"type": "dxf", "pages": 1, "embedded_images": 0, "metadata": {}}
    with open(path, "rb") as f:
        data = f.read(8_000_000)
    text = data.decode("utf-8", errors="replace")
    if "SECTION" not in text[:200000].upper() and "ENTITIES" not in text[:500000].upper():
        text = data.decode("latin-1", errors="replace")

    layers = re.findall(r"\n\s*8\s*\n([^\n\r]{1,120})", text)
    entity_types = re.findall(r"\n\s*0\s*\n(LINE|LWPOLYLINE|POLYLINE|CIRCLE|ARC|TEXT|MTEXT|INSERT|DIMENSION|HATCH|SPLINE)\b", text, flags=re.I)
    readable = re.findall(r"[A-Za-z0-9][A-Za-z0-9\s\.,;:\-_/@$%#&()\[\]{}'\"!?]{4,}", text[:1_500_000])
    lines = ["VeriFYD DXF CAD Evidence Extraction"]
    if layers:
        seen_layers = []
        for layer in layers:
            layer = layer.strip()
            if layer and layer not in seen_layers:
                seen_layers.append(layer)
            if len(seen_layers) >= 60:
                break
        lines.append("Layers: " + ", ".join(seen_layers))
    if entity_types:
        counts = Counter(x.upper() for x in entity_types)
        lines.append("Entity counts: " + ", ".join(f"{k}:{v}" for k, v in counts.most_common(20)))
    for frag in readable[:800]:
        clean = re.sub(r"\s+", " ", frag).strip()
        if clean and clean not in lines:
            lines.append(clean[:500])
    meta["metadata"] = {
        "byte_count": str(len(data)),
        "layer_count_sample": str(len(set(layers))),
        "entity_type_sample": ", ".join(sorted(set(x.upper() for x in entity_types))[:25]),
    }
    return "\n".join(lines), meta


def _read_binary_evidence(path: str, doc_type: str) -> Tuple[str, Dict[str, Any]]:
    """
    Lightweight binary evidence handler for PST/OST/DWG.

    These formats are intentionally treated as package-level evidence here, not
    deeply decoded archives. VeriFYD records file size, hash is handled upstream,
    magic bytes, and recoverable string fragments without executing anything.
    """
    meta: Dict[str, Any] = {"type": doc_type, "pages": 1, "embedded_images": 0, "metadata": {}}
    size = os.path.getsize(path)
    with open(path, "rb") as f:
        head = f.read(64)
        f.seek(0)
        sample = f.read(3_000_000)

    fragments: List[str] = []
    for enc in ("utf-16-le", "latin-1"):
        try:
            decoded = sample.decode(enc, errors="ignore").replace("\x00", " ")
        except Exception:
            continue
        pieces = re.findall(r"[A-Za-z0-9][A-Za-z0-9\s\.,;:\-_/@$%#&()\[\]{}'\"!?]{5,}", decoded)
        for piece in pieces:
            clean = re.sub(r"\s+", " ", piece).strip()
            if len(clean) >= 6:
                fragments.append(clean[:500])
            if len(fragments) >= 500:
                break
        if len(fragments) >= 500:
            break

    seen = set()
    unique: List[str] = []
    for frag in fragments:
        key = frag.lower()[:160]
        if key not in seen:
            seen.add(key)
            unique.append(frag)
        if len(unique) >= 250:
            break

    type_label = {
        "pst": "Outlook PST email archive",
        "ost": "Outlook OST offline mailbox archive",
        "dwg": "AutoCAD DWG drawing",
    }.get(doc_type, doc_type.upper())
    meta["metadata"] = {
        "format": type_label,
        "byte_count": str(size),
        "magic_hex": head.hex()[:128],
        "readable_string_sample_count": str(len(unique)),
        "analysis_mode": "lightweight binary evidence manifest; deep native parsing not enabled",
    }
    lines = [
        f"VeriFYD {type_label} Evidence Record",
        f"File size: {size} bytes",
        f"Magic/header hex: {head.hex()[:128]}",
        "Analysis mode: lightweight evidence certification. Deep native parsing/rendering is not enabled for this format.",
    ]
    if unique:
        lines.append("\nReadable string samples:")
        lines.extend(unique)
    else:
        lines.append("\nNo readable text strings found in scanned binary sample.")
    return "\n".join(lines), meta

def _extract_document(path: str) -> Tuple[str, Dict[str, Any]]:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return _read_pdf(path)
    if ext == ".docx":
        return _read_docx(path)
    if ext == ".doc":
        return _read_doc(path)
    if ext == ".xlsx":
        return _read_xlsx(path)
    if ext == ".xls":
        return _read_xls(path)
    if ext == ".pptx":
        return _read_pptx(path)
    if ext == ".ppt":
        return _read_ppt(path)
    if ext in (".odt", ".ods", ".odp"):
        return _read_odf(path)
    if ext in (".txt", ".md", ".csv"):
        return _read_txt(path)
    if ext in (".yaml", ".yml", ".ini", ".log", ".sql"):
        return _read_textlike_config(path, ext.lstrip("."))
    if ext == ".rtf":
        return _read_rtf(path)
    if ext == ".eml":
        return _read_eml(path)
    if ext == ".msg":
        return _read_msg(path)
    if ext in (".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tif", ".tiff", ".webp", ".heic", ".heif"):
        return _read_image(path)
    if ext in (".html", ".htm"):
        return _read_html(path)
    if ext in (".mhtml", ".mht"):
        return _read_mhtml(path)
    if ext == ".xml":
        return _read_xml(path)
    if ext == ".svg":
        return _read_svg(path)
    if ext == ".json":
        return _read_json(path)
    if ext == ".vsdx":
        return _read_vsdx(path)
    if ext == ".dxf":
        return _read_dxf(path)
    if ext in (".pst", ".ost", ".dwg"):
        return _read_binary_evidence(path, ext.lstrip("."))
    if ext == ".zip":
        return _read_zip(path)
    raise RuntimeError(f"Unsupported document format: {ext}")


def _metadata_score(meta: Dict[str, Any]) -> Tuple[int, List[str]]:
    score = 0
    flags: List[str] = []
    md = {str(k).lower(): str(v).lower() for k, v in (meta.get("metadata") or {}).items()}
    joined = " ".join(md.values())

    for kw in AI_PRODUCER_KEYWORDS:
        if kw in joined:
            score += 35
            flags.append(f"metadata references possible AI tool: {kw}")
            break

    producer = md.get("producer", "") or md.get("creator", "")
    if producer and not any(k in producer for k in COMMON_GENERATOR_KEYWORDS):
        if any(k in producer for k in ("pdfkit", "wkhtml", "weasyprint", "reportlab", "chromium", "headless")):
            score += 10
            flags.append("document was generated by an automated PDF/rendering pipeline")

    created = md.get("creationdate", "") or md.get("created", "")
    modified = md.get("moddate", "") or md.get("modified", "")
    if created and modified and created != modified:
        score += 5
        flags.append("creation and modification timestamps differ")

    if meta.get("type") == "pdf" and meta.get("pages", 0) > 0 and meta.get("embedded_images", 0) >= meta.get("pages", 0):
        score += 8
        flags.append("PDF appears image-heavy or scanned; visual tamper review recommended")

    return min(score, 60), flags


def _text_stats_score(text: str) -> Tuple[int, List[str], Dict[str, Any]]:
    flags: List[str] = []
    clean = re.sub(r"\s+", " ", text or "").strip()
    words = re.findall(r"[A-Za-z']+", clean.lower())
    sentences = [s.strip() for s in re.split(r"[.!?]+", clean) if len(s.strip()) > 2]
    stats = {
        "word_count": len(words),
        "sentence_count": len(sentences),
        "unique_word_ratio": 0.0,
        "avg_sentence_len": 0.0,
        "sentence_len_cv": 0.0,
        "ai_phrase_hits": [],
    }

    if not words:
        return 20, ["no extractable text; document may be scanned or image-only"], stats

    counts = Counter(words)
    unique_ratio = len(counts) / max(1, len(words))
    stats["unique_word_ratio"] = round(unique_ratio, 3)

    sent_lens = [len(re.findall(r"[A-Za-z']+", s)) for s in sentences]
    if sent_lens:
        avg = sum(sent_lens) / len(sent_lens)
        var = sum((x - avg) ** 2 for x in sent_lens) / len(sent_lens)
        cv = math.sqrt(var) / avg if avg else 0
        stats["avg_sentence_len"] = round(avg, 2)
        stats["sentence_len_cv"] = round(cv, 3)
    else:
        cv = 0

    score = 0
    lower = clean.lower()
    phrase_hits = [p for p in AI_TEXT_PATTERNS if p in lower]
    stats["ai_phrase_hits"] = phrase_hits[:8]
    if len(phrase_hits) >= 3:
        score += 16
        flags.append(f"multiple AI-style phrases detected ({len(phrase_hits)})")
    elif len(phrase_hits) >= 1:
        score += 6
        flags.append("some AI-style phrasing detected")

    if len(words) >= 250 and unique_ratio < 0.34:
        score += 10
        flags.append("low vocabulary variation for document length")

    if len(sent_lens) >= 8 and cv < 0.38:
        score += 10
        flags.append("unusually uniform sentence rhythm")

    # Very short docs are difficult to authenticate using writing style.
    if len(words) < 80:
        score = min(score, 10)
        flags.append("short document; text-style AI confidence reduced")

    return min(score, 40), flags, stats


def _gpt_document_score(text: str, meta: Dict[str, Any], ext: str) -> Tuple[int, str, List[str], bool]:
    """Use OpenAI if configured. Safe fallback returns unavailable."""
    api_key = os.environ.get("OPENAI_API_KEY", "")
    if not api_key:
        return 50, "GPT document analysis unavailable: OPENAI_API_KEY not set.", [], False

    try:
        from openai import OpenAI
        client = OpenAI(api_key=api_key)
        prompt = {
            "task": "Document authenticity analysis",
            "instructions": (
                "Return JSON only. Estimate probability this document text/metadata was AI-generated, "
                "synthetically assembled, or materially tampered. Do not judge whether claims are true unless internal inconsistencies are present. "
                "Use 0=confident authentic human/original, 100=confident AI/generated/tampered."
            ),
            "file_type": ext,
            "metadata": meta.get("metadata", {}),
            "pages": meta.get("pages", 0),
            "embedded_images": meta.get("embedded_images", 0),
            "text_excerpt": _clip_text(text),
            "schema": {"ai_probability": "integer 0-100", "flags": ["short strings"], "reasoning": "one concise paragraph"},
        }
        resp = client.chat.completions.create(
            model=os.environ.get("OPENAI_TEXT_MODEL", "gpt-4o-mini"),
            temperature=0,
            messages=[
                {"role": "system", "content": "You are a forensic document authentication analyst. Return strict JSON."},
                {"role": "user", "content": json.dumps(prompt)},
            ],
        )
        raw = resp.choices[0].message.content or "{}"
        raw = raw.strip().removeprefix("```json").removeprefix("```").removesuffix("```").strip()
        data = json.loads(raw)
        score = max(0, min(100, _safe_int(data.get("ai_probability", 50), 50)))
        flags = [str(x) for x in data.get("flags", [])][:8]
        reasoning = str(data.get("reasoning", ""))[:1200]
        return score, reasoning, flags, True
    except Exception as e:
        log.warning("GPT document analysis failed: %s", e)
        return 50, f"GPT document analysis failed: {str(e)[:120]}", [], False


def _label_from_ai_score(ai_score: int) -> Tuple[int, str]:
    authenticity = 100 - int(round(ai_score))
    if authenticity >= THRESHOLD_REAL:
        label = "REAL"
    elif authenticity >= THRESHOLD_UNDETERMINED:
        label = "UNDETERMINED"
    else:
        label = "AI"
    return authenticity, label


def _build_reasoning(label: str, authenticity: int, flags: List[str], gpt_reasoning: str, meta: Dict[str, Any], text_stats: Dict[str, Any]) -> str:
    doc_type = meta.get("type", "document").upper()
    if label == "REAL":
        lead = f"This {doc_type} document shows no strong indicators of AI generation or material tampering."
    elif label == "UNDETERMINED":
        lead = f"This {doc_type} document contains mixed or incomplete authenticity signals."
    else:
        lead = f"This {doc_type} document shows indicators consistent with AI generation, automated assembly, or tampering."

    evidence = "; ".join(flags[:4]) if flags else "metadata and text signals were reviewed"
    if gpt_reasoning and "unavailable" not in gpt_reasoning.lower() and "failed" not in gpt_reasoning.lower():
        return f"{lead} Key evidence: {evidence}. GPT review: {gpt_reasoning} Authenticity score: {authenticity}%."
    return f"{lead} Key evidence: {evidence}. Authenticity score: {authenticity}%."




# ─────────────────────────────────────────────
# Document Risk Report / Metadata Consistency Engine
# ─────────────────────────────────────────────

def _metadata_lookup(md: Dict[str, Any], candidates: Tuple[str, ...]) -> str:
    """Find a metadata value using loose key matching across formats."""
    if not md:
        return ""
    lowered = {str(k).lower().replace("_", "").replace("-", "").replace("/", "").strip(): v for k, v in md.items()}
    for candidate in candidates:
        key = candidate.lower().replace("_", "").replace("-", "").replace("/", "").strip()
        if key in lowered and str(lowered[key]).strip():
            return str(lowered[key]).strip()
    for k, v in lowered.items():
        for candidate in candidates:
            key = candidate.lower().replace("_", "").replace("-", "").replace("/", "").strip()
            if key and key in k and str(v).strip():
                return str(v).strip()
    return ""


def _parse_any_date(value: Any) -> datetime | None:
    """Parse common PDF/Office/EXIF/ISO timestamp formats into a naive UTC-ish datetime."""
    if not value:
        return None
    raw = str(value).strip()
    if not raw:
        return None
    s = raw
    if s.startswith("D:"):
        # PDF format: D:YYYYMMDDHHmmSS+/-HH'mm'
        s = s[2:16]
        try:
            return datetime.strptime(s, "%Y%m%d%H%M%S")
        except Exception:
            pass
    # Normalize common ExifTool timezone form: 2026:04:09 13:52:46-04:00
    tz_normalized = re.sub(r"^(\d{4}:\d{2}:\d{2} \d{2}:\d{2}:\d{2})([+-]\d{2}:?\d{2})$", r"\1", raw)
    # EXIF-like timestamps.
    for fmt in (
        "%Y:%m:%d %H:%M:%S",
        "%Y-%m-%dT%H:%M:%S.%f%z",
        "%Y-%m-%dT%H:%M:%S%z",
        "%Y-%m-%dT%H:%M:%S.%f",
        "%Y-%m-%dT%H:%M:%S",
        "%Y-%m-%d %H:%M:%S",
        "%Y-%m-%d",
        "%m/%d/%Y %H:%M:%S",
        "%m/%d/%Y",
    ):
        try:
            dt = datetime.strptime(tz_normalized.replace("Z", "+0000"), fmt)
            if dt.tzinfo:
                return dt.astimezone().replace(tzinfo=None)
            return dt
        except Exception:
            continue
    try:
        dt = datetime.fromisoformat(tz_normalized.replace("Z", "+00:00"))
        if dt.tzinfo:
            return dt.astimezone().replace(tzinfo=None)
        return dt
    except Exception:
        return None


def _risk_level(score: int) -> str:
    if score >= 70:
        return "HIGH"
    if score >= 35:
        return "MEDIUM"
    return "LOW"


def _build_document_risk_report(path: str, ext: str, meta: Dict[str, Any], meta_score: int,
                                text_score: int, gpt_score: int, flags: List[str]) -> Dict[str, Any]:
    """
    Build a business-friendly risk report for legal/insurance/compliance users.

    This report explains why a document looks clean or suspicious. It is separate
    from the AI score so VeriFYD can certify both authenticity indicators and
    evidence integrity indicators.
    """
    md = dict(meta.get("metadata") or {})
    doc_type = str(meta.get("type") or ext.lstrip(".") or "document")
    pages = _safe_int(meta.get("pages", 0), 0)
    embedded_images = _safe_int(meta.get("embedded_images", 0), 0)

    created_raw = _metadata_lookup(md, (
        "created", "creationdate", "creation-date", "createdate", "datecreated",
        "created_time", "creation_time", "create_date", "datetimeoriginal",
    ))
    modified_raw = _metadata_lookup(md, (
        "modified", "moddate", "modificationdate", "lastmodified", "last_modified_by",
        "date", "modify_date", "metadata_date", "lastsaved", "lastsavedby",
    ))
    creator = _metadata_lookup(md, ("creator", "producer", "generator", "application", "software", "author"))
    author = _metadata_lookup(md, ("author", "creator", "initial-creator", "last_modified_by", "lastmodifiedby"))

    created_dt = _parse_any_date(created_raw)
    modified_dt = _parse_any_date(modified_raw)
    now = datetime.utcnow()
    future_threshold = now.replace(microsecond=0)

    risk_points = 0
    reasons: List[str] = []
    passed: List[str] = []
    warnings: List[str] = []

    if created_dt and modified_dt:
        if modified_dt < created_dt:
            risk_points += 45
            reasons.append("Modification timestamp predates creation timestamp.")
        else:
            passed.append("Creation/modification timestamp order is consistent.")
    elif created_raw or modified_raw:
        warnings.append("Only one primary document timestamp was found.")
        risk_points += 8
    else:
        warnings.append("No primary creation/modification timestamp was found.")
        risk_points += 12

    actual_future_timestamp = False
    for label, dt in (("creation", created_dt), ("modification", modified_dt)):
        # Allow a small clock/time-zone tolerance. This prevents GPT or metadata
        # phrasing from creating a false future-date warning when the parsed
        # document dates are actually in the past.
        if dt and dt > future_threshold:
            actual_future_timestamp = True
            risk_points += 25
            reasons.append(f"Document {label} timestamp appears to be in the future.")

    joined_md = " ".join(str(v).lower() for v in md.values() if v is not None)
    ai_hits = [kw for kw in AI_PRODUCER_KEYWORDS if kw in joined_md]
    if ai_hits:
        risk_points += 35
        reasons.append("Metadata references AI/generative software: " + ", ".join(ai_hits[:5]))

    if meta_score >= 35:
        risk_points += 25
        reasons.append("Metadata analysis produced elevated risk indicators.")
    elif meta_score >= 15:
        risk_points += 10
        warnings.append("Metadata analysis produced moderate risk indicators.")
    else:
        passed.append("No high-risk metadata producer indicators were found.")

    if gpt_score >= 65:
        risk_points += 25
        reasons.append("Semantic analysis found elevated AI/manipulation indicators.")
    elif gpt_score >= 40:
        risk_points += 10
        warnings.append("Semantic analysis found moderate AI/manipulation indicators.")

    if not md:
        warnings.append("Document metadata is empty or stripped.")
        risk_points += 10
    elif len([v for v in md.values() if str(v).strip()]) <= 2:
        warnings.append("Document has very limited metadata.")
        risk_points += 6

    if not author:
        warnings.append("Author/owner metadata is missing.")
    if not creator:
        warnings.append("Creator/producer metadata is missing.")

    if doc_type == "pdf" and embedded_images > 0 and text_score == 0:
        risk_points += 12
        warnings.append("PDF appears image-heavy with little or no extractable text.")

    if flags:
        # Existing detector/GPT flags are useful, but validate future-date claims
        # against parsed metadata so the report does not show false positives.
        filtered_flags: List[str] = []
        for f in flags[:8]:
            f_text = str(f)
            f_low = f_text.lower()
            if "future" in f_low and "date" in f_low and not actual_future_timestamp:
                continue
            filtered_flags.append(f_text)
        if filtered_flags:
            flag_text = "; ".join(filtered_flags[:5])
            warnings.append("Detector flags: " + flag_text)

    risk_points = max(0, min(100, int(round(risk_points))))
    overall = _risk_level(risk_points)
    if reasons:
        metadata_integrity = "FAIL" if any("timestamp" in r.lower() or "future" in r.lower() for r in reasons) else "WARN"
    elif warnings:
        metadata_integrity = "WARN"
    else:
        metadata_integrity = "PASS"

    if not reasons:
        if warnings:
            reasons.append("No critical manipulation indicators found; minor metadata limitations noted.")
        else:
            reasons.append("No material metadata or structure inconsistencies were detected.")

    return {
        "overall_risk": overall,
        "risk_score": risk_points,
        "created": created_raw or "Not found",
        "modified": modified_raw or "Not found",
        "creator": creator or "Not found",
        "author": author or "Not found",
        "metadata_integrity": metadata_integrity,
        "file_structure": "PASS",
        "ai_indicators": "PRESENT" if (ai_hits or gpt_score >= 65) else "NONE DETECTED",
        "hidden_revisions": "NOT CHECKED",
        "digital_signature": "NOT CHECKED",
        "metadata_tool": (meta.get("external_metadata_tool") or {}).get("name", "Python extractors"),
        "metadata_tool_status": (meta.get("external_metadata_tool") or {}).get("status", "built_in"),
        "reasons": reasons[:8],
        "warnings": warnings[:8],
        "passed_checks": passed[:8],
        "signals_checked": [
            "Creation timestamp",
            "Modification timestamp",
            "Timestamp ordering",
            "Future-dated metadata",
            "Author / creator fields",
            "Producer / software fields",
            "Metadata stripping",
            "Embedded object count",
            "AI / manipulation indicators",
            "File structure extraction",
            "Optional ExifTool metadata extraction",
        ],
        "document_type": doc_type,
        "pages": pages,
        "embedded_images": embedded_images,
    }

def run_document_detection(path: str) -> Tuple[int, str, Dict[str, Any]]:
    """Run document authentication and return (authenticity, label, detail)."""
    ext = os.path.splitext(path)[1].lower()
    sha256 = _sha256_file(path)
    size_bytes = os.path.getsize(path)

    text, meta = _extract_document(path)
    meta = _merge_exiftool_metadata(path, meta)
    meta_score, meta_flags = _metadata_score(meta)
    text_score, text_flags, text_stats = _text_stats_score(text)
    gpt_score, gpt_reasoning, gpt_flags, gpt_available = _gpt_document_score(text, meta, ext)

    # Weighting: metadata/text are stable and cheap; GPT is semantic but not final authority.
    # If GPT unavailable, use metadata/text only and pull toward uncertain when evidence is weak.
    if gpt_available:
        combined = meta_score * 0.30 + text_score * 0.25 + gpt_score * 0.45
        w_gpt = 0.45
    else:
        base = meta_score * 0.55 + text_score * 0.45
        # Avoid over-certifying text-only docs without GPT review.
        combined = max(12.0, base) if base < 35 else base
        w_gpt = 0.0

    # Strong metadata AI flag should be respected even if writing looks ordinary.
    if meta_score >= 35:
        combined = max(combined, 62.0)

    # No extractable text + image-heavy PDF should not be certified REAL in MVP.
    if not text.strip() and meta.get("type") == "pdf":
        combined = max(combined, 45.0)

    combined = max(0, min(100, int(round(combined))))
    authenticity, label = _label_from_ai_score(combined)

    flags = meta_flags + text_flags + gpt_flags
    reasoning = _build_reasoning(label, authenticity, flags, gpt_reasoning, meta, text_stats)
    risk_report = _build_document_risk_report(path, ext, meta, meta_score, text_score, gpt_score, flags)

    detail = {
        "ai_score": combined,
        "authenticity": authenticity,
        "label": label,
        "signal_ai_score": int(round(meta_score * 0.55 + text_score * 0.45)),
        "gpt_ai_score": gpt_score,
        "gpt_available": gpt_available,
        "gpt_reasoning": reasoning,
        "gpt_flags": flags[:10],
        "metadata_score": meta_score,
        "text_score": text_score,
        "weight_signal": 1.0 - w_gpt,
        "weight_gpt": w_gpt,
        "blend_mode": "document metadata+text+gpt" if gpt_available else "document metadata+text only",
        "content_type": "document",
        "document_type": meta.get("type", ext.lstrip(".")),
        "sha256": sha256,
        "file_size_bytes": size_bytes,
        "pages": meta.get("pages", 0),
        "embedded_images": meta.get("embedded_images", 0),
        "text_stats": text_stats,
        "metadata": meta.get("metadata", {}),
        "external_metadata_tool": meta.get("external_metadata_tool", {}),
        "document_risk_report": risk_report,
        "risk_report": risk_report,
        "overall_risk": risk_report.get("overall_risk"),
        "risk_score": risk_report.get("risk_score"),
        "metadata_integrity": risk_report.get("metadata_integrity"),
        "threshold_real": THRESHOLD_REAL,
        "threshold_undet": THRESHOLD_UNDETERMINED,
    }
    log.info(
        "Document detection complete | type=%s meta=%d text=%d gpt=%d combined=%d auth=%d label=%s",
        meta.get("type"), meta_score, text_score, gpt_score, combined, authenticity, label,
    )
    return authenticity, label, detail
